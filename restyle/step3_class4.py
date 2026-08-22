"""Step 3 (class 4): approximate every class-4 element as a composition of
native, editable PPTX shapes (rectangles, ellipses, arrows, chevrons, lines,
...).

For each class-4 component:
  1. Claude returns a JSON decomposition into <= MAX_SHAPES primitive PPTX
     shapes, with positions/sizes as bbox-relative fractions.
  2. We build a one-element PPTX with those shapes, hard-clamped to the bbox.
  3. Render -> verify with Claude. Iterate up to MAX_VERIFY_ROUNDS times.
  4. Save the final ``<page>/<stem>.pptx`` under step3_class4_dir.
"""

from __future__ import annotations

import argparse
import html
import json
import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

import os

from common import (
    CLAUDE_OPUS,
    DPI,
    MIN_READABLE_PT,
    PIPELINE_DIR,
    SLIDECODER_PYTHON,
    anthropic_client,
    call_claude,
    case_paths,
    crop_bbox,
    encode_image_b64,
    fit_font_size_pt,
    carlito_metrics_unreliable,
    font_sizing_mode,
    hex_to_rgb,
    list_front_bg_pages,
    parse_json_loose,
    px_to_emu,
    render_pptx_to_png,
)


TARGET_CLASS = 4
MODEL = CLAUDE_OPUS
MAX_VERIFY_ROUNDS = 2
MAX_SHAPES = 32
MIN_EXPECTED_TEXT_NORM_LEN = 3
MAX_COVERAGE_MISSING = 0


SHAPE_TYPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL, "ellipse": MSO_SHAPE.OVAL, "circle": MSO_SHAPE.OVAL,
    "isoceles_triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "isosceles_triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "right_triangle": MSO_SHAPE.RIGHT_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "trapezoid": MSO_SHAPE.TRAPEZOID,
    "pentagon": MSO_SHAPE.PENTAGON, "hexagon": MSO_SHAPE.HEXAGON,
    "octagon": MSO_SHAPE.OCTAGON,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "left_arrow": MSO_SHAPE.LEFT_ARROW,
    "up_arrow": MSO_SHAPE.UP_ARROW, "down_arrow": MSO_SHAPE.DOWN_ARROW,
    "left_right_arrow": MSO_SHAPE.LEFT_RIGHT_ARROW,
    "up_down_arrow": MSO_SHAPE.UP_DOWN_ARROW,
    "chevron": MSO_SHAPE.CHEVRON,
    "star_5": MSO_SHAPE.STAR_5_POINT, "star": MSO_SHAPE.STAR_5_POINT,
    "heart": MSO_SHAPE.HEART, "cloud": MSO_SHAPE.CLOUD, "sun": MSO_SHAPE.SUN,
    "lightning_bolt": MSO_SHAPE.LIGHTNING_BOLT,
    "plaque": MSO_SHAPE.PLAQUE,
    "can": MSO_SHAPE.CAN, "cylinder": MSO_SHAPE.CAN,
}
ALLOWED_TYPES_DOC = (
    "rectangle, rounded_rectangle, oval, isoceles_triangle, right_triangle, "
    "diamond, parallelogram, trapezoid, pentagon, hexagon, octagon, "
    "right_arrow, left_arrow, up_arrow, down_arrow, left_right_arrow, "
    "up_down_arrow, chevron, star_5, heart, cloud, sun, lightning_bolt, "
    "plaque, cylinder, line"
)


EXTRACT_PROMPT = """You are looking at one foreground element extracted from a PowerPoint slide.

The FIRST image is the ORIGINAL FULL SLIDE (for context).
The SECOND image is the FOREGROUND ELEMENT to approximate using a composition
of NATIVE, EDITABLE PPTX SHAPES.

Class 4 = "an object that can be approximated by a composition of multiple
simple shapes". Decompose into the SMALLEST set of primitive PPTX shapes
that recognizably reproduces it. At most <<MAX_SHAPES>> shapes.

The element's bbox on the slide is (x0, y0, x1, y1) = <<BBOX>> pixels.
Slide size is <<SLIDE_WIDTH>> x <<SLIDE_HEIGHT>> pixels.
The bbox is <<BBOX_W>> x <<BBOX_H>> pixels.

Coordinates are FRACTIONS of the bbox (origin = bbox top-left). For
non-line shapes use (x, y, w, h) in [0, 1] with x+w <= 1, y+h <= 1. For
"line" use (x1, y1, x2, y2). Order back-to-front (first = bottom).

Allowed type values: <<ALLOWED_TYPES>>.

For non-line shapes:
  type, x, y, w, h, rotation_deg (optional),
  fill_color_hex ("#RRGGBB" or "none"),
  line_color_hex ("#RRGGBB" or "none"),
  line_width_pt (number, default 0.75),
  text (optional, may include "\\n"),
  font_size_pt, font_color_hex, bold, italic, h_align, v_align.

For lines:
  type="line", x1, y1, x2, y2, line_color_hex, line_width_pt,
  arrow_start (bool), arrow_end (bool).

STYLE / PALETTE CONTEXT: this deck has been restyled. Pick every
fill_color_hex, line_color_hex, and font_color_hex from the NEW DECK
PALETTE below so the shape composition harmonises with the restyled
deck. Use FONT-palette entries for text colours and THEME-palette
entries for fills / outlines. Do not invent colours outside the
palette unless the original element is a chart / legend whose specific
colours encode data (in that case keep the source colour).

<<PALETTE>>

TEXT COVERAGE REQUIREMENT:
<<TEXT_COVERAGE>>

Reply with STRICT JSON only:
{"shapes": [ ...shape entries... ]}
"""


VERIFY_PROMPT = """You are verifying a PPTX shape-composition reconstruction.

Image A - the ORIGINAL foreground element.
Image B - the CURRENT shape composition rendered from the PPTX (same bbox).

Hard constraint: rendered shapes must NOT extend beyond the bbox. Aim for
a recognisable approximation, not pixel-perfect fidelity.

Current shapes: <<CURRENT_PARAMS>>

If Image B is a faithful, in-bounds reconstruction, reply EXACTLY: {"ok": true}

Otherwise reply with a COMPLETE REPLACEMENT (same schema; top-level
"ok": false and a "shapes" array, at most <<MAX_SHAPES>>). Strict JSON,
no markdown.
Allowed types: <<ALLOWED_TYPES>>.
"""


TEXT_INVENTORY_PROMPT = """You are reading one foreground element extracted from a PowerPoint slide.

Return every visible semantic text label/callout that should survive in a
reconstruction. Include short labels such as DIP1/VIP2 and longer callout
sentences. Ignore purely decorative strokes. For repeated numeric queue
positions, include them only if they are attached to a named label.

Output STRICT JSON only:
{"texts": ["..."]}
"""


def _f(v, default=0.0):
    try:
        return float(v)
    except (TypeError, ValueError):
        return float(default)


def _c01(v):
    return max(0.0, min(1.0, v))


def _match_ocr_line_height(text, ocr_lines):
    """Median h_px of inventory OCR lines whose normalized text matches the
    shape's text (containment either way). None when nothing matches."""
    n = _norm_text(str(text))
    if not n or not ocr_lines:
        return None
    hits = []
    for ln in ocr_lines:
        ln_n = _norm_text(ln.get("text") or "")
        if not ln_n:
            continue
        if n in ln_n or ln_n in n:
            h = ln.get("h_px")
            if h:
                hits.append(float(h))
    if not hits:
        return None
    hits.sort()
    return hits[len(hits) // 2]


def _apply_text(shape, sh, bbox_h_px, h_frac, shape_w_px=None, ocr_lines=None):
    text = sh.get("text")
    if not text:
        try:
            shape.text_frame.text = ""
        except Exception:
            pass
        return
    tf = shape.text_frame
    tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
    # R2: text in a vertically-elongated shape (arrow labels etc.) would
    # wrap at the short axis into 3-4 char stacked lines. Run the line
    # along the LONG axis instead (vert270) with wrapping off; sizing
    # below swaps the width/height budget accordingly.
    shape_h_px_early = max(1.0, h_frac * bbox_h_px)
    vert_layout = bool(
        shape_w_px and len(str(text).replace("\n", " ")) >= 6
        and shape_h_px_early > 2.2 * float(shape_w_px)
    )
    if vert_layout:
        tf.word_wrap = False
        # Pick the vertical direction that reads upright after the shape's
        # own rotation: vert270 runs bottom-to-top, vert top-to-bottom.
        rot_deg = _f(sh.get("rotation_deg"), 0.0) % 360.0
        eff = (rot_deg + 270.0) % 360.0
        vert_val = "vert" if 90.0 < eff < 270.0 else "vert270"
        try:
            tf._txBody.bodyPr.set("vert", vert_val)
        except Exception:
            vert_layout = False
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get((sh.get("v_align") or "middle").lower(), MSO_ANCHOR.MIDDLE)
    halign = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get((sh.get("h_align") or "center").lower(), PP_ALIGN.CENTER)
    lines = str(text).split("\n")
    shape_h_px = max(1.0, h_frac * bbox_h_px)
    # R2: for vert270 layout the line runs along the shape's HEIGHT, so the
    # fit budget swaps axes; single physical line, so join explicit breaks.
    if vert_layout:
        fit_w, fit_h = max(4.0, shape_h_px - 8.0), max(4.0, float(shape_w_px) - 4.0)
        fit_text = str(text).replace("\n", " ")
    else:
        fit_w, fit_h = max(4.0, (shape_w_px or 1.0) - 8.0), max(4.0, shape_h_px - 4.0)
        fit_text = str(text)
    if font_sizing_mode() == "deterministic" and shape_w_px:
        # E1: same canvas-scale bug as class1/2 — size class4 shape text by
        # code. Margins (~4px x, ~2px y at 96dpi) subtracted from the box.
        ink = _match_ocr_line_height(text, ocr_lines)
        pt, _dbg = fit_font_size_pt(
            fit_text, fit_w, fit_h,
            bold=bool(sh.get("bold", False)),
            italic=bool(sh.get("italic", False)),
            ink_h_px=ink,
        )
    else:
        pt = _f(sh.get("font_size_pt"), 12)
        # vert270: the line's vertical extent is the shape WIDTH.
        budget_h_px = float(shape_w_px) if (vert_layout and shape_w_px) else shape_h_px
        max_pt = budget_h_px / DPI * 72.0 / max(1, 1 if vert_layout else len(lines))
        pt = max(MIN_READABLE_PT, min(pt, max_pt))
    color = hex_to_rgb(sh.get("font_color_hex")) or (0, 0, 0)
    if vert_layout:
        lines = [str(text).replace("\n", " ")]
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for i, line in enumerate(lines):
        p = p0 if i == 0 else tf.add_paragraph()
        p.alignment = halign
        run = p.add_run()
        run.text = line
        run.font.size = Pt(pt)
        run.font.bold = bool(sh.get("bold", False))
        run.font.italic = bool(sh.get("italic", False))
        run.font.color.rgb = RGBColor(*color)


def _apply_fill(shape, sh):
    rgb = hex_to_rgb(sh.get("fill_color_hex"))
    if rgb is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*rgb)
    else:
        shape.fill.background()


def _apply_line(line_obj, sh, default_pt=0.75):
    rgb = hex_to_rgb(sh.get("line_color_hex"))
    if rgb is None:
        line_obj.fill.background()
        return
    line_obj.color.rgb = RGBColor(*rgb)
    w = max(0.0, _f(sh.get("line_width_pt"), default_pt))
    if w > 0:
        line_obj.width = Pt(w)


def _add_primitive(slide, sh, bx0, by0, bw, bh, ocr_lines=None):
    typ = str(sh.get("type", "rectangle")).strip().lower()
    if typ == "line":
        x1 = _c01(_f(sh.get("x1"), 0.0))
        y1 = _c01(_f(sh.get("y1"), 0.0))
        x2 = _c01(_f(sh.get("x2"), 1.0))
        y2 = _c01(_f(sh.get("y2"), 1.0))
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(bx0 + x1 * bw), px_to_emu(by0 + y1 * bh),
            px_to_emu(bx0 + x2 * bw), px_to_emu(by0 + y2 * bh),
        )
        _apply_line(conn.line, sh, default_pt=1.0)
        if bool(sh.get("arrow_end", False)) or bool(sh.get("arrow_start", False)):
            try:
                from lxml import etree
                from pptx.oxml.ns import qn
                spPr = conn._element.spPr
                ln = spPr.find(qn("a:ln"))
                if ln is None:
                    ln = etree.SubElement(spPr, qn("a:ln"))
                if bool(sh.get("arrow_start", False)):
                    head = etree.SubElement(ln, qn("a:headEnd"))
                    head.set("type", "triangle")
                if bool(sh.get("arrow_end", False)):
                    tail = etree.SubElement(ln, qn("a:tailEnd"))
                    tail.set("type", "triangle")
            except Exception:
                pass
        return conn
    mso = SHAPE_TYPE_MAP.get(typ, MSO_SHAPE.RECTANGLE)
    x = _c01(_f(sh.get("x"), 0.0))
    y = _c01(_f(sh.get("y"), 0.0))
    w = _c01(_f(sh.get("w"), 0.1))
    h = _c01(_f(sh.get("h"), 0.1))
    if x + w > 1.0:
        w = max(0.001, 1.0 - x)
    if y + h > 1.0:
        h = max(0.001, 1.0 - y)
    w = max(0.001, w); h = max(0.001, h)
    left = px_to_emu(bx0 + x * bw); top = px_to_emu(by0 + y * bh)
    width = px_to_emu(w * bw); height = px_to_emu(h * bh)
    shape = slide.shapes.add_shape(mso, left, top, width, height)
    shape.left, shape.top = left, top
    shape.width, shape.height = width, height
    rot = _f(sh.get("rotation_deg"), 0.0)
    if rot:
        try:
            shape.rotation = float(rot) % 360.0
        except Exception:
            pass
    _apply_fill(shape, sh)
    _apply_line(shape.line, sh, default_pt=0.75)
    _apply_text(shape, sh, bh, h, shape_w_px=w * bw, ocr_lines=ocr_lines)
    return shape


def collect_text_inventory(client, original_path, component_path) -> list[str]:
    content = [
        {"type": "text", "text": "Original full slide for context:"},
        encode_image_b64(original_path),
        {"type": "text", "text": "Foreground element to read:"},
        encode_image_b64(component_path),
        {"type": "text", "text": TEXT_INVENTORY_PROMPT},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=2048)
    parsed = parse_json_loose(raw)
    texts = []
    if isinstance(parsed, dict):
        raw_texts = parsed.get("texts") or []
        if isinstance(raw_texts, list):
            texts = [str(t).strip() for t in raw_texts if str(t).strip()]
    return _important_texts(texts)


def _norm_text(s: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", (s or "").lower())


def _important_texts(texts: list[str]) -> list[str]:
    out = []
    seen = set()
    for t in texts:
        t = re.sub(r"\s+", " ", str(t)).strip()
        n = _norm_text(t)
        if len(n) < MIN_EXPECTED_TEXT_NORM_LEN:
            continue
        if n.isdigit():
            continue
        if n not in seen:
            out.append(t)
            seen.add(n)
    return out


def extract_pptx_texts(pptx_path: Path) -> list[str]:
    texts = []
    if not pptx_path.exists():
        return texts
    try:
        with zipfile.ZipFile(pptx_path) as zf:
            names = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
            for name in names:
                xml = zf.read(name).decode("utf-8", errors="ignore")
                for m in re.finditer(r"<a:t[^>]*>(.*?)</a:t>", xml, flags=re.S):
                    txt = html.unescape(re.sub(r"<[^>]+>", "", m.group(1))).strip()
                    if txt:
                        texts.append(txt)
    except Exception:
        return texts
    return texts


def coverage_missing(expected_texts: list[str], pptx_path: Path) -> list[str]:
    expected = _important_texts(expected_texts)
    if not expected:
        return []
    got_norm = _norm_text(" ".join(extract_pptx_texts(pptx_path)))
    missing = []
    for text in expected:
        n = _norm_text(text)
        if n and n not in got_norm:
            missing.append(text)
    return missing


def text_coverage_message(expected_texts: list[str], missing_texts: list[str] | None = None) -> str:
    expected = _important_texts(expected_texts)
    if not expected:
        return "No reliable text inventory was available; preserve any visible labels/callouts you can read."
    msg = (
        "You MUST include all readable semantic labels/callouts as editable PPTX "
        "text inside the shapes. Required text inventory:\n"
        + "\n".join(f"- {t}" for t in expected)
    )
    if missing_texts:
        msg += (
            "\n\nThe previous attempt was rejected because these required "
            "texts were missing. Add them explicitly:\n"
            + "\n".join(f"- {t}" for t in missing_texts)
        )
    return msg


# ---------------------------------------------------------------------------
# H3: hard coverage gate — deterministic OCR inventory + textbox injection
# ---------------------------------------------------------------------------

def coverage_gate_enabled() -> bool:
    return os.environ.get("COVERAGE_GATE", "").strip().lower() == "hard"


def load_ocr_inventory(page_dir: Path) -> dict:
    """Per-component OCR line inventory (class4_ocr_inventory.py, easyocr
    env). Generated on demand; {} when unavailable."""
    out = page_dir / "class4_ocr_inventory.json"
    if not out.exists():
        try:
            subprocess.run(
                [SLIDECODER_PYTHON, str(PIPELINE_DIR / "class4_ocr_inventory.py"),
                 str(page_dir)],
                check=True, timeout=600,
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
            )
        except Exception as e:
            print(f"  [gate] OCR inventory generation failed: {e}",
                  file=sys.stderr, flush=True)
            return {}
    try:
        return json.loads(out.read_text())
    except Exception:
        return {}


def _sample_ink_color(page_img, bbox) -> str:
    """Dominant ink colour inside a text-line bbox on the original page:
    background = median border colour, ink = median of pixels far from it."""
    x0, y0, x1, y1 = (int(v) for v in bbox)
    crop = page_img.crop((max(0, x0), max(0, y0),
                          min(page_img.width, max(x0 + 1, x1)),
                          min(page_img.height, max(y0 + 1, y1))))
    px = list(crop.getdata())
    if not px:
        return "#000000"
    w, h = crop.size
    border = [px[i] for i in range(len(px))
              if i < w or i >= len(px) - w or i % w == 0 or i % w == w - 1]
    med = lambda vals: tuple(sorted(c[k] for c in vals)[len(vals) // 2] for k in range(3))
    bg = med(border or px)
    ink = [p for p in px if sum(abs(p[k] - bg[k]) for k in range(3)) > 120]
    if not ink:
        # fall back on bg luminance: dark bg -> white text, light bg -> black
        lum = 0.299 * bg[0] + 0.587 * bg[1] + 0.114 * bg[2]
        return "#FFFFFF" if lum < 128 else "#000000"
    r, g, b = med(ink)
    return f"#{r:02X}{g:02X}{b:02X}"


def inject_missing_lines(pptx_path: Path, ocr_lines: list[dict],
                         original_path: Path, stem: str, out_dir: Path):
    """Append every OCR line missing from the recon pptx as an editable
    textbox at its source bbox (H1 deterministic sizing). Returns the
    injected line records."""
    import difflib
    pptx_texts = extract_pptx_texts(pptx_path)
    got_norm = _norm_text(" ".join(pptx_texts))
    got_norm_items = [_norm_text(t) for t in pptx_texts if _norm_text(t)]

    def _covered(n: str) -> bool:
        if n in got_norm:
            return True
        # OCR misreads ("Contraller" for a recon'd "Controller") must not
        # trigger a near-duplicate injection: fuzzy-match against each pptx
        # text and against same-length windows of the concatenation.
        for g in got_norm_items:
            if difflib.SequenceMatcher(None, n, g).ratio() >= 0.8:
                return True
        step = max(1, len(n) // 2)
        for i in range(0, max(1, len(got_norm) - len(n) + 1), step):
            win = got_norm[i:i + len(n)]
            if difflib.SequenceMatcher(None, n, win).ratio() >= 0.85:
                return True
        return False

    missing = []
    for ln in ocr_lines:
        n = _norm_text(ln.get("text") or "")
        if len(n) < MIN_EXPECTED_TEXT_NORM_LEN or n.isdigit():
            continue
        if not _covered(n):
            missing.append(ln)
    if not missing:
        return []
    from PIL import Image
    page_img = Image.open(original_path).convert("RGB")
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    log = []
    for ln in missing:
        x0, y0, x1, y1 = ln["bbox_xyxy"]
        pt, dbg = fit_font_size_pt(ln["text"], x1 - x0, y1 - y0,
                                   ink_h_px=ln.get("h_px"))
        color = _sample_ink_color(page_img, ln["bbox_xyxy"])
        tb = slide.shapes.add_textbox(px_to_emu(x0), px_to_emu(y0),
                                      px_to_emu(x1 - x0), px_to_emu(y1 - y0))
        tb.fill.background()
        tb.line.fill.background()
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = ln["text"]
        run.font.size = Pt(pt)
        run.font.color.rgb = RGBColor(*(hex_to_rgb(color) or (0, 0, 0)))
        log.append({**ln, "font_pt": round(pt, 2), "font_color": color,
                    "sizing": dbg})
    prs.save(str(pptx_path))
    (out_dir / f"{stem}_coverage_injections.json").write_text(
        json.dumps(log, indent=1, ensure_ascii=False))
    return log


def extract_params(client, original_path, component_path, bbox, slide_w, slide_h,
                   palette_summary: str = "", expected_texts: list[str] | None = None,
                   missing_texts: list[str] | None = None):
    bw = bbox[2] - bbox[0]
    bh = bbox[3] - bbox[1]
    prompt = (
        EXTRACT_PROMPT
        .replace("<<BBOX>>", f"({bbox[0]:.1f}, {bbox[1]:.1f}, {bbox[2]:.1f}, {bbox[3]:.1f})")
        .replace("<<SLIDE_WIDTH>>", str(slide_w))
        .replace("<<SLIDE_HEIGHT>>", str(slide_h))
        .replace("<<BBOX_W>>", f"{bw:.1f}")
        .replace("<<BBOX_H>>", f"{bh:.1f}")
        .replace("<<MAX_SHAPES>>", str(MAX_SHAPES))
        .replace("<<ALLOWED_TYPES>>", ALLOWED_TYPES_DOC)
        .replace("<<PALETTE>>", palette_summary or "(no palette)")
        .replace("<<TEXT_COVERAGE>>", text_coverage_message(expected_texts or [], missing_texts))
    )
    content = [
        {"type": "text", "text": "Original slide:"},
        encode_image_b64(original_path),
        {"type": "text", "text": "Foreground element to approximate:"},
        encode_image_b64(component_path),
        {"type": "text", "text": prompt},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=8192)
    return raw, parse_json_loose(raw)


def verify_and_adjust(client, component_path, rendered_crop, params):
    prompt = (
        VERIFY_PROMPT
        .replace("<<CURRENT_PARAMS>>", json.dumps(params, ensure_ascii=False))
        .replace("<<MAX_SHAPES>>", str(MAX_SHAPES))
        .replace("<<ALLOWED_TYPES>>", ALLOWED_TYPES_DOC)
    )
    content = [
        {"type": "text", "text": "Image A - original:"},
        encode_image_b64(component_path),
        {"type": "text", "text": "Image B - current rendered shapes:"},
        encode_image_b64(rendered_crop),
        {"type": "text", "text": prompt},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=8192)
    return raw, parse_json_loose(raw)


def build_pptx(slide_w_px, slide_h_px, bbox, params, out_path, ocr_lines=None):
    x0, y0, x1, y1 = bbox
    x0 = max(0.0, min(x0, slide_w_px))
    y0 = max(0.0, min(y0, slide_h_px))
    x1 = max(0.0, min(x1, slide_w_px))
    y1 = max(0.0, min(y1, slide_h_px))
    if x1 <= x0 or y1 <= y0:
        raise ValueError(f"degenerate bbox: {bbox}")
    bw = x1 - x0; bh = y1 - y0
    prs = Presentation()
    prs.slide_width = px_to_emu(slide_w_px)
    prs.slide_height = px_to_emu(slide_h_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    raw_shapes = params.get("shapes") or []
    if not isinstance(raw_shapes, list):
        raw_shapes = []
    shapes = raw_shapes[:MAX_SHAPES]
    n_added = 0; n_failed = 0
    for sh in shapes:
        if not isinstance(sh, dict):
            n_failed += 1
            continue
        try:
            _add_primitive(slide, sh, x0, y0, bw, bh, ocr_lines=ocr_lines)
            n_added += 1
        except Exception as e:
            print(f"      [warn] shape failed: {e}")
            n_failed += 1
    prs.save(str(out_path))
    return {"n_added": n_added, "n_failed": n_failed, "n_requested": len(shapes)}


def _validate(p):
    return isinstance(p, dict) and isinstance(p.get("shapes"), list) and p.get("shapes")


def _default_params():
    return {"shapes": [
        {"type": "rectangle", "x": 0.0, "y": 0.0, "w": 1.0, "h": 1.0,
         "fill_color_hex": "none", "line_color_hex": "#888888",
         "line_width_pt": 0.5}
    ]}


def build_raster_fallback_pptx(slide_w_px, slide_h_px, bbox, component_path: Path,
                               out_path: Path):
    x0, y0, x1, y1 = bbox
    prs = Presentation()
    prs.slide_width = px_to_emu(slide_w_px)
    prs.slide_height = px_to_emu(slide_h_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        str(component_path),
        px_to_emu(x0), px_to_emu(y0),
        width=px_to_emu(x1 - x0), height=px_to_emu(y1 - y0),
    )
    prs.save(str(out_path))
    return {"n_added": 1, "n_failed": 0, "n_requested": 1, "raster_fallback": True}


def process_component(client, page_dir, slide_w, slide_h, original_path,
                      comp_meta, comp_entry, out_root,
                      palette_summary: str = ""):
    page_name = page_dir.name
    cf = comp_entry["component_file"]
    bbox = comp_meta["bbox_xyxy"]
    stem = Path(cf).stem
    out_dir = out_root / page_name
    out_dir.mkdir(parents=True, exist_ok=True)
    component_path = page_dir / cf
    pptx_path = out_dir / f"{stem}.pptx"
    if pptx_path.exists():
        print(f"  [skip] {pptx_path.name} exists")
        # The hard gate is idempotent (already-covered lines inject nothing),
        # so existing recons can be gated without redoing the VLM recon.
        if coverage_gate_enabled():
            try:
                inv = (load_ocr_inventory(page_dir).get(cf) or {}).get("lines", [])
                injected = inject_missing_lines(pptx_path, inv, original_path,
                                                stem, out_dir)
                if injected:
                    print(f"  [gate] injected {len(injected)} missing OCR line(s) "
                          f"into existing {pptx_path.name}")
                    try:
                        final = render_pptx_to_png(pptx_path, slide_w, slide_h)
                        crop_bbox(final, comp_meta["bbox_xyxy"],
                                  out_dir / f"{stem}_rendered_crop.png")
                    except Exception:
                        pass
            except Exception as e:
                print(f"  [gate] injection failed: {e}")
        return pptx_path
    print(f"--- [{page_name}] {stem} ---")
    inventory_path = out_dir / f"{stem}_text_inventory.json"
    expected_texts = []
    if inventory_path.exists():
        try:
            expected_texts = _important_texts(json.loads(inventory_path.read_text()).get("texts") or [])
        except Exception:
            expected_texts = []
    else:
        try:
            expected_texts = collect_text_inventory(client, original_path, component_path)
            inventory_path.write_text(json.dumps({"texts": expected_texts}, indent=2, ensure_ascii=False))
            if expected_texts:
                print(f"  text inventory: {len(expected_texts)} item(s)")
        except Exception as e:
            inventory_path.write_text(json.dumps({"texts": [], "error": str(e)}, indent=2, ensure_ascii=False))
            print(f"  [warn] text inventory failed: {e}")
    gate_lines = []
    if coverage_gate_enabled() or font_sizing_mode() == "deterministic":
        gate_lines = (load_ocr_inventory(page_dir).get(cf) or {}).get("lines", [])
    missing_texts = []
    raw, params = extract_params(
        client, original_path, component_path, bbox, slide_w, slide_h,
        palette_summary, expected_texts,
    )
    if not _validate(params):
        print(f"  [warn] extract failed; defaults. raw: {raw[:200]!r}")
        params = _default_params()
    summary = None
    for r in range(MAX_VERIFY_ROUNDS + 1):
        summary = build_pptx(slide_w, slide_h, bbox, params, pptx_path,
                             ocr_lines=gate_lines)
        if r == MAX_VERIFY_ROUNDS:
            print(f"  [round {r}] max rounds")
            break
        missing_texts = coverage_missing(expected_texts, pptx_path)
        if len(missing_texts) > MAX_COVERAGE_MISSING:
            print(f"  [round {r}] text coverage missing {len(missing_texts)}; retrying")
            raw_retry, retry_params = extract_params(
                client, original_path, component_path, bbox, slide_w, slide_h,
                palette_summary, expected_texts, missing_texts,
            )
            if _validate(retry_params):
                params = retry_params
                continue
            print(f"  [round {r}] coverage retry invalid; continuing verifier")
        try:
            rendered = render_pptx_to_png(pptx_path, slide_w, slide_h)
        except (subprocess.TimeoutExpired, Exception) as e:
            print(f"  [warn] render failed: {e}")
            break
        crop = out_dir / f"{stem}_rendered_crop_r{r}.png"
        crop_bbox(rendered, bbox, crop)
        raw_v, adj = verify_and_adjust(client, component_path, crop, params)
        if not isinstance(adj, dict):
            print(f"  [round {r}] verify parse failed; stopping")
            break
        if adj.get("ok") is True:
            print(f"  [round {r}] OK")
            break
        if not _validate(adj):
            print(f"  [round {r}] invalid replacement; stopping")
            break
        params = adj
    missing_texts = coverage_missing(expected_texts, pptx_path)
    if len(missing_texts) > MAX_COVERAGE_MISSING:
        fallback_path = out_dir / f"{stem}_shape_failed.bak"
        try:
            shutil.copyfile(pptx_path, fallback_path)
        except Exception:
            pass
        summary = build_raster_fallback_pptx(slide_w, slide_h, bbox, component_path, pptx_path)
        (out_dir / f"{stem}_coverage_failure.json").write_text(json.dumps({
            "expected_texts": expected_texts,
            "missing_texts": missing_texts,
            "fallback": "raster_preserve",
            "failed_shape_pptx": str(fallback_path),
        }, indent=2, ensure_ascii=False))
        print(f"  [coverage] fallback to raster preserve; missing={missing_texts}")
    if coverage_gate_enabled() and not (summary or {}).get("raster_fallback"):
        try:
            injected = inject_missing_lines(pptx_path, gate_lines, original_path,
                                            stem, out_dir)
            if injected:
                print(f"  [gate] injected {len(injected)} missing OCR line(s): "
                      + "; ".join(ln['text'][:30] for ln in injected[:5]))
        except Exception as e:
            print(f"  [gate] injection failed: {e}")
    # final render preview
    try:
        final = render_pptx_to_png(pptx_path, slide_w, slide_h)
        crop_bbox(final, bbox, out_dir / f"{stem}_rendered_crop.png")
    except Exception:
        pass
    print(f"  saved: {pptx_path}  (added={summary['n_added']} "
          f"failed={summary['n_failed']})")
    return pptx_path


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--case", required=True)
    ap.add_argument("--out-dir", default=None,
                    help="override per-case output root (default: output/<case>/)")
    args = ap.parse_args()
    paths = case_paths(args.case, out_dir=args.out_dir)
    paths.step3_class4_dir.mkdir(parents=True, exist_ok=True)
    client = anthropic_client()
    # F-08: load the deck palette so shape colours pick from the palette.
    from step3_class5 import _load_palette, _palette_summary
    palette = _load_palette(paths)
    palette_summary = _palette_summary(palette)
    pages = list_front_bg_pages(paths)
    if not pages:
        print(f"no front_bg pages under {paths.front_bg_dir}", file=sys.stderr)
        return 1
    for page in pages:
        cat_path = page / "categorization.json"
        meta_path = page / "metadata.json"
        if not cat_path.exists() or not meta_path.exists():
            continue
        cat = json.loads(cat_path.read_text())
        meta = json.loads(meta_path.read_text())
        slide_w = int(meta["query_image_size"]["width"])
        slide_h = int(meta["query_image_size"]["height"])
        original_path = Path(meta["query_image_path"])
        bbox_lookup = {c["component_file"]: c for c in meta["components"]}
        targets = [r for r in cat["results"] if r.get("class") == TARGET_CLASS]
        if not targets:
            continue
        print(f"\n=== Page {page.name}: {len(targets)} class-4 ===")
        for r in targets:
            cm = bbox_lookup.get(r["component_file"])
            if cm is None:
                continue
            try:
                process_component(client, page, slide_w, slide_h, original_path,
                                  cm, r, paths.step3_class4_dir,
                                  palette_summary)
            except Exception as e:
                print(f"  [error] {r['component_file']}: {e}")
    return 0


if __name__ == "__main__":
    sys.exit(main())

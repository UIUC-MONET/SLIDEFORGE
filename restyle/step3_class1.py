"""Step 3 (class 1): convert every class-1 foreground element into an
editable PPTX textbox, with a self-verify-and-adjust loop driven by Claude.

For each class-1 component:
  1. Claude proposes a textbox parameterization (text, font size/colour,
     alignment, fill, line count).
  2. We build a one-shape PPTX whose slide size matches the source slide
     and place the textbox at the element's bbox.
  3. Render the PPTX via headless LibreOffice, crop the bbox, and ask
     Claude to verify: (a) text stays inside the bbox, (b) the rendered
     line count matches the original visible line count. If the line
     count is too high (text wrapped), prefer EXTENDING the bbox
     horizontally (within slide canvas) over shrinking the font.
  4. Save the final ``<page>/<stem>.pptx`` under step3_class1_dir.
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from common import (
    CLAUDE_SONNET,
    DPI,
    MIN_READABLE_PT,
    anthropic_client,
    call_claude,
    case_paths,
    crop_bbox,
    encode_image_b64,
    fit_font_size_pt,
    carlito_metrics_unreliable,
    font_sizing_mode,
    hex_to_rgb,
    list_front_bg_pages,
    load_text_metrics,
    parse_json_loose,
    px_to_emu,
    render_pptx_to_png,
)


TARGET_CLASS = 1
MAX_VERIFY_ROUNDS = 3
MODEL = CLAUDE_SONNET


EXTRACT_PROMPT = """You are looking at one foreground element extracted from a PowerPoint slide.

The FIRST image is the ORIGINAL FULL SLIDE (for context).
The SECOND image is the FOREGROUND ELEMENT to encode as an EDITABLE TEXT BOX.

This element was classified as "class 1 - pure text or formula WITHOUT any
background colour or shape". Your job is to reproduce it faithfully as a
single textbox that fits EXACTLY inside its bounding box. Decide:

- The verbatim text content of the element. Preserve casing, punctuation,
  spaces, and any line breaks (use "\\n" to encode line breaks).
- The TOTAL number of visible text lines in the original element.
- Approximate point size: estimate from visible glyph height relative to
  the bbox HEIGHT in pixels we give below. For multi-line text, estimate
  the size of ONE line.
- Dominant text colour (hex).
- Whether the text looks bold and/or italic.
- Horizontal alignment ("left"|"center"|"right").
- Vertical alignment ("top"|"middle"|"bottom").
- Fill / background colour ("none" for class 1 unless you see one).

The element's bbox on the slide is (x0, y0, x1, y1) = <<BBOX>> pixels.
Slide size is <<SLIDE_WIDTH>> x <<SLIDE_HEIGHT>> pixels.
The bbox is <<BBOX_W>> x <<BBOX_H>> pixels.

Reply with STRICT JSON only - no markdown, no commentary:
{
  "text": "<verbatim text, may include \\n>",
  "num_lines": <integer>,
  "font_size_pt": <number>,
  "font_color_hex": "#RRGGBB",
  "bold": <true|false>,
  "italic": <true|false>,
  "h_align": "left|center|right",
  "v_align": "top|middle|bottom",
  "fill_color_hex": "#RRGGBB or none"
}

HARD RULE: the "text" field MUST be a non-empty verbatim transcription
of what you actually see in the element. Returning "" is a hard failure.
If you genuinely cannot read any text in the element, return exactly
{"text": "[unreadable]", "num_lines": 1} so the orchestrator can log it
for review.
"""


VERIFY_PROMPT = """You are verifying a PPTX textbox reconstruction.

Image A - the ORIGINAL foreground element (ground truth).
Image B - the CURRENT textbox rendered from the PPTX, cropped to the same
region as the current (possibly extended) textbox bbox.

Constraints, in priority order:
  1. (HARD) The rendered text must NOT extend beyond the textbox bbox.
  2. (HARD) The rendered font size must NOT go below <<MIN_PT>> pt.
  3. (SOFT) Aim to match the original line count <<ORIG_NUM_LINES>>.
     If matching it would require dropping the font below <<MIN_PT>> pt
     (constraint 2), accept +/-1 line of soft-wrap instead. Legibility
     wins over exact line count.

If Image B has MORE lines than the original (the textbox is too narrow,
so text soft-wrapped), the PREFERRED fix is to EXTEND the textbox bbox
horizontally - request via "extend_width_px" (positive integer = extra
pixels to add to bbox WIDTH; the orchestrator distributes them across
sides that have slack).

Hard cap: the resulting effective bbox MUST stay inside the slide canvas
(<<SLIDE_W>> x <<SLIDE_H>> pixels). Current bbox = <<CURRENT_BBOX>>;
slack: left=<<SLACK_LEFT>>, right=<<SLACK_RIGHT>>, total=<<SLACK_TOTAL>>.
Never request extend_width_px larger than <<SLACK_TOTAL>>. If total
slack is 0, keep "font_size_pt" at or above <<MIN_PT>> and accept the
smallest readable soft-wrap difference.

Prefer reducing font size only down to <<MIN_PT>> over extending the
bbox when the element is near other text or when the bbox is already
large. Only request extend_width_px for isolated short text labels.

Current textbox parameters: <<CURRENT_PARAMS>>

If Image B already (a) stays inside the bbox AND (b) has either the same
line count or the smallest readable soft-wrap difference allowed above
AND (c) reproduces the text, reply EXACTLY: {"ok": true}

Otherwise reply with adjustments using this schema:
{
  "ok": false,
  "observed_lines_original": <int>,
  "observed_lines_rendered": <int>,
  "extend_width_px": <int, 0 if none>,
  "text": "...",
  "font_size_pt": <number>,
  "font_color_hex": "#RRGGBB",
  "bold": <bool>, "italic": <bool>,
  "h_align": "left|center|right",
  "v_align": "top|middle|bottom",
  "fill_color_hex": "#RRGGBB or none"
}

Only include keys you actually want to change. Strict JSON, no markdown."""


DEFAULT_PARAMS = {
    "text": "", "num_lines": 1, "font_size_pt": 12,
    "font_color_hex": "#000000",
    "bold": False, "italic": False,
    "h_align": "left", "v_align": "middle",
    "fill_color_hex": "none",
}

ADJUSTABLE_KEYS = (
    "text", "font_size_pt", "font_color_hex", "bold", "italic",
    "h_align", "v_align", "fill_color_hex",
)


OCR_UNDEREXTRACT_MIN_CHARS = 120
OCR_UNDEREXTRACT_RATIO = 0.72
SPLIT_REGION_MIN_AREA_FRAC = 0.08


SPLIT_PROMPT = """You are checking whether one extracted text element is actually
multiple spatial text regions, such as two columns, separated bullet groups,
or multiple text blocks that should not be flattened into one textbox.

The FIRST image is the original full slide. The SECOND image is the extracted
foreground text element. Its bbox is <<BBOX>> on a <<SLIDE_W>>x<<SLIDE_H>>
slide. OCR text, if available, is:

<<OCR_TEXT>>

If the element is a single text block, return {"use_single": true}.

If it has multiple spatial regions, return their positions as fractions of
the element bbox. Preserve the text per region and keep reading order.

Output STRICT JSON only:
{
  "use_single": <true|false>,
  "regions": [
    {
      "x": <0..1>, "y": <0..1>, "w": <0..1>, "h": <0..1>,
      "text": "<verbatim text for this region, may contain \\n>",
      "font_size_pt": <number>,
      "font_color_hex": "#RRGGBB",
      "bold": <true|false>,
      "italic": <true|false>,
      "h_align": "left|center|right",
      "v_align": "top|middle|bottom"
    }
  ],
  "reason": "<short sentence>"
}
"""


def _is_valid_extract(params) -> bool:
    return (
        isinstance(params, dict)
        and isinstance(params.get("text"), str)
        and params["text"].strip() not in ("", "[unreadable]")
    )


def _latex_ocr_to_text(raw: str | None) -> str:
    if not raw:
        return ""
    s = str(raw)
    s = s.replace("\\\\", "\n")
    s = s.replace("\\left", "").replace("\\right", "")
    s = s.replace("\\cdot", "*").replace("\\times", "*")
    s = s.replace("\\leq", "<=").replace("\\geq", ">=")
    s = s.replace("\\neq", "!=").replace("\\approx", "~")
    s = s.replace("\\infty", "infinity")
    for _ in range(8):
        ns = re.sub(r"\\frac\{([^{}]+)\}\{([^{}]+)\}", r"(\1)/(\2)", s)
        if ns == s:
            break
        s = ns
    s = re.sub(r"_\{([^{}]+)\}", r"_\1", s)
    s = re.sub(r"\^\{([^{}]+)\}", r"^(\1)", s)
    s = re.sub(r"_([A-Za-z0-9])", r"_\1", s)
    s = re.sub(r"\^([A-Za-z0-9])", r"^(\1)", s)
    s = s.replace("\\(", "").replace("\\)", "")
    s = s.replace("\\[", "").replace("\\]", "")
    s = s.replace("$", "")
    s = re.sub(r"\\begin\{itemize\}|\\end\{itemize\}", "\n", s)
    s = re.sub(r"\\item\[[^\]]+\]", "\n- ", s)
    s = re.sub(r"\\item\b", "\n• ", s)
    # Keep text inside the common formatting commands emitted by OCR.
    for _ in range(4):
        s = re.sub(r"\\textcolor\{[^{}]*\}\{([^{}]*)\}", r"\1", s)
        s = re.sub(
            r"\\(?:textbf|textit|emph|underline|textsuperscript)\{([^{}]*)\}",
            r"\1", s,
        )
    s = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?(?:\{[^{}]*\})?", " ", s)
    s = s.replace("{", "").replace("}", "").replace("\\", "")
    lines = [re.sub(r"\s+", " ", line).strip() for line in s.splitlines()]
    return "\n".join(line for line in lines if line)


def _norm_text_len(text: str) -> int:
    return len(re.sub(r"\W+", "", (text or "").lower()))


def _bullet_count(text: str) -> int:
    return len(re.findall(r"(?m)^\s*(?:[•\-]|\d+[.)])\s+", text or ""))


def _looks_multi_region_candidate(bbox, slide_w, slide_h, ocr_text: str) -> bool:
    x0, y0, x1, y1 = bbox
    bw = max(1.0, x1 - x0)
    bh = max(1.0, y1 - y0)
    area_frac = (bw * bh) / max(1.0, slide_w * slide_h)
    wide = bw >= 0.42 * slide_w and bw / bh >= 1.25
    many_bullets = _bullet_count(ocr_text) >= 4
    many_lines = len([ln for ln in (ocr_text or "").splitlines() if ln.strip()]) >= 6
    return area_frac >= SPLIT_REGION_MIN_AREA_FRAC and (wide or many_bullets or many_lines)


def _valid_split_regions(parsed) -> list[dict]:
    if not isinstance(parsed, dict) or parsed.get("use_single") is True:
        return []
    regs = parsed.get("regions")
    if not isinstance(regs, list) or len(regs) < 2:
        return []
    out = []
    for r in regs:
        if not isinstance(r, dict) or not str(r.get("text") or "").strip():
            continue
        try:
            x = max(0.0, min(1.0, float(r.get("x", 0))))
            y = max(0.0, min(1.0, float(r.get("y", 0))))
            w = max(0.01, min(1.0 - x, float(r.get("w", 0.1))))
            h = max(0.01, min(1.0 - y, float(r.get("h", 0.1))))
        except (TypeError, ValueError):
            continue
        rr = dict(DEFAULT_PARAMS)
        rr.update(r)
        rr.update({"x": x, "y": y, "w": w, "h": h})
        out.append(rr)
    return out if len(out) >= 2 else []


def detect_split_regions(client, original_path, component_path, bbox,
                         slide_w, slide_h, ocr_text: str) -> list[dict]:
    prompt = (
        SPLIT_PROMPT
        .replace("<<BBOX>>", f"({bbox[0]:.1f}, {bbox[1]:.1f}, {bbox[2]:.1f}, {bbox[3]:.1f})")
        .replace("<<SLIDE_W>>", str(slide_w))
        .replace("<<SLIDE_H>>", str(slide_h))
        .replace("<<OCR_TEXT>>", ocr_text[:3000] if ocr_text else "(none)")
    )
    content = [
        {"type": "text", "text": "Original slide (full image):"},
        encode_image_b64(original_path),
        {"type": "text", "text": "Foreground text element:"},
        encode_image_b64(component_path),
        {"type": "text", "text": prompt},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=4096)
    parsed = parse_json_loose(raw)
    return _valid_split_regions(parsed)


def _looks_under_extracted(params, ocr_text: str) -> bool:
    if not _is_valid_extract(params) or not ocr_text.strip():
        return False
    got = params.get("text") or ""
    ocr_len = _norm_text_len(ocr_text)
    got_len = _norm_text_len(got)
    if ocr_len >= OCR_UNDEREXTRACT_MIN_CHARS and got_len < OCR_UNDEREXTRACT_RATIO * ocr_len:
        return True
    ocr_bullets = _bullet_count(ocr_text)
    got_bullets = _bullet_count(got)
    return ocr_bullets >= 4 and got_bullets <= ocr_bullets - 2


def _ocr_fallback_font_size(ocr_text: str, bbox) -> float:
    lines = [line for line in ocr_text.splitlines() if line.strip()]
    n_lines = max(1, len(lines))
    max_line_chars = max((len(line) for line in lines), default=0)
    bbox_h_pt = max(1.0, (bbox[3] - bbox[1]) / DPI * 72.0)
    by_height = bbox_h_pt / n_lines * 0.70
    if max_line_chars > 150:
        cap = 14.0
    elif max_line_chars > 100:
        cap = 18.0
    elif max_line_chars > 70:
        cap = 28.0
    else:
        cap = 34.0
    return max(MIN_READABLE_PT, min(cap, by_height))


def _apply_ocr_text(params, ocr_text: str, bbox) -> dict:
    fixed = dict(DEFAULT_PARAMS)
    if isinstance(params, dict):
        fixed.update(params)
    fixed["text"] = ocr_text
    fixed["num_lines"] = max(1, len(ocr_text.splitlines()))
    fixed["font_size_pt"] = _ocr_fallback_font_size(ocr_text, bbox)
    return fixed


def extract_params(client, original_path, component_path, bbox, slide_w, slide_h):
    bw = bbox[2] - bbox[0]
    bh = bbox[3] - bbox[1]
    prompt = (
        EXTRACT_PROMPT
        .replace("<<BBOX>>", f"({bbox[0]:.1f}, {bbox[1]:.1f}, {bbox[2]:.1f}, {bbox[3]:.1f})")
        .replace("<<SLIDE_WIDTH>>", str(slide_w))
        .replace("<<SLIDE_HEIGHT>>", str(slide_h))
        .replace("<<BBOX_W>>", f"{bw:.1f}")
        .replace("<<BBOX_H>>", f"{bh:.1f}")
    )
    content = [
        {"type": "text", "text": "Original slide (full image):"},
        encode_image_b64(original_path),
        {"type": "text", "text": "Foreground element to encode as a textbox:"},
        encode_image_b64(component_path),
        {"type": "text", "text": prompt},
    ]
    max_tokens = 4096 if bw * bh > 0.12 * slide_w * slide_h else 1024
    raw = call_claude(client, content, model=MODEL, max_tokens=max_tokens)
    return raw, parse_json_loose(raw)


def verify_and_adjust(client, component_path, rendered_crop, params,
                      target_lines, eff_bbox, slide_w, slide_h):
    x0, y0, x1, y1 = eff_bbox
    sl = max(0.0, x0)
    sr = max(0.0, slide_w - x1)
    prompt = (
        VERIFY_PROMPT
        .replace("<<CURRENT_PARAMS>>", json.dumps(
            {k: v for k, v in params.items() if not k.startswith("_")},
            ensure_ascii=False))
        .replace("<<ORIG_NUM_LINES>>", str(target_lines))
        .replace("<<MIN_PT>>", f"{MIN_READABLE_PT:.0f}")
        .replace("<<SLIDE_W>>", str(slide_w))
        .replace("<<SLIDE_H>>", str(slide_h))
        .replace("<<CURRENT_BBOX>>",
                 f"({x0:.1f}, {y0:.1f}, {x1:.1f}, {y1:.1f})")
        .replace("<<SLACK_LEFT>>", f"{sl:.0f}")
        .replace("<<SLACK_RIGHT>>", f"{sr:.0f}")
        .replace("<<SLACK_TOTAL>>", f"{sl + sr:.0f}")
    )
    content = [
        {"type": "text", "text": "Image A - original foreground element:"},
        encode_image_b64(component_path),
        {"type": "text", "text": "Image B - current rendered textbox (cropped):"},
        encode_image_b64(rendered_crop),
        {"type": "text", "text": prompt},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=512)
    return raw, parse_json_loose(raw)


def extend_bbox_horizontally(bbox, extend_px, slide_w):
    if extend_px <= 0:
        return bbox
    x0, y0, x1, y1 = bbox
    sl = max(0.0, x0)
    sr = max(0.0, slide_w - x1)
    if sl + sr <= 0:
        return bbox
    half = extend_px / 2.0
    add_l = min(half, sl)
    add_r = min(half, sr)
    rem = extend_px - (add_l + add_r)
    if rem > 0:
        extra_l = sl - add_l
        extra_r = sr - add_r
        if extra_l > 0:
            add = min(rem, extra_l)
            add_l += add
            rem -= add
        if rem > 0 and extra_r > 0:
            add_r += min(rem, extra_r)
    return (x0 - add_l, y0, x1 + add_r, y1)

TEXT_OVERLAP_THRESH = 0.12
TEXT_OVERLAP_MARGIN_PX = 10


def bbox_area(b):
    x0, y0, x1, y1 = b
    return max(0.0, x1 - x0) * max(0.0, y1 - y0)


def bbox_intersection_area(a, b):
    x0 = max(a[0], b[0])
    y0 = max(a[1], b[1])
    x1 = min(a[2], b[2])
    y1 = min(a[3], b[3])
    return max(0.0, x1 - x0) * max(0.0, y1 - y0)


def bbox_overlap_ratios(a, b):
    inter = bbox_intersection_area(a, b)
    if inter <= 0:
        return 0.0, 0.0
    return inter / max(1.0, bbox_area(a)), inter / max(1.0, bbox_area(b))


def resolve_text_bbox_overlap(bbox, placed_bboxes, slide_h):
    """
    If this text bbox overlaps previously accepted text bboxes on the same page,
    move it just below/above the conflicting bbox. Normal non-overlapping pages
    are unchanged.
    """
    cur = tuple(bbox)

    for _ in range(5):
        changed = False

        for prev in placed_bboxes:
            ratio_cur, ratio_prev = bbox_overlap_ratios(cur, prev)
            if ratio_cur <= TEXT_OVERLAP_THRESH and ratio_prev <= TEXT_OVERLAP_THRESH:
                continue

            x0, y0, x1, y1 = cur
            h = y1 - y0
            cur_cy = (y0 + y1) / 2.0
            prev_cy = (prev[1] + prev[3]) / 2.0

            # If current text is visually below previous text, push it down.
            if cur_cy >= prev_cy:
                new_y0 = prev[3] + TEXT_OVERLAP_MARGIN_PX
                new_y1 = new_y0 + h

                if new_y1 <= slide_h:
                    cur = (x0, new_y0, x1, new_y1)
                    changed = True
                    continue

                # If no room below, try moving above.
                new_y1 = prev[1] - TEXT_OVERLAP_MARGIN_PX
                new_y0 = new_y1 - h
                if new_y0 >= 0:
                    cur = (x0, new_y0, x1, new_y1)
                    changed = True
                    continue

            # If current text is visually above previous text, push it up.
            else:
                new_y1 = prev[1] - TEXT_OVERLAP_MARGIN_PX
                new_y0 = new_y1 - h

                if new_y0 >= 0:
                    cur = (x0, new_y0, x1, new_y1)
                    changed = True
                    continue

                # If no room above, try moving below.
                new_y0 = prev[3] + TEXT_OVERLAP_MARGIN_PX
                new_y1 = new_y0 + h
                if new_y1 <= slide_h:
                    cur = (x0, new_y0, x1, new_y1)
                    changed = True
                    continue

        if not changed:
            break

    return cur


def _add_textbox(slide, bbox, params):
    x0, y0, x1, y1 = bbox
    tb = slide.shapes.add_textbox(
        px_to_emu(x0), px_to_emu(y0),
        px_to_emu(x1 - x0), px_to_emu(y1 - y0),
    )
    fill_rgb = hex_to_rgb(params.get("fill_color_hex"))
    if fill_rgb is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor(*fill_rgb)
    else:
        tb.fill.background()
    tb.line.fill.background()

    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.auto_size = MSO_AUTO_SIZE.NONE

    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get((params.get("v_align") or "middle").lower(), MSO_ANCHOR.MIDDLE)
    halign = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get((params.get("h_align") or "left").lower(), PP_ALIGN.LEFT)

    text = params.get("text") or ""
    lines = text.split("\n") if text else [""]
    per_line_pt = None
    if font_sizing_mode() == "deterministic":
        # H1+H2: size from code (ink prior + max-fit), not VLM eyeballing.
        bold = bool(params.get("bold", False))
        italic = bool(params.get("italic", False))
        boxes_y_h = params.get("_ink_boxes_y_h") or []
        heights = [h for _y, h in boxes_y_h]
        # Multi-scale component (e.g. title + presenter line in one crop):
        # uniform sizing would inflate the small line to the big line's size.
        # When OCR lines map 1:1 to explicit text lines and heights differ
        # >1.35x, size each line by its own ink height.
        hetero = (len(lines) >= 2 and len(heights) == len(lines)
                  and min(heights) > 0 and max(heights) / min(heights) > 1.35)
        if hetero:
            per_line_pt = []
            for line, (_by, bh) in zip(lines, boxes_y_h):
                lpt, _ = fit_font_size_pt(
                    line or " ", x1 - x0, bh * 1.6,
                    bold=bold, italic=italic, ink_h_px=bh,
                )
                per_line_pt.append(lpt)
            # keep total stacked height within the bbox
            total_px = sum(p * DPI / 72.0 * 1.22 for p in per_line_pt)
            if total_px > (y1 - y0):
                shrink = (y1 - y0) / total_px
                per_line_pt = [max(MIN_READABLE_PT, p * shrink)
                               for p in per_line_pt]
            pt_size = min(per_line_pt)
            params["_sizing_debug"] = {
                "per_line_pt": [round(p, 2) for p in per_line_pt],
                "ink_heights": heights,
                "vlm_pt": params.get("font_size_pt"),
            }
        else:
            pt_size, dbg = fit_font_size_pt(
                text, x1 - x0, y1 - y0,
                bold=bold, italic=italic,
                ink_h_px=params.get("_ink_h_px"),
            )
            dbg["vlm_pt"] = params.get("font_size_pt")
            params["_sizing_debug"] = dbg
    else:
        pt_size = float(params.get("font_size_pt") or 12)
        bbox_h_pt = (y1 - y0) / DPI * 72.0
        max_pt = bbox_h_pt / max(1, len(lines))
        pt_size = max(MIN_READABLE_PT, min(pt_size, max_pt))
    color = hex_to_rgb(params.get("font_color_hex")) or (0, 0, 0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = halign
        run = p.add_run()
        run.text = line
        run.font.size = Pt(per_line_pt[i] if per_line_pt else pt_size)
        run.font.bold = bool(params.get("bold", False))
        run.font.italic = bool(params.get("italic", False))
        run.font.color.rgb = RGBColor(*color)

    return pt_size


def build_pptx(slide_w_px, slide_h_px, bbox, params, output_path):
    x0, y0, x1, y1 = bbox
    x0 = max(0.0, min(x0, slide_w_px))
    y0 = max(0.0, min(y0, slide_h_px))
    x1 = max(0.0, min(x1, slide_w_px))
    y1 = max(0.0, min(y1, slide_h_px))
    if x1 <= x0 or y1 <= y0:
        raise ValueError(f"degenerate bbox: {bbox}")

    prs = Presentation()
    prs.slide_width = px_to_emu(slide_w_px)
    prs.slide_height = px_to_emu(slide_h_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    regions = params.get("regions") if isinstance(params, dict) else None
    if isinstance(regions, list) and regions:
        bw = x1 - x0
        bh = y1 - y0
        used = []
        ink_boxes = params.get("_ink_boxes_y_h") or []
        for r in regions:
            rx0 = x0 + float(r.get("x", 0)) * bw
            ry0 = y0 + float(r.get("y", 0)) * bh
            rx1 = rx0 + float(r.get("w", 1)) * bw
            ry1 = ry0 + float(r.get("h", 1)) * bh
            # per-region ink prior: OCR line boxes whose vertical center
            # falls inside the region (regions stack vertically). Without
            # this, region maxfit inflates text to the VLM-drawn box.
            hits = sorted(h for by, h in ink_boxes
                          if ry0 <= by + h / 2.0 <= ry1)
            if hits:
                r["_ink_h_px"] = hits[len(hits) // 2]
            used.append(_add_textbox(slide, (rx0, ry0, rx1, ry1), r))
        prs.save(str(output_path))
        return min(used) if used else None

    pt_size = _add_textbox(slide, (x0, y0, x1, y1), params)
    prs.save(str(output_path))
    return pt_size


def build_raster_fallback_pptx(slide_w_px, slide_h_px, bbox,
                               component_path, output_path):
    """Preserve the original component as a raster when text extraction
    fails. This sacrifices editability but prevents blank slide bodies."""
    x0, y0, x1, y1 = bbox
    prs = Presentation()
    prs.slide_width = px_to_emu(slide_w_px)
    prs.slide_height = px_to_emu(slide_h_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        str(component_path),
        px_to_emu(x0), px_to_emu(y0),
        width=px_to_emu(x1 - x0), height=px_to_emu(y1 - y0),
    )
    prs.save(str(output_path))


def process_component(client, page_dir, slide_w, slide_h, original_path,
                      comp_meta, comp_entry, out_root, placed_text_bboxes=None,
                      ink_h_px=None, ink_boxes_y_h=None):
    page_name = page_dir.name
    cf = comp_entry["component_file"]
    orig_bbox = tuple(comp_meta["bbox_xyxy"])
    stem = Path(cf).stem
    out_dir = out_root / page_name
    out_dir.mkdir(parents=True, exist_ok=True)
    component_path = page_dir / cf

    pptx_path = out_dir / f"{stem}.pptx"
    if pptx_path.exists():
        print(f"  [skip] {pptx_path.name} exists")
        return pptx_path, orig_bbox

    print(f"--- [{page_name}] {stem} ---")
    ocr_text = _latex_ocr_to_text(comp_meta.get("ocr_latex"))
    if _looks_multi_region_candidate(orig_bbox, slide_w, slide_h, ocr_text):
        try:
            regions = detect_split_regions(
                client, original_path, component_path, orig_bbox,
                slide_w, slide_h, ocr_text,
            )
        except Exception as e:
            regions = []
            print(f"  [split] detector failed: {e}")
        if regions:
            params = {"regions": regions}
            if ink_boxes_y_h:
                params["_ink_boxes_y_h"] = ink_boxes_y_h
            split_log = out_dir / f"{stem}.SPLIT_REGIONS.json"
            split_log.write_text(json.dumps({
                "stem": stem,
                "component_file": cf,
                "region_count": len(regions),
                "regions": regions,
            }, indent=2, ensure_ascii=False))
            pt_used = build_pptx(slide_w, slide_h, orig_bbox, params, pptx_path)
            try:
                rendered = render_pptx_to_png(pptx_path, slide_w, slide_h)
                crop_bbox(rendered, orig_bbox, out_dir / f"{stem}_rendered_crop_split.png")
            except Exception:
                pass
            pt_str = f"{pt_used:.2f}" if pt_used is not None else "n/a"
            print(f"  [split] {len(regions)} regions; saved {pptx_path} (pt={pt_str})")
            return pptx_path, orig_bbox
    raw, params = extract_params(
        client, original_path, component_path, orig_bbox, slide_w, slide_h
    )
    if not _is_valid_extract(params):
        print(f"  [retry] empty/invalid extract; second attempt")
        raw, params = extract_params(
            client, original_path, component_path, orig_bbox, slide_w, slide_h
        )
    if not _is_valid_extract(params):
        if ocr_text:
            params = _apply_ocr_text(params, ocr_text, orig_bbox)
            ocr_log = out_dir / f"{stem}.OCR_FALLBACK.json"
            ocr_log.write_text(json.dumps({
                "stem": stem,
                "component_file": cf,
                "reason": "extract returned empty/invalid text; used metadata ocr_latex",
                "ocr_chars": len(ocr_text),
                "raw_response": raw[:2000],
            }, indent=2, ensure_ascii=False))
            print(f"  [ocr fallback] {stem}: extract invalid; "
                  f"using metadata OCR + {ocr_log.name}")
        else:
            fallback_log = out_dir / f"{stem}.FALLBACK.json"
            fallback_log.write_text(json.dumps({
                "stem": stem,
                "component_file": cf,
                "reason": "extract returned empty/invalid text; preserved original raster",
                "raw_response": raw[:2000],
            }, indent=2, ensure_ascii=False))
            build_raster_fallback_pptx(
                slide_w, slide_h, orig_bbox, component_path, pptx_path
            )
            print(f"  [fallback] {stem}: extract failed after retry; "
                  f"wrote raster fallback + {fallback_log.name}")
            return pptx_path, orig_bbox
    elif _looks_under_extracted(params, ocr_text):
        extracted_chars = _norm_text_len(params.get("text") or "")
        ocr_chars = _norm_text_len(ocr_text)
        params = _apply_ocr_text(params, ocr_text, orig_bbox)
        ocr_log = out_dir / f"{stem}.OCR_FALLBACK.json"
        ocr_log.write_text(json.dumps({
            "stem": stem,
            "component_file": cf,
            "reason": "extract text was substantially shorter than metadata ocr_latex",
            "extracted_norm_chars": extracted_chars,
            "ocr_norm_chars": ocr_chars,
            "raw_response": raw[:2000],
        }, indent=2, ensure_ascii=False))
        print(f"  [ocr fallback] {stem}: under-extracted "
              f"({extracted_chars}/{ocr_chars} chars); using metadata OCR")
    for k, v in DEFAULT_PARAMS.items():
        params.setdefault(k, v)
    if ink_h_px:
        params["_ink_h_px"] = ink_h_px
    if ink_boxes_y_h:
        params["_ink_boxes_y_h"] = ink_boxes_y_h
    try:
        target_lines = max(1, int(params.get("num_lines") or 1))
    except (TypeError, ValueError):
        target_lines = 1

    eff_bbox = orig_bbox
    if placed_text_bboxes is None:
        placed_text_bboxes = []
    pt_used = None
    for r in range(MAX_VERIFY_ROUNDS + 1):
        pt_used = build_pptx(slide_w, slide_h, eff_bbox, params, pptx_path)
        if r == MAX_VERIFY_ROUNDS:
            print(f"  [round {r}] max rounds, accepting")
            break
        try:
            rendered = render_pptx_to_png(pptx_path, slide_w, slide_h)
        except (subprocess.TimeoutExpired, Exception) as e:
            print(f"  [warn] render failed: {e}; stopping verify")
            break
        crop_out = out_dir / f"{stem}_rendered_crop_r{r}.png"
        crop_bbox(rendered, eff_bbox, crop_out)
        raw_v, adj = verify_and_adjust(
            client, component_path, crop_out, params, target_lines,
            eff_bbox, slide_w, slide_h,
        )
        if not isinstance(adj, dict):
            print(f"  [round {r}] verify parse failed; stopping")
            break
        if adj.get("ok") is True:
            print(f"  [round {r}] OK")
            break
        try:
            ext_px = int(adj.get("extend_width_px") or 0)
        except (TypeError, ValueError):
            ext_px = 0
        old_eff = eff_bbox
        if ext_px > 0:
            slack = max(0.0, old_eff[0]) + max(0.0, slide_w - old_eff[2])
            ext_px = min(ext_px, int(slack))
            if ext_px > 0:
                eff_bbox = extend_bbox_horizontally(eff_bbox, ext_px, slide_w)
                print(f"  [round {r}] extend bbox by {ext_px}px -> {eff_bbox}")
        changed = False
        for k in ADJUSTABLE_KEYS:
            if k in adj and adj[k] != params.get(k):
                params[k] = adj[k]; changed = True
        if not changed and eff_bbox == old_eff:
            print(f"  [round {r}] no changes proposed; stopping")
            break

    # Page-level layout fix: after the component fits itself, avoid overlap
    # with previously accepted class-1 text boxes on the same slide.
    resolved_bbox = resolve_text_bbox_overlap(eff_bbox, placed_text_bboxes, slide_h)

    if resolved_bbox != eff_bbox:
        print(f"  [layout] move bbox to avoid text overlap: {eff_bbox} -> {resolved_bbox}")
        eff_bbox = resolved_bbox
        pt_used = build_pptx(slide_w, slide_h, eff_bbox, params, pptx_path)

    if isinstance(params.get("_sizing_debug"), dict):
        (out_dir / f"{stem}.sizing.json").write_text(json.dumps(
            {"component_file": cf, **params["_sizing_debug"]}, indent=2))
    print(f"  saved: {pptx_path}  (pt={pt_used:.2f})")
    return pptx_path, eff_bbox


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--case", required=True)
    ap.add_argument("--out-dir", default=None,
                    help="override per-case output root (default: output/<case>/)")
    args = ap.parse_args()
    paths = case_paths(args.case, out_dir=args.out_dir)
    paths.step3_class1_dir.mkdir(parents=True, exist_ok=True)
    client = anthropic_client()
    pages = list_front_bg_pages(paths)
    if not pages:
        print(f"no front_bg pages under {paths.front_bg_dir}", file=sys.stderr)
        return 1
    for page in pages:
        cat_path = page / "categorization.json"
        meta_path = page / "metadata.json"
        if not cat_path.exists() or not meta_path.exists():
            continue
        cat = json.loads(cat_path.read_text())
        meta = json.loads(meta_path.read_text())
        slide_w = int(meta["query_image_size"]["width"])
        slide_h = int(meta["query_image_size"]["height"])
        original_path = Path(meta["query_image_path"])
        bbox_lookup = {c["component_file"]: c for c in meta["components"]}
        targets = [r for r in cat["results"] if r.get("class") == TARGET_CLASS]
        if not targets:
            continue
        
        print(f"\n=== Page {page.name}: {len(targets)} class-1 ===")

        # Process top-to-bottom, left-to-right, so later text boxes can avoid
        # earlier accepted text boxes on the same page.
        targets = sorted(
            targets,
            key=lambda r: (
                bbox_lookup.get(r["component_file"], {}).get("bbox_xyxy", [0, 0, 0, 0])[1],
                bbox_lookup.get(r["component_file"], {}).get("bbox_xyxy", [0, 0, 0, 0])[0],
            )
        )

        placed_text_bboxes = []
        text_metrics = (
            load_text_metrics(page) if font_sizing_mode() == "deterministic" else {}
        )

        for r in targets:
            comp_meta = bbox_lookup.get(r["component_file"])
            if comp_meta is None:
                continue
            try:
                result = process_component(
                    client, page, slide_w, slide_h, original_path,
                    comp_meta, r, paths.step3_class1_dir,
                    placed_text_bboxes=placed_text_bboxes,
                    ink_h_px=(text_metrics.get(r["component_file"]) or {}).get(
                        "median_box_h_px"),
                    ink_boxes_y_h=(text_metrics.get(r["component_file"]) or {}).get(
                        "boxes_y_h"),
                )
                if result is not None:
                    _, final_bbox = result
                    placed_text_bboxes.append(final_bbox)
            except Exception as e:
                print(f"  [error] {r['component_file']}: {e}")
    return 0


if __name__ == "__main__":
    sys.exit(main())

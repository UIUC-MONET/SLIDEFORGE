"""Step 4: COLOUR-ONLY restyling of every step-3 PPTX.

For every step-3 PPTX (class 1, 3, 4) we rewrite ONLY colours so the
element harmonises with the new styled background. Geometry, font sizes,
row heights, column widths, and element bboxes are NEVER touched.

  Class 1 (text-only):
    Set every text run to the highest-contrast colour in the font palette.
  Class 3 (tables):
    Ask Claude for fill_map / border_map (old hex -> new theme-palette
    hex) and a text colour (override to highest-contrast in font palette).
    Apply via OXML edits to cell fills + borders, plus recolour every
    text run.
  Class 4 (shape compositions):
    Same as class 3 but for shape fills + lines.

Identical input hexes always map to the same output, so the original
colour distribution across shapes / cells is preserved.

Outputs mirror the source layout under
``output/<case>/step4_recolor/{class1,class3,class4}/<page>/<comp>.pptx``,
plus a ``decisions.json`` log of every Claude decision and best-effort
LibreOffice-rendered PNG previews next to each new pptx.
"""

from __future__ import annotations

import argparse
import json
import shutil
import sys
from collections import Counter
from pathlib import Path

from lxml import etree
from pptx import Presentation
from dataclasses import replace

from common import (
    CLAUDE_OPUS,
    anthropic_client,
    call_claude,
    case_paths,
    encode_image_b64,
    hex_norm,
    list_front_bg_pages,
    page_role,
    parse_json_loose,
    render_pptx_to_png,
    styled_bg_for_role,
)


MODEL = CLAUDE_OPUS

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def qn(tag: str) -> str:
    prefix, name = tag.split(":")
    return f"{{{NSMAP[prefix]}}}{name}"


A_R = qn("a:r"); A_RPR = qn("a:rPr"); A_SOLID = qn("a:solidFill")
A_SRGB = qn("a:srgbClr"); A_TCPR = qn("a:tcPr"); A_LN = qn("a:ln")
A_LNL = qn("a:lnL"); A_LNR = qn("a:lnR"); A_LNT = qn("a:lnT"); A_LNB = qn("a:lnB")
P_SP = qn("p:sp"); P_CXNSP = qn("p:cxnSp"); P_SPPR = qn("p:spPr")
A_TBL = qn("a:tbl"); A_TC = qn("a:tc")


# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------

def load_palette(paths):
    pj = paths.styled_bg_dir / "palette.json"
    j = json.loads(pj.read_text())
    res = j["result"]
    return res["theme_palette"], res["font_palette"]


def highest_contrast_font(font_palette):
    return max(font_palette, key=lambda x: float(x["contrast_ratio"]))


# WCAG-AA contrast against the destination background. Text-runs that
# already meet this threshold are LEFT ALONE so source semantic colours
# (red highlights, blue keywords, header-vs-body distinction) survive
# step 4. Runs that fall below the threshold get re-coloured to the
# highest-contrast palette entry.
MIN_CONTRAST = 4.5


def _rel_luminance(rgb):
    def chan(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * chan(rgb[0]) + 0.7152 * chan(rgb[1]) + 0.0722 * chan(rgb[2])


def _contrast_ratio(rgb_a, rgb_b):
    la = _rel_luminance(rgb_a)
    lb = _rel_luminance(rgb_b)
    return (max(la, lb) + 0.05) / (min(la, lb) + 0.05)


def _bg_dominant_rgb(theme_palette):
    """Pick the palette entry whose role is 'background'/'base'; else
    return entry 0."""
    from common import hex_to_rgb as _h
    for e in theme_palette:
        role = (e.get("role") or "").lower()
        if "background" in role or "base" in role or "ground" in role:
            rgb = _h(e["hex"])
            if rgb is not None:
                return rgb
    if theme_palette:
        rgb = _h(theme_palette[0]["hex"])
        if rgb is not None:
            return rgb
    return (255, 255, 255)


def set_text_runs_contrast_aware(root, target_hex, bg_rgb,
                                 min_contrast=MIN_CONTRAST):
    """Re-colour ONLY text runs whose current colour fails contrast
    against ``bg_rgb``. Runs that already meet ``min_contrast`` (the
    typical case for source-side red highlights / blue keywords that
    happened to already be high-contrast) are preserved.

    Returns ``(n_changed, n_kept)`` counts."""
    from common import hex_to_rgb as _h
    target_rgb = _h(target_hex) or (0, 0, 0)
    n_changed = n_kept = 0
    for rpr in iter_run_rPr(root):
        sf = rpr.find(A_SOLID)
        cur_hex = None
        if sf is not None:
            srgb = sf.find(A_SRGB)
            if srgb is not None:
                cur_hex = hex_norm(srgb.get("val"))
        cur_rgb = _h(cur_hex) if cur_hex else None
        if cur_rgb is not None and _contrast_ratio(cur_rgb, bg_rgb) >= min_contrast:
            n_kept += 1
            continue
        if sf is None:
            sf = etree.SubElement(rpr, A_SOLID)
        _set_solidfill_to_hex(sf, target_hex)
        n_changed += 1
    return n_changed, n_kept


# ---------------------------------------------------------------------------
# XML helpers
# ---------------------------------------------------------------------------

def _set_solidfill_to_hex(solid_el, new_hex):
    for child in list(solid_el):
        solid_el.remove(child)
    srgb = etree.SubElement(solid_el, A_SRGB)
    srgb.set("val", new_hex)


def _direct_solidfill(parent):
    return None if parent is None else parent.find(A_SOLID)


def _solidfill_srgb_hex(parent):
    sf = _direct_solidfill(parent)
    if sf is None:
        return None
    s = sf.find(A_SRGB)
    return hex_norm(s.get("val")) if s is not None else None


def iter_run_rPr(root):
    for r in root.iter(A_R):
        rpr = r.find(A_RPR)
        if rpr is not None:
            yield rpr


def set_all_text_runs_to(root, target_hex):
    n = 0
    for rpr in iter_run_rPr(root):
        sf = rpr.find(A_SOLID)
        if sf is None:
            sf = etree.SubElement(rpr, A_SOLID)
        _set_solidfill_to_hex(sf, target_hex)
        n += 1
    return n


def iter_table_cells(root):
    for tbl in root.iter(A_TBL):
        for tc in tbl.iter(A_TC):
            yield tc


def collect_table_fill_and_border_hexes(root):
    fills = Counter(); borders = Counter()
    for tc in iter_table_cells(root):
        tcpr = tc.find(A_TCPR)
        if tcpr is None:
            continue
        ch = _solidfill_srgb_hex(tcpr)
        if ch:
            fills[ch] += 1
        for tag in (A_LNL, A_LNR, A_LNT, A_LNB):
            ln = tcpr.find(tag)
            if ln is None:
                continue
            lh = _solidfill_srgb_hex(ln)
            if lh:
                borders[lh] += 1
    return fills, borders


def apply_table_fill_and_border_map(root, fill_map, border_map):
    nf = nb = 0
    for tc in iter_table_cells(root):
        tcpr = tc.find(A_TCPR)
        if tcpr is None:
            continue
        sf = _direct_solidfill(tcpr)
        if sf is not None:
            srgb = sf.find(A_SRGB)
            if srgb is not None:
                cur = hex_norm(srgb.get("val"))
                if cur and cur in fill_map:
                    _set_solidfill_to_hex(sf, fill_map[cur])
                    nf += 1
        for tag in (A_LNL, A_LNR, A_LNT, A_LNB):
            ln = tcpr.find(tag)
            if ln is None:
                continue
            sf2 = _direct_solidfill(ln)
            if sf2 is None:
                continue
            srgb2 = sf2.find(A_SRGB)
            if srgb2 is None:
                continue
            cur2 = hex_norm(srgb2.get("val"))
            if cur2 and cur2 in border_map:
                _set_solidfill_to_hex(sf2, border_map[cur2])
                nb += 1
    return nf, nb


def iter_shape_spPr(root):
    for sp in root.iter(P_SP):
        spPr = sp.find(P_SPPR)
        if spPr is not None:
            yield spPr
    for cx in root.iter(P_CXNSP):
        spPr = cx.find(P_SPPR)
        if spPr is not None:
            yield spPr


def collect_shape_fill_and_line_hexes(root):
    fills = Counter(); lines = Counter()
    for spPr in iter_shape_spPr(root):
        fh = _solidfill_srgb_hex(spPr)
        ln = spPr.find(A_LN)
        lh = _solidfill_srgb_hex(ln) if ln is not None else None
        if fh:
            fills[fh] += 1
        if lh:
            lines[lh] += 1
    return fills, lines


def apply_shape_fill_and_line_map(root, fill_map, line_map):
    nf = nl = 0
    for spPr in iter_shape_spPr(root):
        sf = _direct_solidfill(spPr)
        if sf is not None:
            srgb = sf.find(A_SRGB)
            if srgb is not None:
                cur = hex_norm(srgb.get("val"))
                if cur and cur in fill_map:
                    _set_solidfill_to_hex(sf, fill_map[cur])
                    nf += 1
        ln = spPr.find(A_LN)
        if ln is not None:
            sf2 = _direct_solidfill(ln)
            if sf2 is not None:
                srgb2 = sf2.find(A_SRGB)
                if srgb2 is not None:
                    cur2 = hex_norm(srgb2.get("val"))
                    if cur2 and cur2 in line_map:
                        _set_solidfill_to_hex(sf2, line_map[cur2])
                        nl += 1
    return nf, nl


def edit_pptx(src_pptx: Path, dst_pptx: Path, mutate_fn) -> dict:
    dst_pptx.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(src_pptx, dst_pptx)
    prs = Presentation(str(dst_pptx))
    totals = {}
    for slide in prs.slides:
        info = mutate_fn(slide._element) or {}
        for k, v in info.items():
            totals[k] = totals.get(k, 0) + v
    prs.save(str(dst_pptx))
    return totals


# ---------------------------------------------------------------------------
# Render helpers
# ---------------------------------------------------------------------------

def safe_render(pptx_path: Path):
    try:
        prs = Presentation(str(pptx_path))
        w = int(round(prs.slide_width / 914400 * 96))
        h = int(round(prs.slide_height / 914400 * 96))
        return render_pptx_to_png(pptx_path, w, h)
    except Exception as e:
        print(f"    [warn] render failed for {pptx_path.name}: {e}")
        return None


def find_existing_rendered(pptx_path: Path):
    stem = pptx_path.stem
    for c in (pptx_path.with_name(stem + "_rendered_crop_final.png"),
              pptx_path.with_name(stem + "_rendered_crop.png"),
              pptx_path.with_name(stem + "_rendered.png")):
        if c.exists():
            return c
    cands = sorted(pptx_path.parent.glob(stem + "*.png"))
    return cands[0] if cands else None


# ---------------------------------------------------------------------------
# Prompts
# ---------------------------------------------------------------------------

TABLE_PROMPT = """You are restyling a PPTX TABLE so it harmonises with a NEW slide background.

Given:
  * Image A - the new slide background.
  * Image B - the table BEFORE restyling.
  * THEME palette (available cell background colours).
  * FONT palette (available text colours, ranked by contrast).
  * Distinct current cell-fill hex values (with counts).
  * Distinct current cell-border hex values (with counts).

Map EACH DISTINCT current cell-fill colour to a NEW hex from the THEME
palette, and EACH DISTINCT cell-border colour to a NEW hex from the THEME
palette.

Hard rules:
  * Pick fills with MODERATE contrast on the new background (not too
    dark, not too loud).
  * Identical input colours must map to the same output colour (a
    function, not a relation). Different input colours SHOULD map to
    different output colours when possible.
  * Border colours should be visible but not overpowering.
  * Only use hex values from THEME palette.
  * Do NOT modify geometry.

For text colour, return the hex of the HIGHEST-CONTRAST entry in FONT
palette.

THEME palette:
<<THEME_PALETTE>>

FONT palette:
<<FONT_PALETTE>>

Distinct cell-fill hexes:
<<FILL_HEXES>>

Distinct cell-border hexes:
<<BORDER_HEXES>>

Reply with STRICT JSON only:
{
  "fill_map": {"<old_hex_no_hash>": "<new_hex_no_hash>", ...},
  "border_map": {"<old_hex_no_hash>": "<new_hex_no_hash>", ...},
  "text_color_hex": "<new_text_hex_no_hash>",
  "rationale": "<one sentence>"
}
"""


SHAPE_PROMPT = """You are restyling a PPTX SHAPE COMPOSITION so it harmonises with a NEW slide background.

Given:
  * Image A - the new slide background.
  * Image B - the shape composition BEFORE restyling.
  * THEME palette (available fill colours).
  * FONT palette (text colours ranked by contrast).
  * Distinct shape-fill hex values (with counts).
  * Distinct shape-line hex values (with counts).

Map EACH DISTINCT shape-fill colour to a NEW hex from THEME palette, and
EACH DISTINCT shape-line colour to a NEW hex from THEME palette.

Hard rules:
  * Moderate contrast against the new background.
  * Identical input -> identical output (preserve original colour
    distribution).
  * Different inputs SHOULD map to different outputs.
  * Border (line) colours should match / complement their fill -
    typically darker or slightly more saturated cousins.
  * Only THEME-palette hexes allowed.
  * Do NOT modify geometry.

Text colour = HIGHEST-CONTRAST entry of FONT palette.

THEME palette:
<<THEME_PALETTE>>

FONT palette:
<<FONT_PALETTE>>

Distinct shape-fill hexes:
<<FILL_HEXES>>

Distinct shape-line hexes:
<<LINE_HEXES>>

Reply with STRICT JSON only:
{
  "fill_map": {"<old_hex_no_hash>": "<new_hex_no_hash>", ...},
  "line_map": {"<old_hex_no_hash>": "<new_hex_no_hash>", ...},
  "text_color_hex": "<new_text_hex_no_hash>",
  "rationale": "<one sentence>"
}
"""


def _fmt_hex_counter(c: Counter) -> str:
    if not c:
        return "(none)"
    items = sorted(c.items(), key=lambda kv: -kv[1])
    return ", ".join(f"#{h} (x{n})" for h, n in items)


def _fmt_palette_list(palette):
    return "\n".join(
        "  - " + ", ".join(f"{k}={e.get(k)!r}" for k in ("hex", "name", "role") if k in e)
        for e in palette
    )


def _fmt_font_palette(font_palette):
    return "\n".join(
        f"  - hex={e.get('hex')!r}, name={e.get('name')!r}, "
        f"contrast_ratio={e.get('contrast_ratio')}, use={e.get('use')!r}"
        for e in font_palette
    )


def _normalise_color_map(raw, valid_inputs, valid_theme):
    out = {}
    if not isinstance(raw, dict):
        return out
    for k, v in raw.items():
        nk = hex_norm(k); nv = hex_norm(v)
        if nk is None or nv is None:
            continue
        if nk not in valid_inputs:
            continue
        if nv not in valid_theme:
            continue
        out[nk] = nv
    return out


# ---------------------------------------------------------------------------
# Per-class drivers
# ---------------------------------------------------------------------------

def process_class1(src: Path, dst: Path, target_hex: str,
                   bg_rgb: tuple) -> dict:
    def mutate(root):
        n_changed, n_kept = set_text_runs_contrast_aware(
            root, target_hex, bg_rgb
        )
        return {"runs_recolored": n_changed, "runs_kept": n_kept}
    info = edit_pptx(src, dst, mutate)
    info["target_hex"] = "#" + target_hex
    return info


def process_class3(client, src: Path, dst: Path, theme_palette, font_palette,
                   bg_image: Path, bg_rgb: tuple) -> dict:
    rendered = find_existing_rendered(src)
    if rendered is None:
        raise RuntimeError(f"no rendered preview near {src}")
    prs = Presentation(str(src))
    fill_c = Counter(); border_c = Counter()
    for slide in prs.slides:
        f, b = collect_table_fill_and_border_hexes(slide._element)
        fill_c.update(f); border_c.update(b)
    valid_theme = {hex_norm(e["hex"]) for e in theme_palette if hex_norm(e["hex"])}
    top_font_hex = hex_norm(highest_contrast_font(font_palette)["hex"])
    prompt = (
        TABLE_PROMPT
        .replace("<<THEME_PALETTE>>", _fmt_palette_list(theme_palette))
        .replace("<<FONT_PALETTE>>", _fmt_font_palette(font_palette))
        .replace("<<FILL_HEXES>>", _fmt_hex_counter(fill_c))
        .replace("<<BORDER_HEXES>>", _fmt_hex_counter(border_c))
    )
    content = [
        {"type": "text", "text": "Image A - new slide background:"},
        encode_image_b64(bg_image),
        {"type": "text", "text": "Image B - current table BEFORE restyling:"},
        encode_image_b64(rendered),
        {"type": "text", "text": prompt},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=1024)
    decision = parse_json_loose(raw)
    if not isinstance(decision, dict):
        raise RuntimeError(f"class-3 Claude reply unparseable: {raw[:200]!r}")
    fmap = _normalise_color_map(decision.get("fill_map"), set(fill_c), valid_theme)
    bmap = _normalise_color_map(decision.get("border_map"), set(border_c), valid_theme)
    text_hex = hex_norm(decision.get("text_color_hex")) or top_font_hex
    if text_hex != top_font_hex:
        text_hex = top_font_hex
    print(f"    fill_map={fmap}  border_map={bmap}  text=#{text_hex}")

    def mutate(root):
        nf, nb = apply_table_fill_and_border_map(root, fmap, bmap)
        nr_c, nr_k = set_text_runs_contrast_aware(root, text_hex, bg_rgb)
        return {"fills_changed": nf, "borders_changed": nb,
                "runs_recolored": nr_c, "runs_kept": nr_k}
    info = edit_pptx(src, dst, mutate)
    info["decision"] = {
        "fill_map": {f"#{k}": f"#{v}" for k, v in fmap.items()},
        "border_map": {f"#{k}": f"#{v}" for k, v in bmap.items()},
        "text_color_hex": f"#{text_hex}",
        "rationale": decision.get("rationale", ""),
    }
    return info


def process_class4(client, src: Path, dst: Path, theme_palette, font_palette,
                   bg_image: Path, bg_rgb: tuple) -> dict:
    rendered = find_existing_rendered(src)
    if rendered is None:
        raise RuntimeError(f"no rendered preview near {src}")
    prs = Presentation(str(src))
    fill_c = Counter(); line_c = Counter()
    for slide in prs.slides:
        f, l = collect_shape_fill_and_line_hexes(slide._element)
        fill_c.update(f); line_c.update(l)
    valid_theme = {hex_norm(e["hex"]) for e in theme_palette if hex_norm(e["hex"])}
    top_font_hex = hex_norm(highest_contrast_font(font_palette)["hex"])
    prompt = (
        SHAPE_PROMPT
        .replace("<<THEME_PALETTE>>", _fmt_palette_list(theme_palette))
        .replace("<<FONT_PALETTE>>", _fmt_font_palette(font_palette))
        .replace("<<FILL_HEXES>>", _fmt_hex_counter(fill_c))
        .replace("<<LINE_HEXES>>", _fmt_hex_counter(line_c))
    )
    content = [
        {"type": "text", "text": "Image A - new slide background:"},
        encode_image_b64(bg_image),
        {"type": "text", "text": "Image B - current shape composition BEFORE restyling:"},
        encode_image_b64(rendered),
        {"type": "text", "text": prompt},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=1024)
    decision = parse_json_loose(raw)
    if not isinstance(decision, dict):
        raise RuntimeError(f"class-4 Claude reply unparseable: {raw[:200]!r}")
    fmap = _normalise_color_map(decision.get("fill_map"), set(fill_c), valid_theme)
    lmap = _normalise_color_map(decision.get("line_map"), set(line_c), valid_theme)
    text_hex = hex_norm(decision.get("text_color_hex")) or top_font_hex
    if text_hex != top_font_hex:
        text_hex = top_font_hex
    print(f"    fill_map={fmap}  line_map={lmap}  text=#{text_hex}")

    def mutate(root):
        nf, nl = apply_shape_fill_and_line_map(root, fmap, lmap)
        nr_c, nr_k = set_text_runs_contrast_aware(root, text_hex, bg_rgb)
        return {"fills_changed": nf, "lines_changed": nl,
                "runs_recolored": nr_c, "runs_kept": nr_k}
    info = edit_pptx(src, dst, mutate)
    info["decision"] = {
        "fill_map": {f"#{k}": f"#{v}" for k, v in fmap.items()},
        "line_map": {f"#{k}": f"#{v}" for k, v in lmap.items()},
        "text_color_hex": f"#{text_hex}",
        "rationale": decision.get("rationale", ""),
    }
    return info


def _iter_pptx_under(root: Path):
    if not root.is_dir():
        return
    yield from sorted(root.rglob("*.pptx"))


# Step 2 records `semantic_role` per component. Step 4 must skip recolour
# on roles that encode meaning in their colour (chart legends, data
# tables, technical flow diagrams). For these we copy the step-3 PPTX
# through unchanged so step 5 still finds a file at the expected path.
PROTECTED_ROLES = {
    "chart", "legend", "table", "technical_diagram",
    "shape_encoded_chart", "data_diagram", "flow_diagram",
}


def _load_categorization(paths, page_name: str) -> dict | None:
    cat_path = paths.front_bg_dir / page_name / "categorization.json"
    if not cat_path.exists():
        return None
    try:
        return json.loads(cat_path.read_text())
    except Exception:
        return None


def _semantic_role_for(paths, page_name: str, stem: str) -> str | None:
    cat = _load_categorization(paths, page_name)
    if not cat:
        return None
    for r in cat.get("results", []):
        cf_stem = Path(r.get("component_file") or "").stem
        if cf_stem == stem:
            return (r.get("semantic_role") or "").lower() or None
    return None


def _entry_for(paths, page_name: str, stem: str) -> dict | None:
    cat = _load_categorization(paths, page_name)
    if not cat:
        return None
    for r in cat.get("results", []):
        cf_stem = Path(r.get("component_file") or "").stem
        if cf_stem == stem:
            return r
    return None


def _is_protected_component(paths, page_name: str, stem: str) -> tuple[bool, str | None]:
    entry = _entry_for(paths, page_name, stem)
    role = ((entry or {}).get("semantic_role") or "").lower() or None
    if bool((entry or {}).get("protected_content")) and role in PROTECTED_ROLES:
        return True, role
    if role in PROTECTED_ROLES:
        return True, role
    return False, role


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--case", required=True)
    ap.add_argument("--out-dir", default=None)
    args = ap.parse_args()

    paths = case_paths(args.case)

    if args.out_dir is not None:
        out = Path(args.out_dir)
        paths = replace(
            paths,
            front_bg_dir=out / "front_bg",
            styled_bg_dir=out / "styled_bg",
            step3_class1_dir=out / "step3_class1",
            step3_class3_dir=out / "step3_class3",
            step3_class4_dir=out / "step3_class4",
            step4_dir=out / "step4_recolor",
        )
    palette_json = paths.styled_bg_dir / "palette.json"
    if not palette_json.exists():
        print(f"missing palette.json: {palette_json}", file=sys.stderr)
        return 1
    theme_palette, font_palette = load_palette(paths)
    pages = list_front_bg_pages(paths)
    page_to_role = {p.name: page_role(i, len(pages)) for i, p in enumerate(pages)}

    def bg_for_rel(rel: Path) -> Path:
        page_name = rel.parts[0] if rel.parts else ""
        role = page_to_role.get(page_name, "content")
        bg_path = styled_bg_for_role(paths, role)
        if not bg_path.exists():
            bg_path = paths.styled_bg_dir / "content.png"
        return bg_path

    fallback_bg = paths.styled_bg_dir / "content.png"
    if not fallback_bg.exists():
        print(f"missing fallback background: {fallback_bg}", file=sys.stderr)
        return 1
    top_font_hex = hex_norm(highest_contrast_font(font_palette)["hex"])
    bg_rgb = _bg_dominant_rgb(theme_palette)
    print(f"theme palette ({len(theme_palette)}): "
          + ", ".join(f"#{hex_norm(e['hex'])}" for e in theme_palette))
    print(f"top-font-color: #{top_font_hex}")
    print(f"deck background RGB (for contrast checks): {bg_rgb}")

    paths.step4_dir.mkdir(parents=True, exist_ok=True)
    client = anthropic_client()
    decisions = []

    for cls, src_root, label in (
        (1, paths.step3_class1_dir, "class1"),
        (3, paths.step3_class3_dir, "class3"),
        (4, paths.step3_class4_dir, "class4"),
    ):
        print(f"\n=== {label} ===")
        for src in _iter_pptx_under(src_root):
            rel = src.relative_to(src_root)
            dst = paths.step4_dir / label / rel
            bg = bg_for_rel(rel)
            if not bg.exists():
                print(f"  [skip] {rel}: missing role background {bg}")
                continue
            if dst.exists():
                print(f"  [skip] {dst.relative_to(paths.step4_dir)} exists")
                continue

            # F-09: copy through unchanged for protected (data-bearing)
            # components so colour-coded semantics are preserved.
            page_name = rel.parts[0] if rel.parts else ""
            stem = rel.stem
            is_protected, role = _is_protected_component(paths, page_name, stem)
            if is_protected:
                print(f"  [protect] {rel}: semantic_role={role!r} (copied through)")
                dst.parent.mkdir(parents=True, exist_ok=True)
                shutil.copyfile(src, dst)
                safe_render(dst)
                decisions.append({
                    "class": cls,
                    "src": str(src),
                    "dst": str(dst),
                    "bg_image": str(bg),
                    "page_role": page_to_role.get(page_name, "content"),
                    "protected": True,
                    "semantic_role": role,
                })
                continue

            print(f"  [{label}] {rel}")
            try:
                if cls == 1:
                    info = process_class1(src, dst, top_font_hex, bg_rgb)
                elif cls == 3:
                    info = process_class3(client, src, dst, theme_palette,
                                          font_palette, bg, bg_rgb)
                else:
                    info = process_class4(client, src, dst, theme_palette,
                                          font_palette, bg, bg_rgb)
            except Exception as e:
                print(f"    [error] {e}")
                continue
            safe_render(dst)
            decisions.append({
                "class": cls,
                "src": str(src),
                "dst": str(dst),
                "bg_image": str(bg),
                "page_role": page_to_role.get(rel.parts[0], "content") if rel.parts else "content",
                **info,
            })

    (paths.step4_dir / "decisions.json").write_text(json.dumps({
        "model": MODEL,
        "bg_image": "role-specific",
        "palette_json": str(palette_json),
        "top_font_hex": f"#{top_font_hex}",
        "decisions": decisions,
    }, indent=2, ensure_ascii=False))
    print(f"\nDone -> {paths.step4_dir}")
    return 0


if __name__ == "__main__":
    sys.exit(main())

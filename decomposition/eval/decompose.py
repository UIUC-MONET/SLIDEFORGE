import os
import re
import json
import math
import subprocess
import shutil
import glob
from concurrent.futures import ThreadPoolExecutor
from pptx import Presentation
from pptx.dml.color import RGBColor
from PIL import Image

def sanitize_filename(name):
    sanitized = re.sub(r'[^a-zA-Z0-9_\-]', '_', name)
    sanitized = re.sub(r'_+', '_', sanitized)
    return sanitized.strip('_')

def run_command(args):
    result = subprocess.run(args, capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(result.stderr)
    return result.stdout

def get_rotated_bounding_box(left, top, width, height, rotation_degrees):
    """
    Computes the exact absolute bounding box of a rotated rectangle.
    All parameters are in the same unit (EMUs or Pixels).
    """
    if not rotation_degrees:
        return left, top, width, height
        
    cx = left + width / 2.0
    cy = top + height / 2.0
    
    rad = math.radians(rotation_degrees)
    cos_val = math.cos(rad)
    sin_val = math.sin(rad)
    
    hw = width / 2.0
    hh = height / 2.0
    
    corners = [
        (-hw, -hh),
        (hw, -hh),
        (hw, hh),
        (-hw, hh)
    ]
    
    rotated_xs = []
    rotated_ys = []
    for x, y in corners:
        rx = x * cos_val - y * sin_val
        ry = x * sin_val + y * cos_val
        rotated_xs.append(rx)
        rotated_ys.append(ry)
        
    min_rx = min(rotated_xs)
    max_rx = max(rotated_xs)
    min_ry = min(rotated_ys)
    max_ry = max(rotated_ys)
    
    r_left = cx + min_rx
    r_top = cy + min_ry
    r_width = max_rx - min_rx
    r_height = max_ry - min_ry
    
    return r_left, r_top, r_width, r_height

def make_background_transparent(img, bg_color=(255, 0, 255), tolerance=180):
    """
    Makes the background transparent by keying out all pixels close to bg_color globally.
    This perfectly transparentizes the inside holes of characters (like 'o', 'a', 'd') 
    since pure magenta is a unique chroma-key color that does not appear in the shapes.
    Uses squared Euclidean distance in RGB space to avoid slow square root operations.
    """
    img = img.convert("RGBA")
    data = img.load()
    w, h = img.size
    
    bg_r, bg_g, bg_b = bg_color
    tolerance_sq = tolerance ** 2
    
    for y in range(h):
        for x in range(w):
            r, g, b, a = data[x, y]
            sq_dist = (r - bg_r)**2 + (g - bg_g)**2 + (b - bg_b)**2
            if sq_dist <= tolerance_sq:
                data[x, y] = (0, 0, 0, 0)
                
    return img

def find_all_shapes(shape_container, parent_group=None):
    shapes_list = []
    for shape in shape_container:
        is_group = False
        try:
            is_group = shape.shape_type == 6 or hasattr(shape, "shapes")
        except Exception:
            pass
        
        shapes_list.append((shape, parent_group, is_group))
        
        if is_group:
            try:
                shapes_list.extend(find_all_shapes(shape.shapes, parent_group=shape))
            except Exception as e:
                print(f"Warning: Failed to recurse group shape {shape.name}: {e}")
    return shapes_list

def convert_pdf_to_png(args):
    """
    Parallel worker to run pdftoppm on a single PDF.
    """
    temp_pdf, temp_png_base, dpi = args
    try:
        run_command(["pdftoppm", "-png", "-r", str(dpi), temp_pdf, temp_png_base])
        return True
    except Exception as e:
        print(f"Error converting PDF {temp_pdf} to PNG: {e}")
        return False

def crop_and_transparentize_worker(args):
    """
    Parallel worker to crop and key out a single shape image.
    """
    (
        idx, shape_id, shape_name, shape_type_str, parent_id, rotation, text_content,
        temp_png_file, output_img_path, crop_box, is_full_bleed, emu_coords, inch_coords
    ) = args
    
    left_px_int, top_px_int, right_px_int, bottom_px_int = crop_box
    crop_width_px = right_px_int - left_px_int
    crop_height_px = bottom_px_int - top_px_int
    
    try:
        with Image.open(temp_png_file) as img:
            cropped_img = img.crop(crop_box)
            if not is_full_bleed:
                cropped_img = make_background_transparent(cropped_img, bg_color=(255, 0, 255), tolerance=180)
            cropped_img.save(output_img_path)
            
        return {
            "index": idx,
            "shape_id": shape_id,
            "name": shape_name,
            "type": shape_type_str,
            "is_group": False,  # Expanded
            "parent_group_id": parent_id,
            "rotation_degrees": rotation,
            "text": text_content,
            "image_path": os.path.relpath(output_img_path, os.path.dirname(os.path.dirname(output_img_path))),
            "position_pixels": {
                "left": left_px_int,
                "top": top_px_int,
                "width": crop_width_px,
                "height": crop_height_px,
                "right": right_px_int,
                "bottom": bottom_px_int
            },
            "position_pixels_unrotated": {
                "left": round(emu_coords[0] * (300/914400)),
                "top": round(emu_coords[1] * (300/914400)),
                "width": round(emu_coords[2] * (300/914400)),
                "height": round(emu_coords[3] * (300/914400))
            },
            "position_inches_rotated": {
                "left": round(inch_coords[0], 4),
                "top": round(inch_coords[1], 4),
                "width": round(inch_coords[2], 4),
                "height": round(inch_coords[3], 4)
            },
            "position_emu_rotated": {
                "left": round(emu_coords[4]),
                "top": round(emu_coords[5]),
                "width": round(emu_coords[6]),
                "height": round(emu_coords[7])
            }
        }
    except Exception as e:
        print(f"Error cropping shape {idx} ({shape_name}): {e}")
        return None

def decompose_slide_single_rendering(pptx_path, output_dir, json_path):
    print("=== Ultra-fast Batch isolated Rendering & Decomposition Pipeline ===")
    
    base_dir = os.path.dirname(os.path.abspath(pptx_path))
    temp_dir = os.path.join(base_dir, "temp_render")
    os.makedirs(temp_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)
    
    dpi = 300
    emu_per_inch = 914400
    scale = dpi / emu_per_inch
    
    # Load original PPTX to get sizes and shapes
    prs = Presentation(pptx_path)
    if not prs.slides:
        raise ValueError("The presentation contains no slides.")
    
    slide = prs.slides[0]
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    slide_width_px = round(slide_width_emu * scale)
    slide_height_px = round(slide_height_emu * scale)
    
    print(f"Slide Size: {slide_width_px}x{slide_height_px} px ({slide_width_emu}x{slide_height_emu} EMUs)")
    
    all_shapes_info = find_all_shapes(slide.shapes)
    print(f"Total shapes to decompose: {len(all_shapes_info)}")
    
    valid_shapes_list = []
    temp_pptx_files = []
    
    # Step 1: Create isolated PPTX files in a single fast loop
    print("Step 1: Cloned PPTX files XML generation...")
    for idx, (original_shape, parent, is_group) in enumerate(all_shapes_info):
        shape_name = original_shape.name
        shape_id = original_shape.shape_id
        rotation = getattr(original_shape, "rotation", 0) or 0
        
        try:
            shape_type_str = original_shape.shape_type.name
        except Exception:
            shape_type_str = str(original_shape.shape_type)
            
        # Coordinates in EMUs
        try:
            left_emu = original_shape.left
            top_emu = original_shape.top
            width_emu = original_shape.width
            height_emu = original_shape.height
        except AttributeError:
            continue
            
        if left_emu is None or top_emu is None or width_emu is None or height_emu is None:
            continue
            
        # Compute rotated bounding box
        r_left_emu, r_top_emu, r_width_emu, r_height_emu = get_rotated_bounding_box(
            left_emu, top_emu, width_emu, height_emu, rotation
        )
        
        # Convert to pixels
        left_px_rot = r_left_emu * scale
        top_px_rot = r_top_emu * scale
        width_px_rot = r_width_emu * scale
        height_px_rot = r_height_emu * scale
        
        left_px_int = max(0, min(round(left_px_rot), slide_width_px))
        top_px_int = max(0, min(round(top_px_rot), slide_height_px))
        right_px_int = max(0, min(round(left_px_rot + width_px_rot), slide_width_px))
        bottom_px_int = max(0, min(round(top_px_rot + height_px_rot), slide_height_px))
        
        # Ensure horizontal/vertical lines (width or height close to 0) have a minimum visible thickness (8px)
        min_thickness_px = 8
        if (right_px_int - left_px_int) < min_thickness_px:
            cx_px = (left_px_int + right_px_int) / 2.0
            left_px_int = max(0, round(cx_px - min_thickness_px / 2.0))
            right_px_int = min(slide_width_px, left_px_int + min_thickness_px)
            
        if (bottom_px_int - top_px_int) < min_thickness_px:
            cy_px = (top_px_int + bottom_px_int) / 2.0
            top_px_int = max(0, round(cy_px - min_thickness_px / 2.0))
            bottom_px_int = min(slide_height_px, top_px_int + min_thickness_px)
            
        crop_width_px = right_px_int - left_px_int
        crop_height_px = bottom_px_int - top_px_int
        
        if crop_width_px <= 0 or crop_height_px <= 0:
            print(f"  Shape {idx} ({shape_name}): Skipped due to zero size.")
            continue
            
        # Text
        text_content = ""
        if original_shape.has_text_frame:
            text_content = original_shape.text.strip()
            
        # Copy presentation
        temp_pptx = os.path.join(temp_dir, f"temp_shape_{idx}.pptx")
        shutil.copyfile(pptx_path, temp_pptx)
        
        # Modify the XML copy
        temp_prs = Presentation(temp_pptx)
        temp_slide = temp_prs.slides[0]
        
        # Set slide background to chroma key magenta (255, 0, 255) unless it's full bleed
        is_full_bleed = (left_px_int == 0 and top_px_int == 0 and right_px_int == slide_width_px and bottom_px_int == slide_height_px)
        if not is_full_bleed:
            try:
                fill = temp_slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 0, 255)
            except Exception:
                pass
                
        # Delete other shapes
        temp_shapes_to_delete = []
        for temp_s in temp_slide.shapes:
            if temp_s.shape_id != shape_id:
                temp_shapes_to_delete.append(temp_s)
                
        for s_del in temp_shapes_to_delete:
            try:
                el = s_del._element
                parent_el = el.getparent()
                if parent_el is not None:
                    parent_el.remove(el)
            except Exception:
                pass
                
        temp_prs.save(temp_pptx)
        temp_pptx_files.append(temp_pptx)
        
        sanitized_name = sanitize_filename(shape_name)
        filename = f"shape_{idx}_{shape_id}_{sanitized_name}.png"
        output_img_path = os.path.join(output_dir, filename)
        
        # Queue for cropping/rendering
        valid_shapes_list.append({
            "idx": idx,
            "shape_id": shape_id,
            "shape_name": shape_name,
            "shape_type_str": shape_type_str,
            "parent_id": parent.shape_id if parent else None,
            "rotation": rotation,
            "text_content": text_content,
            "temp_pptx": temp_pptx,
            "temp_pdf": os.path.join(temp_dir, f"temp_shape_{idx}.pdf"),
            "temp_png_base": os.path.join(temp_dir, f"temp_page_{idx}"),
            "temp_png_file": os.path.join(temp_dir, f"temp_page_{idx}-1.png"),
            "output_img_path": output_img_path,
            "crop_box": (left_px_int, top_px_int, right_px_int, bottom_px_int),
            "is_full_bleed": is_full_bleed,
            "emu_coords": (left_emu, top_emu, width_emu, height_emu, r_left_emu, r_top_emu, r_width_emu, r_height_emu),
            "inch_coords": (r_left_emu/emu_per_inch, r_top_emu/emu_per_inch, r_width_emu/emu_per_inch, r_height_emu/emu_per_inch)
        })
        
    if not valid_shapes_list:
        print("No shapes to render.")
        shutil.rmtree(temp_dir, ignore_errors=True)
        return
        
    # Step 2: Batch render all PPTX files to PDF in a SINGLE soffice command!
    print(f"\nStep 2: Batch rendering {len(temp_pptx_files) + 1} PPTX files to PDF in a single LibreOffice session...")
    user_inst_dir = f"file://{temp_dir}/soffice_profile_batch"
    
    full_slide_temp_pptx = os.path.join(temp_dir, "full_slide_original.pptx")
    shutil.copyfile(pptx_path, full_slide_temp_pptx)
    
    soffice_args = [
        "soffice", f"-env:UserInstallation={user_inst_dir}",
        "--headless", "--convert-to", "pdf", "--outdir", temp_dir
    ] + temp_pptx_files + [full_slide_temp_pptx]
    
    run_command(soffice_args)
    print("  Batch LibreOffice PDF conversion completed successfully!")
    
    # Step 3: Convert all PDFs to PNGs concurrently (pdftoppm is extremely fast)
    print(f"\nStep 3: Converting PDFs to PNGs concurrently using ThreadPoolExecutor...")
    png_args_list = []
    for shape in valid_shapes_list:
        png_args_list.append((shape["temp_pdf"], shape["temp_png_base"], dpi))
        
    full_slide_pdf = os.path.join(temp_dir, "full_slide_original.pdf")
    full_slide_png_base = os.path.join(temp_dir, "full_slide_original")
    full_slide_png_file = os.path.join(temp_dir, "full_slide_original-1.png")
    png_args_list.append((full_slide_pdf, full_slide_png_base, dpi))
        
    max_workers = 16  # High level of concurrency for extremely fast I/O
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        executor.map(convert_pdf_to_png, png_args_list)
        
    # Step 4: Crop and transparentize all shapes in parallel
    print(f"\nStep 4: Cropping and transparentizing all shapes concurrently...")
    crop_args_list = []
    for shape in valid_shapes_list:
        # Check if the PNG file was actually generated
        png_file = shape["temp_png_file"]
        if not os.path.exists(png_file):
            # Fallback check
            base_png = shape["temp_png_base"]
            pngs = glob.glob(f"{base_png}*.png")
            if pngs:
                png_file = pngs[0]
            else:
                print(f"  Warning: PNG not found for shape {shape['idx']} ({shape['shape_name']}). Skipping crop.")
                continue
                
        crop_args_list.append((
            shape["idx"], shape["shape_id"], shape["shape_name"], shape["shape_type_str"],
            shape["parent_id"], shape["rotation"], shape["text_content"],
            png_file, shape["output_img_path"], shape["crop_box"],
            shape["is_full_bleed"], shape["emu_coords"], shape["inch_coords"]
        ))
        
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        results = list(executor.map(crop_and_transparentize_worker, crop_args_list))
        
    # Step 4.5: Save full slide screenshot
    print("\nStep 4.5: Saving full slide screenshot...")
    if os.path.exists(full_slide_png_file):
        base_name_pptx = os.path.splitext(os.path.basename(pptx_path))[0]
        reconstruction_output_path = os.path.join(base_dir, f"{base_name_pptx}_reconstruction.png")
        shutil.copyfile(full_slide_png_file, reconstruction_output_path)
        print(f"  Saved full slide screenshot to {reconstruction_output_path}")
    else:
        print("  Warning: Full slide PNG not found. Skipping screenshot save.")
        
    # Filter out None and sort by index
    components_meta = [r for r in results if r is not None]
    components_meta.sort(key=lambda x: x["index"])
    
    # Write JSON metadata
    json_output = {
        "slide_info": {
            "file": os.path.basename(pptx_path),
            "dimensions_emu": {
                "width": slide_width_emu,
                "height": slide_height_emu
            },
            "dimensions_pixels": {
                "width": slide_width_px,
                "height": slide_height_px
            },
            "dpi": dpi
        },
        "components": components_meta
    }
    
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(json_output, f, indent=2, ensure_ascii=False)
        
    # Step 5: Clean up all temp files
    print("\nStep 5: Cleaning up all temporary profiles and render cache...")
    shutil.rmtree(temp_dir, ignore_errors=True)
    
    print(f"\nSaved metadata JSON to {json_path}")
    print("=== Ultra-fast Parallel Decomposition Completed Successfully! ===")

def process_path(input_path):
    if os.path.isdir(input_path):
        pptx_files = glob.glob(os.path.join(input_path, "**", "*.pptx"), recursive=True)
        print(f"Found {len(pptx_files)} PPTX files in directory {input_path}")
        for pptx in pptx_files:
            process_single_pptx(pptx)
    elif os.path.isfile(input_path) and input_path.lower().endswith(".pptx"):
        process_single_pptx(input_path)
    else:
        print(f"Error: Invalid path or not a PPTX file: {input_path}")

def process_single_pptx(pptx_file):
    base_name = os.path.splitext(os.path.basename(pptx_file))[0]
    base_dir = os.path.dirname(os.path.abspath(pptx_file))
    
    comp_dir = os.path.join(base_dir, f"{base_name}_components")
    meta_json = os.path.join(base_dir, f"{base_name}_components.json")
    
    try:
        decompose_slide_single_rendering(pptx_file, comp_dir, meta_json)
    except Exception as e:
        print(f"Error processing {pptx_file}: {e}")

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        input_path = sys.argv[1]
    else:
        print("Usage: python decompose.py <path_to_pptx_or_dir>", file=sys.stderr)
        sys.exit(1)
        
    if not os.path.exists(input_path):
        print(f"Error: Path {input_path} does not exist.")
        sys.exit(1)
        
    process_path(input_path)

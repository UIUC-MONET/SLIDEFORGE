"""Step 4b (H4): per-run LOCAL contrast audit over step4 outputs.

step4 tests text contrast only against the deck-level dominant background, so
runs inside dark shape fills (recolored by step4 itself) or over dark styled-bg
regions stay dark-on-dark ("visual loss": text present in pptx, unreadable in
render — main19 p0001/p0004).

For every pptx under step4_recolor/: for each text run, the LOCAL background is
  1. the containing shape's solid fill, else
  2. the topmost solid-filled shape underneath (z-order + center containment), else
  3. the mean color of the styled background under the shape's bbox.
Runs with WCAG contrast < 3.0 are recolored to the argmax-contrast candidate
among the font palette + black/white. Deterministic — zero VLM calls.

usage: python3 step4b_contrast_audit.py --case CASE [--out-dir DIR]
Audit log: step4_recolor/contrast_audit.json
"""
from __future__ import annotations

import argparse
import json
import sys
from dataclasses import replace
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE

from common import (
    EMU_PER_PX,
    case_paths,
    hex_norm,
    hex_to_rgb,
    list_front_bg_pages,
    page_role,
    styled_bg_for_role,
)

MIN_LOCAL_CONTRAST = 3.0


def _rel_luminance(rgb):
    def chan(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * chan(rgb[0]) + 0.7152 * chan(rgb[1]) + 0.0722 * chan(rgb[2])


def contrast_ratio(a, b):
    la, lb = _rel_luminance(a), _rel_luminance(b)
    return (max(la, lb) + 0.05) / (min(la, lb) + 0.05)


def _shape_fill_rgb(shape):
    try:
        if shape.fill.type == MSO_FILL_TYPE.SOLID:
            c = shape.fill.fore_color
            if c.type == MSO_COLOR_TYPE.RGB:
                return tuple(c.rgb)
        return None
    except Exception:
        return None


def _bbox_px(shape):
    try:
        return (shape.left / EMU_PER_PX, shape.top / EMU_PER_PX,
                (shape.left + shape.width) / EMU_PER_PX,
                (shape.top + shape.height) / EMU_PER_PX)
    except Exception:
        return None


def _iter_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == 6:  # group
            yield from _iter_shapes(sh.shapes)
        else:
            yield sh


def _mean_rgb(img, bbox):
    x0, y0, x1, y1 = (int(v) for v in bbox)
    x0 = max(0, min(x0, img.width - 1)); x1 = max(x0 + 1, min(x1, img.width))
    y0 = max(0, min(y0, img.height - 1)); y1 = max(y0 + 1, min(y1, img.height))
    crop = img.crop((x0, y0, x1, y1)).resize((8, 8))
    px = list(crop.getdata())
    n = len(px)
    return tuple(sum(p[k] for p in px) // n for k in range(3))


def _run_rgb(run):
    try:
        c = run.font.color
        if c is not None and c.type == MSO_COLOR_TYPE.RGB:
            return tuple(c.rgb)
    except Exception:
        pass
    return (0, 0, 0)  # pptx default text colour is black


def audit_pptx(pptx_path: Path, bg_img, candidates) -> list[dict]:
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    flat = list(_iter_shapes(slide.shapes))
    fixes = []
    for idx, shape in enumerate(flat):
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        if not any(r.text.strip() for p in tf.paragraphs for r in p.runs):
            continue
        local_bg = _shape_fill_rgb(shape)
        bbox = _bbox_px(shape)
        if local_bg is None and bbox is not None:
            cx, cy = (bbox[0] + bbox[2]) / 2, (bbox[1] + bbox[3]) / 2
            for other in flat[:idx]:
                ob = _bbox_px(other)
                if ob and ob[0] <= cx <= ob[2] and ob[1] <= cy <= ob[3]:
                    f = _shape_fill_rgb(other)
                    if f is not None:
                        local_bg = f  # keep last hit = topmost below
        if local_bg is None:
            if bbox is None or bg_img is None:
                continue
            local_bg = _mean_rgb(bg_img, bbox)
        for p in tf.paragraphs:
            for r in p.runs:
                if not r.text.strip():
                    continue
                cur = _run_rgb(r)
                ratio = contrast_ratio(cur, local_bg)
                if ratio >= MIN_LOCAL_CONTRAST:
                    continue
                best = max(candidates, key=lambda c: contrast_ratio(c, local_bg))
                r.font.color.rgb = RGBColor(*best)
                fixes.append({
                    "text": r.text[:40], "old": "#%02X%02X%02X" % cur,
                    "new": "#%02X%02X%02X" % best,
                    "local_bg": "#%02X%02X%02X" % tuple(local_bg),
                    "old_ratio": round(ratio, 2),
                    "new_ratio": round(contrast_ratio(best, local_bg), 2),
                })
    if fixes:
        prs.save(str(pptx_path))
    return fixes


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--case", required=True)
    ap.add_argument("--out-dir", default=None)
    args = ap.parse_args()
    paths = case_paths(args.case)
    if args.out_dir is not None:
        out = Path(args.out_dir)
        paths = replace(
            paths,
            front_bg_dir=out / "front_bg",
            styled_bg_dir=out / "styled_bg",
            step4_dir=out / "step4_recolor",
        )
    palette_json = paths.styled_bg_dir / "palette.json"
    candidates = [(0, 0, 0), (255, 255, 255)]
    if palette_json.exists():
        pal = json.loads(palette_json.read_text())
        for e in ((pal.get("result") or {}).get("font_palette") or []):
            rgb = hex_to_rgb(e.get("hex"))
            if rgb:
                candidates.append(rgb)
    pages = list_front_bg_pages(paths)
    page_to_role = {p.name: page_role(i, len(pages)) for i, p in enumerate(pages)}
    bg_cache: dict[str, object] = {}
    all_fixes: dict[str, list] = {}
    for pptx in sorted(paths.step4_dir.glob("*/*/*.pptx")):
        page_name = pptx.parent.name
        role = page_to_role.get(page_name, "content")
        if role not in bg_cache:
            bg_path = styled_bg_for_role(paths, role)
            if not bg_path.exists():
                bg_path = paths.styled_bg_dir / "content.png"
            bg_cache[role] = Image.open(bg_path).convert("RGB") if bg_path.exists() else None
        rel = str(pptx.relative_to(paths.step4_dir))
        fixes = audit_pptx(pptx, bg_cache[role], candidates)
        if fixes:
            all_fixes[rel] = fixes
            print(f"  [audit] {rel}: {len(fixes)} run(s) recolored")
    (paths.step4_dir / "contrast_audit.json").write_text(
        json.dumps(all_fixes, indent=1, ensure_ascii=False))
    n = sum(len(v) for v in all_fixes.values())
    print(f"contrast audit: {n} run(s) fixed across {len(all_fixes)} file(s)")
    return 0


if __name__ == "__main__":
    sys.exit(main())

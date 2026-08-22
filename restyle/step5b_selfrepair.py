"""Step 5b (H7): end-to-end OCR-diff self-repair of final.pptx.

Deterministic, zero VLM calls. Per page: easyocr the ORIGINAL page and the
rendered final page; any original word-group that is missing from BOTH the
render OCR and the pptx text layer (fuzzy match, so misreads don't duplicate)
is injected into the final pptx as a textbox at its original position, sized
by the H1 fit engine and colored for local contrast. This is the safety net
for text no per-class fix reaches (class5 raster text, recon occlusion).

Must run under the easyocr environment (see SLIDECODER_PYTHON in common.py):
    $SLIDECODER_PYTHON step5b_selfrepair.py --case CASE --out-dir RUN

Requires RUN/final.pptx and RUN/render/page-N.png (from render_pptx_pages.py).
Re-render after this step to measure. Log: RUN/selfrepair.json
"""
from __future__ import annotations

import argparse
import difflib
import json
import re
import statistics
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from common import (  # noqa: E402
    fit_font_size_pt,
    carlito_metrics_unreliable,
    hex_to_rgb,
    px_to_emu,
)

CONF_MIN = 0.4
MIN_NORM_LEN = 3
FUZZ = 0.8


def _norm(s: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", (s or "").lower())


def _covered(n: str, hay_items: list[str], hay_concat: str) -> bool:
    if n in hay_concat:
        return True
    for g in hay_items:
        if difflib.SequenceMatcher(None, n, g).ratio() >= FUZZ:
            return True
    return False


def _mean_rgb(img, bbox):
    x0, y0, x1, y1 = (int(v) for v in bbox)
    x0 = max(0, min(x0, img.width - 1)); x1 = max(x0 + 1, min(x1, img.width))
    y0 = max(0, min(y0, img.height - 1)); y1 = max(y0 + 1, min(y1, img.height))
    px = list(img.crop((x0, y0, x1, y1)).resize((6, 6)).getdata())
    return tuple(sum(p[k] for p in px) // len(px) for k in range(3))


def _rel_lum(rgb):
    def ch(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * ch(rgb[0]) + 0.7152 * ch(rgb[1]) + 0.0722 * ch(rgb[2])


def _contrast(a, b):
    la, lb = _rel_lum(a), _rel_lum(b)
    return (max(la, lb) + 0.05) / (min(la, lb) + 0.05)


def main() -> int:
    import easyocr
    import numpy as np
    from PIL import Image
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Emu, Pt

    ap = argparse.ArgumentParser()
    ap.add_argument("--case", required=True)
    ap.add_argument("--out-dir", required=True)
    ap.add_argument("--min-page-coverage", type=float, default=0.9,
                    help="only repair pages whose pre-repair OCR coverage is "
                         "below this (safety net where needed; avoids "
                         "injecting misread fragments onto healthy pages)")
    args = ap.parse_args()
    run = Path(args.out_dir)
    final_pptx = run / "final.pptx"
    render_dir = run / "render"
    pages = sorted(p for p in (run / "front_bg").iterdir() if p.is_dir())

    palette_candidates = [(0, 0, 0), (255, 255, 255)]
    pj = run / "styled_bg" / "palette.json"
    if pj.exists():
        pal = json.loads(pj.read_text())
        for e in ((pal.get("result") or {}).get("font_palette") or []):
            rgb = hex_to_rgb(e.get("hex"))
            if rgb:
                palette_candidates.append(rgb)

    prs = Presentation(str(final_pptx))
    reader = easyocr.Reader(["en"], gpu=True, verbose=False)

    # pptx text layer per slide (fuzzy corpus)
    slide_texts: list[list[str]] = []
    for slide in prs.slides:
        texts = []
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False):
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text.strip():
                            texts.append(_norm(r.text))
        slide_texts.append([t for t in texts if t])

    log = {}
    total = 0
    for i, page_dir in enumerate(pages):
        if i >= len(prs.slides):
            break
        rend_png = render_dir / f"page-{i+1}.png"
        orig_png = page_dir / "original.png"
        if not rend_png.exists() or not orig_png.exists():
            continue
        orig_res = reader.readtext(str(orig_png))
        rend_res = reader.readtext(str(rend_png))
        rend_norms = [_norm(t) for _b, t, c in rend_res if c >= CONF_MIN]
        rend_norms = [t for t in rend_norms if t]
        rend_concat = "".join(rend_norms)
        ptx_items = slide_texts[i]
        ptx_concat = "".join(ptx_items)

        # pre-repair page coverage (same word semantics as src/ocr_diff.py)
        def _words(res):
            out = set()
            for _b, t, c in res:
                if c < CONF_MIN:
                    continue
                for w in re.findall(r"[A-Za-z0-9]+", str(t).lower()):
                    if len(w) >= 2:
                        out.add(w)
            return out
        ow, rw = _words(orig_res), _words(rend_res)
        pre_cov = 1 - len(ow - rw) / max(len(ow), 1)
        if pre_cov >= args.min_page_coverage:
            continue

        injections = []
        for bbox, text, conf in orig_res:
            if conf < CONF_MIN:
                continue
            n = _norm(text)
            if len(n) < MIN_NORM_LEN or n.isdigit():
                continue
            if _covered(n, rend_norms, rend_concat) or _covered(n, ptx_items, ptx_concat):
                continue
            xs = [float(p[0]) for p in bbox]
            ys = [float(p[1]) for p in bbox]
            injections.append({"text": str(text).strip(),
                               "bbox": [min(xs), min(ys), max(xs), max(ys)],
                               "conf": round(float(conf), 3),
                               "h_px": round(max(ys) - min(ys), 1)})
        if not injections:
            continue
        rend_img = Image.open(rend_png).convert("RGB")
        slide = prs.slides[i]
        for inj in injections:
            x0, y0, x1, y1 = inj["bbox"]
            pt, dbg = fit_font_size_pt(inj["text"], x1 - x0, y1 - y0,
                                       ink_h_px=inj["h_px"])
            local_bg = _mean_rgb(rend_img, inj["bbox"])
            best = max(palette_candidates, key=lambda c: _contrast(c, local_bg))
            tb = slide.shapes.add_textbox(px_to_emu(x0), px_to_emu(y0),
                                          px_to_emu(x1 - x0), px_to_emu(y1 - y0))
            tb.fill.background(); tb.line.fill.background()
            tf = tb.text_frame
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.margin_left = Emu(0); tf.margin_right = Emu(0)
            tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = inj["text"]
            r.font.size = Pt(pt)
            r.font.color.rgb = RGBColor(*best)
            inj["font_pt"] = round(pt, 2)
            inj["color"] = "#%02X%02X%02X" % best
            inj["local_bg"] = "#%02X%02X%02X" % local_bg
        log[page_dir.name] = injections
        total += len(injections)
        print(f"  [selfrepair] {page_dir.name}: injected {len(injections)} "
              f"({'; '.join(x['text'][:20] for x in injections[:5])})")
    if total:
        prs.save(str(final_pptx))
    (run / "selfrepair.json").write_text(json.dumps(log, indent=1, ensure_ascii=False))
    print(f"selfrepair: {total} injection(s) across {len(log)} page(s)")
    return 0


if __name__ == "__main__":
    sys.exit(main())

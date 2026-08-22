"""Step 3 (class 2): "text with a coloured / shaped background" elements.

Class 2 = "Pure text or formula that comes WITH a coloured or shaped
background" — e.g. an orange call-out panel containing bullet text.

A class-2 element is a HYBRID of class 1 (the text content) and class 4
(the single primitive shape sitting behind the text). We process it in
two layers:

    1. BACKGROUND SHAPE  -- one primitive PPTX shape (usually rectangle
       or rounded rectangle) hard-clamped to the bbox, with the
       observed fill + outline. Driven by Claude in the spirit of
       class 4.

    2. TEXT BOX OVERLAY  -- exactly one textbox placed on top, holding
       the verbatim text. Driven by Claude in the spirit of class 1
       (font size, colour, weight, alignment, line count, horizontal-
       extension permission when the text wraps too much).

The PPTX is built with the shape first (bottom layer) and the textbox
second (top layer). Verify loop renders the composite, crops to the
effective bbox, and asks Claude to check BOTH that the text fits / has
the right line count AND that the background shape's fill/border look
right. Up to MAX_VERIFY_ROUNDS rounds.

Outputs land in ``<step3_class2_dir>/<page>/<stem>.pptx`` and the
per-round crops are saved next to each pptx for inspection.
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from common import (
    CLAUDE_SONNET,
    DPI,
    MIN_READABLE_PT,
    anthropic_client,
    call_claude,
    case_paths,
    crop_bbox,
    encode_image_b64,
    fit_font_size_pt,
    carlito_metrics_unreliable,
    font_sizing_mode,
    hex_to_rgb,
    list_front_bg_pages,
    load_text_metrics,
    parse_json_loose,
    px_to_emu,
    render_pptx_to_png,
)


TARGET_CLASS = 2
MODEL = CLAUDE_SONNET
MAX_VERIFY_ROUNDS = 3


# Allowed shapes for the background layer. Class 2 backgrounds are almost
# always rectangles / rounded rectangles, with the occasional oval. We
# silently coerce anything else to rectangle.
_BG_SHAPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "ellipse": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "trapezoid": MSO_SHAPE.TRAPEZOID,
}
_ALLOWED_BG_TYPES = "rectangle, rounded_rectangle, oval, parallelogram, trapezoid"


EXTRACT_PROMPT = """You are looking at one foreground element extracted from a PowerPoint slide.

The FIRST image is the ORIGINAL FULL SLIDE (for context).
The SECOND image is the FOREGROUND ELEMENT to encode.

This element was classified as "class 2 - pure text or formula WITH a
coloured or shaped background" (e.g. a callout panel that holds bullet
text). It must be reproduced as TWO layered PPTX shapes:

  (a) ONE background SHAPE that fills the element's bbox and provides
      the visible coloured backdrop.
  (b) ONE TEXTBOX placed on top of that shape, holding the verbatim
      text content.

The element's bbox on the slide is (x0, y0, x1, y1) = <<BBOX>> pixels.
Slide size is <<SLIDE_WIDTH>> x <<SLIDE_HEIGHT>> pixels.
The bbox is <<BBOX_W>> x <<BBOX_H>> pixels.

IMPORTANT - STYLE RESTYLING CONTEXT:
The deck this element belongs to has been RESTYLED with a NEW visual palette.
The shape on the new slide must NOT keep the colours of the original
element; instead, choose colours from the NEW DECK PALETTE below so the
class-2 panel harmonises with the restyled deck. The new palette is the
ONLY allowed source for bg_fill_color_hex, bg_line_color_hex, and
font_color_hex - do not invent colours outside it.

NEW DECK PALETTE (theme colours + recommended font colours):
<<PALETTE>>

When picking colours, follow these rules:
  * bg_fill_color_hex : pick the entry from the theme palette whose ROLE
    best matches the role the original element's bg colour played (e.g.
    original was a bright accent panel -> primary or secondary accent
    from the new palette; original was a subtle muted box -> supporting
    tone). It must NOT be "none".
  * bg_line_color_hex : if the original element has a visible outline,
    pick a palette entry that contrasts with bg_fill_color_hex
    (often the theme color itself or a supporting tone). If no outline,
    use "none".
  * bg_line_width_pt : 0 when bg_line_color_hex is "none".
  * font_color_hex : pick from the recommended font palette an entry
    whose contrast against bg_fill_color_hex makes the text comfortably
    legible (primary headline / body text level contrast, not muted).
  * Do not produce a low-contrast pairing (bg and font colours from the
    same end of the palette). If two palette entries would be too
    similar, pick a different font palette entry.

For the BACKGROUND shape, decide:
  * bg_shape_type : one of <<ALLOWED_BG_TYPES>> (rectangle is fine when
    in doubt; pick rounded_rectangle if the corners are visibly rounded).
  * bg_fill_color_hex : "#RRGGBB" chosen from the NEW palette per above.
  * bg_line_color_hex : "#RRGGBB" chosen from the NEW palette, or "none".
  * bg_line_width_pt  : outline width in points; 0 when bg_line_color_hex
    is "none".

For the TEXT layer, decide:
  * text          : the verbatim text content of the element. Preserve
    casing, punctuation, spaces, and line breaks (use "\\n").
  * num_lines     : the TOTAL number of visible text lines you actually
    see in the original element (used downstream as a hard target the
    rendered PPTX must match).
  * font_size_pt  : estimated point size of one line, based on visible
    glyph height vs the bbox height in pixels.
  * font_color_hex: "#RRGGBB" chosen from the NEW font palette per above.
  * bold / italic : booleans.
  * h_align       : "left" | "center" | "right".
  * v_align       : "top" | "middle" | "bottom".

Reply with STRICT JSON only - no markdown, no commentary:
{
  "bg_shape_type": "rectangle|rounded_rectangle|oval|...",
  "bg_fill_color_hex": "#RRGGBB",
  "bg_line_color_hex": "#RRGGBB or none",
  "bg_line_width_pt": <number>,
  "text": "<verbatim text, may include \\n>",
  "num_lines": <integer>,
  "font_size_pt": <number>,
  "font_color_hex": "#RRGGBB",
  "bold": <true|false>,
  "italic": <true|false>,
  "h_align": "left|center|right",
  "v_align": "top|middle|bottom"
}

HARD RULE: the "text" field MUST be a non-empty verbatim transcription
of what you actually see in the element. Returning "" is a hard failure.
If you genuinely cannot read any text in the element, return exactly
{"text": "[unreadable]", "num_lines": 1} so the orchestrator can log it
for review.
"""


VERIFY_PROMPT = """You are verifying a class-2 (text + coloured background) PPTX reconstruction.

Image A - the ORIGINAL foreground element (ground truth FOR LAYOUT / TEXT
content only; its colours are from the OLD deck style and must NOT be
reproduced).
Image B - the CURRENT composite rendered from the PPTX (background shape +
text on top), cropped to the same (possibly extended) bbox region.

The deck has been restyled. Acceptable colours are restricted to the
NEW DECK PALETTE below. Do NOT request to revert any colour to whatever
Image A shows. If you propose a colour change, the new value MUST be
drawn from this palette:

<<PALETTE>>

You MUST check FOUR constraints, in priority order:

  1. TEXT BBOX OVERFLOW: the rendered text must NOT extend beyond the
     textbox bbox shown in Image B.

  2. FONT LEGIBILITY: the rendered font size must NOT go below
     <<MIN_PT>> pt.

  3. LINE COUNT: aim to match the original line count:
     <<ORIG_NUM_LINES>>. If matching it would require dropping the font
     below <<MIN_PT>> pt, accept +/-1 line of soft-wrap instead.
     Legibility wins over exact line count.

  4. BACKGROUND / TEXT COLOUR HARMONY: bg_fill_color_hex, bg_line_color_hex,
     and font_color_hex must all be drawn from the NEW palette above and
     must produce a legible bg/text contrast. The shape type (rectangle
     vs rounded rectangle vs oval, etc.) must still match the original
     element's geometry to a reasonable visual approximation.

If the line counts do NOT match because Image B has MORE lines than the
original (i.e. the textbox is too narrow, so the text soft-wrapped),
the PREFERRED fix is to EXTEND the textbox AND background bbox
horizontally - request via "extend_width_px" (positive integer = extra
pixels to add to the bbox WIDTH; the orchestrator distributes them
across sides that have slack within the slide canvas).

Hard cap: the resulting effective bbox MUST stay inside the slide canvas
(<<SLIDE_W>> x <<SLIDE_H>> pixels). The current effective bbox is
<<CURRENT_BBOX>>; slack: left=<<SLACK_LEFT>>, right=<<SLACK_RIGHT>>,
total=<<SLACK_TOTAL>>. Never request extend_width_px larger than
<<SLACK_TOTAL>>. If total slack is 0, keep "font_size_pt" at or above
<<MIN_PT>> and accept the smallest readable soft-wrap difference.

Prefer extending the bbox over shrinking the font below <<MIN_PT>>.

Current parameters were:
<<CURRENT_PARAMS>>

If Image B is already a faithful, in-bounds reconstruction on all
constraints, reply EXACTLY:
  {"ok": true}

Otherwise reply with the MINIMUM set of adjustments. Always report the
observed line counts. Use this schema; only include keys you actually
want to change:
{
  "ok": false,
  "observed_lines_original": <int>,
  "observed_lines_rendered": <int>,
  "extend_width_px": <int, 0 if none>,
  "bg_shape_type": "rectangle|rounded_rectangle|oval|...",
  "bg_fill_color_hex": "#RRGGBB",
  "bg_line_color_hex": "#RRGGBB or none",
  "bg_line_width_pt": <number>,
  "text": "...",
  "font_size_pt": <number>,
  "font_color_hex": "#RRGGBB",
  "bold": <true|false>,
  "italic": <true|false>,
  "h_align": "left|center|right",
  "v_align": "top|middle|bottom"
}
Strict JSON, no markdown."""


DEFAULT_PARAMS = {
    "bg_shape_type": "rectangle",
    "bg_fill_color_hex": "#EEEEEE",
    "bg_line_color_hex": "none",
    "bg_line_width_pt": 0.0,
    "text": "",
    "num_lines": 1,
    "font_size_pt": 12,
    "font_color_hex": "#000000",
    "bold": False,
    "italic": False,
    "h_align": "left",
    "v_align": "middle",
}

# Keys the verifier is allowed to mutate (num_lines stays a TARGET, not
# adjustable). extend_width_px is handled separately (changes the bbox,
# not the params dict).
ADJUSTABLE_KEYS = (
    "bg_shape_type", "bg_fill_color_hex", "bg_line_color_hex",
    "bg_line_width_pt",
    "text", "font_size_pt", "font_color_hex", "bold", "italic",
    "h_align", "v_align",
)


def _is_valid_extract(params) -> bool:
    return (
        isinstance(params, dict)
        and isinstance(params.get("text"), str)
        and params["text"].strip() not in ("", "[unreadable]")
    )


# ---------------------------------------------------------------------------
# Claude wrappers
# ---------------------------------------------------------------------------

def extract_params(client, original_path, component_path, bbox, slide_w, slide_h,
                   palette_summary: str):
    bw = bbox[2] - bbox[0]
    bh = bbox[3] - bbox[1]
    prompt = (
        EXTRACT_PROMPT
        .replace("<<BBOX>>", f"({bbox[0]:.1f}, {bbox[1]:.1f}, {bbox[2]:.1f}, {bbox[3]:.1f})")
        .replace("<<SLIDE_WIDTH>>", str(slide_w))
        .replace("<<SLIDE_HEIGHT>>", str(slide_h))
        .replace("<<BBOX_W>>", f"{bw:.1f}")
        .replace("<<BBOX_H>>", f"{bh:.1f}")
        .replace("<<ALLOWED_BG_TYPES>>", _ALLOWED_BG_TYPES)
        .replace("<<PALETTE>>", palette_summary)
    )
    content = [
        {"type": "text", "text": "Original slide (full image):"},
        encode_image_b64(original_path),
        {"type": "text", "text": "Foreground element to encode (class 2 - text on coloured background):"},
        encode_image_b64(component_path),
        {"type": "text", "text": prompt},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=1024)
    return raw, parse_json_loose(raw)


def verify_and_adjust(client, component_path, rendered_crop, params,
                      target_lines, eff_bbox, slide_w, slide_h,
                      palette_summary: str):
    x0, y0, x1, y1 = eff_bbox
    sl = max(0.0, x0)
    sr = max(0.0, slide_w - x1)
    prompt = (
        VERIFY_PROMPT
        .replace("<<CURRENT_PARAMS>>", json.dumps(
            {k: v for k, v in params.items() if not k.startswith("_")},
            ensure_ascii=False))
        .replace("<<ORIG_NUM_LINES>>", str(target_lines))
        .replace("<<MIN_PT>>", f"{MIN_READABLE_PT:.0f}")
        .replace("<<SLIDE_W>>", str(slide_w))
        .replace("<<SLIDE_H>>", str(slide_h))
        .replace("<<CURRENT_BBOX>>",
                 f"({x0:.1f}, {y0:.1f}, {x1:.1f}, {y1:.1f})")
        .replace("<<SLACK_LEFT>>", f"{sl:.0f}")
        .replace("<<SLACK_RIGHT>>", f"{sr:.0f}")
        .replace("<<SLACK_TOTAL>>", f"{sl + sr:.0f}")
        .replace("<<PALETTE>>", palette_summary)
    )
    content = [
        {"type": "text", "text": "Image A - original foreground element:"},
        encode_image_b64(component_path),
        {"type": "text", "text": "Image B - current rendered composite (cropped to effective bbox):"},
        encode_image_b64(rendered_crop),
        {"type": "text", "text": prompt},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=1024)
    return raw, parse_json_loose(raw)


# ---------------------------------------------------------------------------
# bbox extension (lifted from step3_class1)
# ---------------------------------------------------------------------------

def extend_bbox_horizontally(bbox, extend_px, slide_w):
    if extend_px <= 0:
        return bbox
    x0, y0, x1, y1 = bbox
    sl = max(0.0, x0)
    sr = max(0.0, slide_w - x1)
    if sl + sr <= 0:
        return bbox
    half = extend_px / 2.0
    add_l = min(half, sl)
    add_r = min(half, sr)
    rem = extend_px - (add_l + add_r)
    if rem > 0:
        extra_l = sl - add_l
        extra_r = sr - add_r
        if extra_l > 0:
            add = min(rem, extra_l)
            add_l += add
            rem -= add
        if rem > 0 and extra_r > 0:
            add_r += min(rem, extra_r)
    return (x0 - add_l, y0, x1 + add_r, y1)


# ---------------------------------------------------------------------------
# PPTX construction
# ---------------------------------------------------------------------------

def build_pptx(slide_w_px, slide_h_px, bbox, params, output_path):
    """Build a one-slide PPTX containing (a) a background primitive shape
    hard-clamped to the bbox and (b) a textbox on top at the same bbox.
    Returns the final pt size used for the text."""
    x0, y0, x1, y1 = bbox
    x0 = max(0.0, min(x0, slide_w_px))
    y0 = max(0.0, min(y0, slide_h_px))
    x1 = max(0.0, min(x1, slide_w_px))
    y1 = max(0.0, min(y1, slide_h_px))
    if x1 <= x0 or y1 <= y0:
        raise ValueError(f"degenerate bbox: {bbox}")

    prs = Presentation()
    prs.slide_width = px_to_emu(slide_w_px)
    prs.slide_height = px_to_emu(slide_h_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = px_to_emu(x0)
    top = px_to_emu(y0)
    width = px_to_emu(x1 - x0)
    height = px_to_emu(y1 - y0)

    # ---- (a) Background shape ----
    bg_type = str(params.get("bg_shape_type", "rectangle")).strip().lower()
    mso = _BG_SHAPE_MAP.get(bg_type, MSO_SHAPE.RECTANGLE)
    bg = slide.shapes.add_shape(mso, left, top, width, height)
    bg.left, bg.top = left, top
    bg.width, bg.height = width, height

    bg_fill_rgb = hex_to_rgb(params.get("bg_fill_color_hex"))
    if bg_fill_rgb is not None:
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*bg_fill_rgb)
    else:
        bg.fill.background()

    bg_line_rgb = hex_to_rgb(params.get("bg_line_color_hex"))
    if bg_line_rgb is not None:
        bg.line.color.rgb = RGBColor(*bg_line_rgb)
        try:
            lw = float(params.get("bg_line_width_pt") or 0.0)
        except (TypeError, ValueError):
            lw = 0.0
        if lw > 0:
            bg.line.width = Pt(lw)
    else:
        bg.line.fill.background()

    # The shape's default placeholder text frame would render an empty
    # paragraph on top; clear it so it can't interfere with the textbox.
    try:
        bg.text_frame.text = ""
    except Exception:
        pass

    # ---- (b) Text box overlay ----
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.fill.background()
    tb.line.fill.background()

    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
    tf.auto_size = MSO_AUTO_SIZE.NONE

    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get((params.get("v_align") or "middle").lower(), MSO_ANCHOR.MIDDLE)
    halign = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get((params.get("h_align") or "left").lower(), PP_ALIGN.LEFT)

    text = params.get("text") or ""
    lines = text.split("\n") if text else [""]

    if font_sizing_mode() == "deterministic":
        # H1+H2: size from code (ink prior + max-fit), not VLM eyeballing.
        # The bbox includes the background shape's padding, so the easyocr
        # ink prior is the load-bearing signal here; max-fit only caps.
        pt_size, dbg = fit_font_size_pt(
            text, x1 - x0, y1 - y0,
            bold=bool(params.get("bold", False)),
            italic=bool(params.get("italic", False)),
            ink_h_px=params.get("_ink_h_px"),
        )
        dbg["vlm_pt"] = params.get("font_size_pt")
        params["_sizing_debug"] = dbg
    else:
        pt_size = float(params.get("font_size_pt") or 12)
        bbox_h_pt = (y1 - y0) / DPI * 72.0
        max_pt = bbox_h_pt / max(1, len(lines))
        pt_size = max(MIN_READABLE_PT, min(pt_size, max_pt))

    text_rgb = hex_to_rgb(params.get("font_color_hex")) or (0, 0, 0)
    bold = bool(params.get("bold", False))
    italic = bool(params.get("italic", False))
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = halign
        run = p.add_run()
        run.text = line
        run.font.size = Pt(pt_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = RGBColor(*text_rgb)

    prs.save(str(output_path))
    return pt_size


def build_raster_fallback_pptx(slide_w_px, slide_h_px, bbox,
                               component_path, output_path):
    """Preserve the original component as a raster when extraction fails."""
    x0, y0, x1, y1 = bbox
    prs = Presentation()
    prs.slide_width = px_to_emu(slide_w_px)
    prs.slide_height = px_to_emu(slide_h_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        str(component_path),
        px_to_emu(x0), px_to_emu(y0),
        width=px_to_emu(x1 - x0), height=px_to_emu(y1 - y0),
    )
    prs.save(str(output_path))


# ---------------------------------------------------------------------------
# Per-component pipeline
# ---------------------------------------------------------------------------

def process_component(client, page_dir, slide_w, slide_h, original_path,
                      comp_meta, comp_entry, out_root, palette_summary: str,
                      ink_h_px=None):
    page_name = page_dir.name
    cf = comp_entry["component_file"]
    orig_bbox = tuple(comp_meta["bbox_xyxy"])
    stem = Path(cf).stem
    out_dir = out_root / page_name
    out_dir.mkdir(parents=True, exist_ok=True)
    component_path = page_dir / cf

    pptx_path = out_dir / f"{stem}.pptx"
    if pptx_path.exists():
        print(f"  [skip] {pptx_path.name} exists")
        return pptx_path

    print(f"--- [{page_name}] {stem} (class=2) ---")
    raw, params = extract_params(
        client, original_path, component_path, orig_bbox, slide_w, slide_h,
        palette_summary,
    )
    if not _is_valid_extract(params):
        print(f"  [retry] empty/invalid extract; second attempt")
        raw, params = extract_params(
            client, original_path, component_path, orig_bbox, slide_w, slide_h,
            palette_summary,
        )
    if not _is_valid_extract(params):
        fallback_log = out_dir / f"{stem}.FALLBACK.json"
        fallback_log.write_text(json.dumps({
            "stem": stem,
            "component_file": cf,
            "reason": "extract returned empty/invalid text; preserved original raster",
            "raw_response": raw[:2000],
        }, indent=2, ensure_ascii=False))
        build_raster_fallback_pptx(
            slide_w, slide_h, orig_bbox, component_path, pptx_path
        )
        print(f"  [fallback] {stem}: extract failed after retry; "
              f"wrote raster fallback + {fallback_log.name}")
        return pptx_path
    for k, v in DEFAULT_PARAMS.items():
        params.setdefault(k, v)
    if ink_h_px:
        params["_ink_h_px"] = ink_h_px
    try:
        target_lines = max(1, int(params.get("num_lines") or 1))
    except (TypeError, ValueError):
        target_lines = 1
    params["num_lines"] = target_lines
    print(f"  initial: bg={params.get('bg_shape_type')} "
          f"fill={params.get('bg_fill_color_hex')} "
          f"text={params.get('text')[:40]!r}... lines={target_lines}")

    eff_bbox = orig_bbox
    pt_used = None
    for r in range(MAX_VERIFY_ROUNDS + 1):
        pt_used = build_pptx(slide_w, slide_h, eff_bbox, params, pptx_path)
        if r == MAX_VERIFY_ROUNDS:
            print(f"  [round {r}] max rounds reached, accepting")
            break
        try:
            rendered = render_pptx_to_png(pptx_path, slide_w, slide_h)
        except (subprocess.TimeoutExpired, Exception) as e:
            print(f"  [warn] render failed: {e}; stopping verify")
            break
        crop_out = out_dir / f"{stem}_rendered_crop_r{r}.png"
        crop_bbox(rendered, eff_bbox, crop_out)

        raw_v, adj = verify_and_adjust(
            client, component_path, crop_out, params, target_lines,
            eff_bbox, slide_w, slide_h, palette_summary,
        )
        if not isinstance(adj, dict):
            print(f"  [round {r}] verify parse failed; stopping. raw: {raw_v[:200]!r}")
            break
        if adj.get("observed_lines_original") is not None \
                or adj.get("observed_lines_rendered") is not None:
            print(f"  [round {r}] observed lines: "
                  f"orig={adj.get('observed_lines_original')} "
                  f"rendered={adj.get('observed_lines_rendered')} "
                  f"(target={target_lines})")
        if adj.get("ok") is True:
            print(f"  [round {r}] OK")
            break

        # ---- bbox extension ----
        try:
            ext_px = int(adj.get("extend_width_px") or 0)
        except (TypeError, ValueError):
            ext_px = 0
        old_eff = eff_bbox
        if ext_px > 0:
            slack = max(0.0, old_eff[0]) + max(0.0, slide_w - old_eff[2])
            ext_px = min(ext_px, int(slack))
            if ext_px > 0:
                eff_bbox = extend_bbox_horizontally(eff_bbox, ext_px, slide_w)
                print(f"  [round {r}] extend bbox by {ext_px}px -> {eff_bbox}")

        # ---- adjustable params ----
        changed = False
        for k in ADJUSTABLE_KEYS:
            if k in adj and adj[k] != params.get(k):
                params[k] = adj[k]
                changed = True
        if not changed and eff_bbox == old_eff:
            print(f"  [round {r}] no changes proposed; stopping")
            break

    # Always save a final rendered crop next to the pptx.
    try:
        final = render_pptx_to_png(pptx_path, slide_w, slide_h)
        crop_bbox(final, eff_bbox, out_dir / f"{stem}_rendered_crop_final.png")
    except Exception:
        pass
    if isinstance(params.get("_sizing_debug"), dict):
        (out_dir / f"{stem}.sizing.json").write_text(json.dumps(
            {"component_file": cf, **params["_sizing_debug"]}, indent=2))
    print(f"  saved: {pptx_path}  (final pt={pt_used:.2f}, "
          f"bbox={tuple(round(v, 1) for v in eff_bbox)})")
    return pptx_path


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def _palette_summary(palette: dict | None) -> str:
    """One-shot formatter shared with step3_class5 -- mirrors that
    function's output so the same colour vocabulary is in front of Claude
    in both places."""
    if not isinstance(palette, dict):
        return "(no palette.json found; fall back to neutral grey bg + black text)"
    res = palette.get("result") or {}
    theme = res.get("theme_color") or {}
    theme_str = f"{theme.get('hex','')} ({theme.get('name','')})".strip()
    entries = []
    for p in res.get("theme_palette") or []:
        if isinstance(p, dict):
            entries.append(
                f"{p.get('hex','')} {p.get('name','')} [{p.get('role','')}]".strip()
            )
    fonts = []
    for p in res.get("font_palette") or []:
        if isinstance(p, dict):
            fonts.append(
                f"{p.get('hex','')} {p.get('name','')} "
                f"contrast={p.get('contrast_ratio','?')} "
                f"use={p.get('use','')}".strip()
            )
    notes = res.get("notes") or ""
    return (
        f"theme color: {theme_str}\n"
        f"theme palette: {'; '.join(entries) if entries else '(none)'}\n"
        f"font palette (high->low contrast): "
        f"{'; '.join(fonts) if fonts else '(none)'}\n"
        f"notes: {notes}"
    )


def _load_palette(paths) -> dict | None:
    p = paths.styled_bg_dir / "palette.json"
    if not p.exists():
        return None
    try:
        return json.loads(p.read_text())
    except json.JSONDecodeError:
        return None


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--case", required=True)
    ap.add_argument("--out-dir", default=None,
                    help="override per-case output root (default: output/<case>/)")
    args = ap.parse_args()
    paths = case_paths(args.case, out_dir=args.out_dir)
    paths.step3_class2_dir.mkdir(parents=True, exist_ok=True)
    client = anthropic_client()

    palette = _load_palette(paths)
    if palette is None:
        print(f"  [warn] no palette.json under {paths.styled_bg_dir}; "
              f"class-2 colours will fall back to neutral defaults")
    palette_summary = _palette_summary(palette)

    pages = list_front_bg_pages(paths)
    if not pages:
        print(f"no front_bg pages under {paths.front_bg_dir}", file=sys.stderr)
        return 1

    n_written = 0
    for page in pages:
        cat_path = page / "categorization.json"
        meta_path = page / "metadata.json"
        if not cat_path.exists() or not meta_path.exists():
            continue
        cat = json.loads(cat_path.read_text())
        meta = json.loads(meta_path.read_text())
        slide_w = int(meta["query_image_size"]["width"])
        slide_h = int(meta["query_image_size"]["height"])
        original_path = Path(meta["query_image_path"])
        bbox_lookup = {c["component_file"]: c for c in meta["components"]}
        targets = [r for r in cat["results"] if r.get("class") == TARGET_CLASS]
        if not targets:
            continue
        print(f"\n=== Page {page.name}: {len(targets)} class-2 element(s) ===")
        text_metrics = (
            load_text_metrics(page) if font_sizing_mode() == "deterministic" else {}
        )
        for r in targets:
            cm = bbox_lookup.get(r["component_file"])
            if cm is None:
                continue
            try:
                result = process_component(client, page, slide_w, slide_h, original_path,
                                           cm, r, paths.step3_class2_dir, palette_summary,
                                           ink_h_px=(text_metrics.get(r["component_file"])
                                                     or {}).get("median_box_h_px"))
                if result is not None:
                    n_written += 1
            except Exception as e:
                print(f"  [error] {r['component_file']}: {e}")
    print(f"\nDone. Wrote {n_written} pptx file(s) under {paths.step3_class2_dir}")
    return 0


if __name__ == "__main__":
    sys.exit(main())

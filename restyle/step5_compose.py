"""Step 5: compose the final restyled deck.

For every page directory under ``front_bg/`` we build ONE slide in the
output PPTX:

  * Slide canvas size = source slide pixel size.
  * Background picture (full canvas) is the styled bg for the page role:
      page index 0          -> opening.png
      page index N-1        -> closing.png
      everything in between -> content.png
  * Every component is placed on top, at its ORIGINAL bbox:
      class 1/3/4 -> deepcopy every shape-bearing top-level child of the
                     step-4 PPTX's first-slide ``p:spTree`` into the
                     destination slide. Falls back to the step-3 PPTX
                     when step 4 has no output.
      class 2     -> deepcopy shapes from the step-3 class-2 PPTX (step 4
                     does not restyle class 2; its colours were already
                     chosen against the styled bg during step 3).
      class 5     -> ``<stem>_gpt.png`` is placed according to its
                     step3 prompt metadata: logo/solid elements may use
                     fit-outside center-crop, while data charts and
                     natural images use fit-inside with panel fallback.
      class 0     -> skipped (pure decoration).

After every slide is composed, a self-critic pass renders the deck and
asks Claude to flag any shape blocks whose Z-order looks wrong against
the original pre-style-transfer slide. Only Z-order is adjusted -- no
shapes are added, removed, moved, or recoloured. The critic can be
disabled with --no-critic.

Result: ``output/<case>/final.pptx``.
"""

from __future__ import annotations

import argparse
import copy
from io import BytesIO
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import replace
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from common import (
    CLAUDE_OPUS,
    LIBREOFFICE_BIN,
    LIBREOFFICE_TIMEOUT_S,
    anthropic_client,
    call_claude,
    case_paths,
    encode_image_b64,
    hex_to_rgb,
    list_front_bg_pages,
    page_role,
    parse_json_loose,
    px_to_emu,
    styled_bg_for_role,
)


NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def qn(tag: str) -> str:
    prefix, name = tag.split(":")
    return f"{{{NSMAP[prefix]}}}{name}"


COPYABLE_TAGS = {
    qn("p:sp"), qn("p:cxnSp"), qn("p:graphicFrame"),
    qn("p:grpSp"), qn("p:pic"),
}

REL_ATTRS = (qn("r:embed"), qn("r:link"), qn("r:id"))


def _next_shape_id(dst_spTree) -> int:
    ids = []
    for cnv in dst_spTree.iter(qn("p:cNvPr")):
        try:
            ids.append(int(cnv.get("id") or 0))
        except (TypeError, ValueError):
            pass
    return (max(ids) if ids else 1) + 1


def _remap_relationships(clone, src_part, dst_part) -> None:
    """Retarget relationship ids referenced by a cloned shape.

    Direct XML copies keep source rIds such as ``rId2``. In the destination
    slide those ids may point to a different part, commonly the styled
    background image, so copied pictures can render as the wrong asset.
    """
    for el in clone.iter():
        for attr in REL_ATTRS:
            r_id = el.get(attr)
            if not r_id:
                continue
            try:
                rel = src_part.rels[r_id]
            except KeyError:
                continue

            if getattr(rel, "is_external", False):
                new_r_id = dst_part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True
                )
            elif rel.reltype == RT.IMAGE:
                image_part = rel.target_part
                _, new_r_id = dst_part.get_or_add_image_part(
                    BytesIO(image_part.blob)
                )
            else:
                continue
            el.set(attr, new_r_id)


def copy_shapes_from_pptx(src_pptx: Path, dst_slide) -> list:
    """Deepcopy every top-level shape from src_pptx's first slide into
    dst_slide. Returns the list of newly-appended XML elements (in the
    order they were appended). Empty list when src has no slides."""
    src_prs = Presentation(str(src_pptx))
    if not src_prs.slides:
        return []
    src_spTree = src_prs.slides[0].shapes._spTree
    dst_spTree = dst_slide.shapes._spTree
    added = []
    next_id = _next_shape_id(dst_spTree)
    for child in list(src_spTree):
        if child.tag not in COPYABLE_TAGS:
            continue
        clone = copy.deepcopy(child)
        _remap_relationships(clone, src_prs.slides[0].part, dst_slide.part)
        for cnv in clone.iter(qn("p:cNvPr")):
            cnv.set("id", str(next_id))
            next_id += 1
        dst_spTree.append(clone)
        added.append(clone)
    return added


def step4_pptx_for(paths, cls: int, page: str, stem: str) -> Path | None:
    sub = {1: "class1", 3: "class3", 4: "class4"}.get(cls)
    if sub is None:
        return None
    p = paths.step4_dir / sub / page / f"{stem}.pptx"
    return p if p.exists() else None


def step3_pptx_for(paths, cls: int, page: str, stem: str) -> Path | None:
    root = {
        1: paths.step3_class1_dir,
        3: paths.step3_class3_dir,
        4: paths.step3_class4_dir,
    }.get(cls)
    if root is None:
        return None
    p = root / page / f"{stem}.pptx"
    return p if p.exists() else None


def find_class5_gpt_png(class5_page_dir: Path, comp_file: str) -> Path | None:
    stem = Path(comp_file).stem
    p = class5_page_dir / f"{stem}_gpt.png"
    return p if p.exists() else None


def class5_prompt_meta(class5_page_dir: Path, comp_file: str) -> dict:
    stem = Path(comp_file).stem
    p = class5_page_dir / f"{stem}_gpt_prompt.meta.json"
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text())
    except Exception:
        return {}


def _class5_can_center_crop(meta: dict) -> bool:
    category = str(meta.get("category") or "").upper()
    approach = str(meta.get("approach") or "").upper()
    return category == "LOGO_OR_SOLID" and approach == "A"


# F-07: when a class-5 image's aspect ratio differs sharply from its
# bbox, the fit-inside placement leaves transparent slack on two sides.
# To stop the styled background's texture from "leaking" through, paint
# a palette-matched solid rectangle behind the image first.
_CLASS5_LETTERBOX_AR_RATIO = 1.5   # trigger when aspect mismatch > 1.5x
_PALETTE_BG_CACHE: dict[str, str | None] = {}


def _palette_bg_hex(paths) -> str | None:
    """Return the dominant background hex from palette.json, or None.
    Result is cached per styled_bg_dir."""
    key = str(paths.styled_bg_dir)
    if key in _PALETTE_BG_CACHE:
        return _PALETTE_BG_CACHE[key]
    pj = paths.styled_bg_dir / "palette.json"
    out = None
    if pj.exists():
        try:
            j = json.loads(pj.read_text())
            theme = (j.get("result") or {}).get("theme_palette") or []
            for e in theme:
                if "background" in (e.get("role") or "").lower():
                    out = e.get("hex")
                    break
            if out is None and theme:
                out = theme[0].get("hex")
        except Exception:
            out = None
    _PALETTE_BG_CACHE[key] = out
    return out


def _make_block_label(cls: int, entry: dict, stem: str, bbox) -> str:
    """Compact human-readable label for the Z-order critic prompt."""
    hint = entry.get("filename_hint") or stem
    x0, y0, x1, y1 = bbox
    return (
        f"class-{cls} '{hint}'  "
        f"bbox=({int(x0)},{int(y0)},{int(x1)},{int(y1)})"
    )


_STRIP_MAX_FRAC = 0.30


def crop_class5_text_overlap(gpt_path, bbox, text_bboxes, out_dir, stem):
    """Shrink a class-5 placement rect past edge strips that overlap peer
    class1/2 text bboxes, cropping the image proportionally (ghosting fix:
    the repaint bakes the original text into those strips). Interior
    overlaps and strips deeper than _STRIP_MAX_FRAC are left alone.
    Returns (image_path, bbox) — unchanged when nothing applies."""
    x0, y0, x1, y1 = (float(v) for v in bbox)
    ow, oh = x1 - x0, y1 - y0
    if ow <= 1 or oh <= 1:
        return gpt_path, bbox
    nx0, ny0, nx1, ny1 = x0, y0, x1, y1
    for tb in text_bboxes:
        tx0, ty0, tx1, ty1 = (float(v) for v in tb)
        ix0, iy0 = max(nx0, tx0), max(ny0, ty0)
        ix1, iy1 = min(nx1, tx1), min(ny1, ty1)
        if ix1 <= ix0 or iy1 <= iy0:
            continue
        # top strip: intersection touches the panel top edge
        if iy0 <= ny0 + 1 and (iy1 - ny0) <= _STRIP_MAX_FRAC * oh:
            ny0 = max(ny0, iy1)
        # bottom strip
        elif iy1 >= ny1 - 1 and (ny1 - iy0) <= _STRIP_MAX_FRAC * oh:
            ny1 = min(ny1, iy0)
        # left strip
        elif ix0 <= nx0 + 1 and (ix1 - nx0) <= _STRIP_MAX_FRAC * ow:
            nx0 = max(nx0, ix1)
        # right strip
        elif ix1 >= nx1 - 1 and (nx1 - ix0) <= _STRIP_MAX_FRAC * ow:
            nx1 = min(nx1, ix0)
    if (nx0, ny0, nx1, ny1) == (x0, y0, x1, y1):
        return gpt_path, bbox
    if nx1 - nx0 < 0.4 * ow or ny1 - ny0 < 0.4 * oh:
        return gpt_path, bbox  # over-shrunk; keep original
    try:
        with Image.open(gpt_path) as im:
            iw, ih = im.size
            cl = int(round((nx0 - x0) / ow * iw))
            ct = int(round((ny0 - y0) / oh * ih))
            cr = int(round((nx1 - x0) / ow * iw))
            cb = int(round((ny1 - y0) / oh * ih))
            cropped = im.crop((max(0, cl), max(0, ct),
                               min(iw, max(cl + 1, cr)),
                               min(ih, max(ct + 1, cb))))
            out = Path(out_dir) / f"{stem}_gpt_stripcrop.png"
            cropped.save(out)
        print(f"  class5 strip-crop: bbox {bbox} -> "
              f"({nx0:.0f},{ny0:.0f},{nx1:.0f},{ny1:.0f}) (text-overlap ghost fix)")
        return out, (nx0, ny0, nx1, ny1)
    except Exception as e:
        print(f"  [warn] strip-crop failed: {e}")
        return gpt_path, bbox


def build_slide(paths, prs, page_dir, slide_w_px, slide_h_px,
                role: str) -> list:
    """Compose one slide; return the list of foreground "blocks" added,
    in current back-to-front order.

    A block is one logical foreground component (one categorization
    result). It may correspond to multiple XML shapes when copied from a
    multi-shape source PPTX (e.g. step3_class4 outputs). Each block dict
    contains: cls, stem, bbox, label, xml_elements (list of XML elements
    inserted into ``slide.shapes._spTree``, in their natural order).
    The styled-bg picture is NOT part of any block; it sits as the very
    first (back-most) child of spTree and is left alone by the critic.
    """
    page_name = page_dir.name
    print(f"\n=== Building slide for page {page_name} (role={role}) ===")
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = styled_bg_for_role(paths, role)
    if not bg.exists():
        bg = paths.styled_bg_dir / "content.png"
    if not bg.exists():
        raise FileNotFoundError(f"styled background missing: {bg}")
    slide.shapes.add_picture(
        str(bg), 0, 0,
        width=px_to_emu(slide_w_px), height=px_to_emu(slide_h_px),
    )
    print(f"  bg: {bg.name}")

    cat = json.loads((page_dir / "categorization.json").read_text())
    meta = json.loads((page_dir / "metadata.json").read_text())
    bbox_lookup = {c["component_file"]: c for c in meta["components"]}

    spTree = slide.shapes._spTree
    blocks = []

    # Deterministic layering: rasters and shape compositions first, text
    # last (topmost). Insertion previously followed categorization.json
    # order, letting a class-5 raster listed after a text component draw
    # over it (title half-clipped behind a panel; the z-order critic only
    # sometimes rescued it). The critic retains authority for genuine
    # image-over-text originals.
    _LAYER = {5: 0, 4: 1, 3: 2, 2: 3, 1: 4}
    ordered_results = sorted(
        cat.get("results", []),
        key=lambda e: _LAYER.get(e.get("class"), 2),
    )

    for entry in ordered_results:
        cf = entry.get("component_file")
        cls = entry.get("class")
        cm = bbox_lookup.get(cf)
        if cm is None:
            print(f"  [skip] {cf}: not in metadata")
            continue
        bbox = cm["bbox_xyxy"]
        stem = Path(cf).stem
        added_elements = []

        if cls in (1, 3, 4):
            src_pptx = step4_pptx_for(paths, cls, page_name, stem)
            tag = "step4"
            if src_pptx is None:
                src_pptx = step3_pptx_for(paths, cls, page_name, stem)
                tag = "step3 (fallback)"
            if src_pptx is None:
                print(f"  [skip] class-{cls} {stem}: no pptx")
                continue
            added_elements = copy_shapes_from_pptx(src_pptx, slide)
            print(f"  class{cls} <- {tag}: {src_pptx.name} "
                  f"({len(added_elements)} shape(s))")
        elif cls == 2:
            src_pptx = paths.step3_class2_dir / page_name / f"{stem}.pptx"
            if not src_pptx.exists():
                print(f"  [skip] class-2 {stem}: no pptx at {src_pptx}")
                continue
            added_elements = copy_shapes_from_pptx(src_pptx, slide)
            print(f"  class2 <- step3: {src_pptx.name} "
                  f"({len(added_elements)} shape(s))")
        elif cls == 5:
            cls5_dir = paths.step3_class5_dir / page_name
            gpt = find_class5_gpt_png(cls5_dir, cf)
            if gpt is None:
                print(f"  [skip] class-5 {stem}: gpt png missing in {cls5_dir}")
                continue
            # Ghosting fix (item 2b): the repaint carries the ORIGINAL text
            # baked into regions where class1/2 text bboxes overlap the
            # panel; the overlaid textbox then ghosts it and OCR collapses.
            # Shrink the placement rect past EDGE-STRIP intersections
            # (≤30% of the panel side) and crop the image proportionally.
            _peer_text_bboxes = [
                bbox_lookup[e2["component_file"]]["bbox_xyxy"]
                for e2 in ordered_results
                if e2.get("class") in (1, 2)
                and e2.get("component_file") in bbox_lookup
                and bbox_lookup[e2["component_file"]].get("bbox_xyxy")
            ]
            gpt, bbox = crop_class5_text_overlap(
                gpt, bbox, _peer_text_bboxes, cls5_dir, stem)
            x0, y0, x1, y1 = bbox
            bw = float(x1 - x0)
            bh = float(y1 - y0)
            try:
                iw, ih = Image.open(gpt).size
            except Exception as e:
                print(f"  [error] class-5 read size: {e}")
                continue
            if iw <= 0 or ih <= 0:
                print(f"  [skip] class-5 {stem}: degenerate gpt image {iw}x{ih}")
                continue
            n_before = len(spTree)
            bbox_ar = bw / max(1e-6, bh)
            img_ar = iw / max(1, ih)
            ar_ratio = max(bbox_ar / img_ar, img_ar / bbox_ar)
            meta = class5_prompt_meta(cls5_dir, cf)
            category = str(meta.get("category") or "").upper()
            can_crop = _class5_can_center_crop(meta)

            if can_crop and 1.5 <= ar_ratio <= 4.0:
                scale = max(bw / iw, bh / ih)
                nw = max(1, int(round(iw * scale)))
                nh = max(1, int(round(ih * scale)))
                crop_w = max(1, int(round(bw)))
                crop_h = max(1, int(round(bh)))
                crop_path = cls5_dir / f"{stem}_gpt_cropped.png"
                try:
                    with Image.open(gpt) as im:
                        resized = im.convert("RGB").resize((nw, nh), Image.LANCZOS)
                    left = max(0, (nw - crop_w) // 2)
                    top = max(0, (nh - crop_h) // 2)
                    cropped = resized.crop((left, top, left + crop_w, top + crop_h))
                    crop_path.parent.mkdir(parents=True, exist_ok=True)
                    cropped.save(crop_path)
                    slide.shapes.add_picture(
                        str(crop_path),
                        px_to_emu(x0), px_to_emu(y0),
                        width=px_to_emu(bw), height=px_to_emu(bh),
                    )
                    print(f"  class5 crop-fill <- {crop_path.name} "
                          f"bbox={bw:.1f}x{bh:.1f} category={category}")
                except Exception as e:
                    print(f"  [warn] class5 crop failed; using fit-inside: {e}")
                    can_crop = False

            if not (can_crop and 1.5 <= ar_ratio <= 4.0):
                scale = min(bw / iw, bh / ih)
                fw = iw * scale
                fh = ih * scale
                offx = x0 + (bw - fw) / 2.0
                offy = y0 + (bh - fh) / 2.0

                # F-07: if aspect ratios differ sharply, draw a palette
                # backing rectangle behind the image so the empty bbox
                # slack matches the deck colour instead of leaking the
                # styled-bg texture through. DATA_CHART/NATURAL_IMAGE
                # keep this full-content path; they are never cropped.
                bg_hex = _palette_bg_hex(paths)
                bg_rgb = hex_to_rgb(bg_hex) if bg_hex else None
                if bg_rgb is not None and (
                    ar_ratio > _CLASS5_LETTERBOX_AR_RATIO
                    or category in {"DATA_CHART", "NATURAL_IMAGE"}
                    or meta.get("fallback")
                ):
                    from pptx.enum.shapes import MSO_SHAPE
                    from pptx.dml.color import RGBColor as _RGB
                    rect = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        px_to_emu(x0), px_to_emu(y0),
                        px_to_emu(bw), px_to_emu(bh),
                    )
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = _RGB(*bg_rgb)
                    rect.line.fill.background()
                    try:
                        rect.text_frame.text = ""
                    except Exception:
                        pass
                    print(f"  class5 panel: aspect bbox={bbox_ar:.2f} "
                          f"vs img={img_ar:.2f} category={category} -> #{bg_hex}")
                slide.shapes.add_picture(
                    str(gpt),
                    px_to_emu(offx), px_to_emu(offy),
                    width=px_to_emu(fw), height=px_to_emu(fh),
                )
                if category == "DATA_CHART" or meta.get("fallback"):
                    try:
                        from pptx.enum.shapes import MSO_SHAPE
                        from pptx.dml.color import RGBColor as _RGB
                        frame = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            px_to_emu(x0), px_to_emu(y0),
                            px_to_emu(bw), px_to_emu(bh),
                        )
                        frame.fill.background()
                        frame.line.color.rgb = _RGB(120, 120, 120)
                        frame.line.width = 12700
                    except Exception:
                        pass
                print(f"  class5 <- {gpt.name} fit {fw:.1f}x{fh:.1f} "
                      f"at ({offx:.1f},{offy:.1f})  bbox={bw:.1f}x{bh:.1f} "
                      f"category={category or 'unknown'}")
            added_elements = list(spTree)[n_before:]
        elif cls == 0:
            print(f"  [skip] class-0 {stem}: decorative")
            continue
        else:
            print(f"  [skip] class-{cls} {stem}: unsupported")
            continue

        if not added_elements:
            continue
        blocks.append({
            "cls": cls,
            "stem": stem,
            "bbox": bbox,
            "label": _make_block_label(cls, entry, stem, bbox),
            "xml_elements": added_elements,
        })

    return blocks


# ---------------------------------------------------------------------------
# Self-critic: render the composed pptx, ask Claude per slide whether the
# Z-order matches the original. Conservative: only reorder when Claude
# explicitly disagrees with the current stacking.
# ---------------------------------------------------------------------------

CRITIC_PROMPT = """You are reviewing the Z-ORDER of foreground shape blocks on ONE slide of a restyled PowerPoint deck.

You will see TWO images:
  Image A -- the ORIGINAL slide (BEFORE style transfer). Use it ONLY as
  ground truth for WHICH element should be ABOVE which.
  Image B -- the CURRENT restyled slide as rendered. New styled
  background, new colours, with all foreground elements placed on top.

The current slide has the following foreground shape blocks, listed from
BACK (drawn first) to FRONT (drawn last). Each block is one logical
component (e.g. one text block, one image, one shape group):

<<BLOCKS>>

Your job: ONLY adjust the Z-order if any block is wrongly hidden,
obscured, or shown out of order compared to the original (Image A).

Strict rules:
  - You are NOT allowed to add, remove, or modify shapes.
  - You are NOT allowed to change positions, sizes, colours, or anything
    other than Z-order.
  - Be CONSERVATIVE. Only flag a reorder if a block is clearly hidden,
    occluded, or visually mis-stacked in Image B in a way that
    contradicts Image A. Cosmetic differences from style transfer (new
    palette, new textures, new decorations on the background) are NOT
    reasons to reorder.

Reply with STRICT JSON only, no markdown:
  If the current Z-order is correct (or there is no clear evidence of a
  problem):
    {"ok": true}
  Otherwise, propose the DESIRED order:
    {"ok": false,
     "new_order": [<block ids back-to-front>],
     "reason": "<one short sentence>"}

new_order MUST be a permutation of the block ids 1..N -- no duplicates,
no omissions, no new ids.
"""


def _natural_slide_key(p: Path) -> int:
    m = re.search(r"-(\d+)\.png$", p.name)
    return int(m.group(1)) if m else 0


def render_pptx_all_slides(pptx_path: Path, out_dir: Path,
                           slide_w_px: int, slide_h_px: int) -> list[Path]:
    """Render every slide of ``pptx_path`` to its own resized PNG under
    ``out_dir``. Returns list of PNG paths in slide order. The pipeline
    is: libreoffice .pptx -> .pdf, then pdftoppm .pdf -> per-page .png,
    then Pillow resize to (slide_w_px, slide_h_px)."""
    work = Path(tempfile.mkdtemp(prefix="pptx_render_all_"))
    profile = Path(tempfile.mkdtemp(prefix="lo_profile_"))
    try:
        env = {k: v for k, v in os.environ.items() if k != "LD_LIBRARY_PATH"}
        subprocess.run(
            [
                LIBREOFFICE_BIN, "--headless",
                f"-env:UserInstallation=file://{profile}",
                "--convert-to", "pdf",
                "--outdir", str(work), str(pptx_path),
            ],
            check=True,
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
            timeout=LIBREOFFICE_TIMEOUT_S, env=env,
        )
        pdfs = list(work.glob("*.pdf"))
        if not pdfs:
            raise RuntimeError(f"libreoffice produced no PDF from {pptx_path}")
        pdf_path = pdfs[0]
        prefix = work / "slide"
        subprocess.run(
            ["pdftoppm", "-r", "96", "-png", str(pdf_path), str(prefix)],
            check=True,
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
            timeout=LIBREOFFICE_TIMEOUT_S,
        )
        pngs = sorted(work.glob("slide-*.png"), key=_natural_slide_key)
        out_dir.mkdir(parents=True, exist_ok=True)
        out_pngs = []
        for i, p in enumerate(pngs, 1):
            dst = out_dir / f"slide_{i:03d}.png"
            with Image.open(p) as im:
                im.convert("RGB").resize(
                    (slide_w_px, slide_h_px), Image.LANCZOS
                ).save(dst)
            out_pngs.append(dst)
        return out_pngs
    finally:
        shutil.rmtree(work, ignore_errors=True)
        shutil.rmtree(profile, ignore_errors=True)


def apply_zorder(slide, blocks, new_order) -> None:
    """Reorder the foreground blocks on a slide to match ``new_order``
    (a permutation of 1..len(blocks), interpreted back-to-front).

    Implementation: detach every block's XML elements from ``spTree``,
    then re-append them in the requested order. The styled-bg picture is
    untouched because it isn't in ``blocks``."""
    spTree = slide.shapes._spTree
    for b in blocks:
        for el in b["xml_elements"]:
            try:
                spTree.remove(el)
            except ValueError:
                pass
    for idx in new_order:
        b = blocks[idx - 1]
        for el in b["xml_elements"]:
            spTree.append(el)


# ---------------------------------------------------------------------------
# F-06: Second self-critic pass — font sizes. The Z-order critic only
# adjusts stacking; this second pass examines visual hierarchy and lets
# Claude propose per-block font-size scales (clamped to [0.6, 1.6]) so
# titles read as titles, body as body, captions as captions. Floor is
# enforced at MIN_READABLE_PT (10 pt) to never undo F-04's legibility.
# ---------------------------------------------------------------------------

FONT_SIZE_CRITIC_PROMPT = """You are reviewing the FONT SIZES of foreground text blocks on ONE slide of a restyled PowerPoint deck.

You will see TWO images:
  Image A -- the ORIGINAL slide (BEFORE style transfer). Use it as
  ground truth for the relative visual hierarchy: which blocks should
  read as the LARGEST text (titles, headlines), which as body text /
  bullets, which as caption / footer / page-number / citation.
  Image B -- the CURRENT restyled slide as rendered.

The current slide has the following editable foreground text blocks.
Each entry shows the block id, a short text snippet, the dominant
current point size used, and the bbox in slide pixels:

<<BLOCKS>>

Your job: per block, propose a multiplicative scale factor that brings
the rendered font size closer to the visual hierarchy of Image A.

Strict rules:
  - You may ONLY change font size. Do NOT change colour, position,
    bold/italic, alignment, or text content.
  - Scale factor MUST be in [0.6, 1.6]. Larger jumps are not allowed.
  - The post-scale font size MUST stay >= <<MIN_PT>> pt. If a proposed
    scale would push the resulting pt below the floor, the orchestrator
    will clamp upward to <<MIN_PT>>.
  - Most blocks should usually keep scale = 1.0. Only flag a block when
    its rendered size noticeably differs from the original's visual
    weight (title rendered as body text, body rendered as caption,
    caption rendered as body, etc.).
  - Be conservative — flag at most the 3 most visually-wrong blocks.

Reply with STRICT JSON only, no markdown:
{"adjustments": [
   {"block_id": <int>, "scale": <float in [0.6, 1.6]>,
    "reason": "<one short sentence>"},
   ... (omit blocks with scale = 1.0)
]}

If no block needs adjustment, reply: {"adjustments": []}
"""

A_RPR = qn("a:rPr")
A_T = qn("a:t")

_FONT_SCALE_MIN = 0.6
_FONT_SCALE_MAX = 1.6
_FONT_MIN_PT = 10.0   # mirrors common.MIN_READABLE_PT


def _block_pt_sizes(block) -> list[int]:
    """All sz attrs (centi-points) on a:rPr nodes inside the block."""
    out = []
    for el in block["xml_elements"]:
        for rpr in el.iter(A_RPR):
            sz = rpr.get("sz")
            if sz:
                try:
                    out.append(int(sz))
                except ValueError:
                    pass
    return out


def _block_text_snippet(block, max_chars: int = 40) -> str:
    chunks = []
    n = 0
    for el in block["xml_elements"]:
        for t in el.iter(A_T):
            if t.text:
                chunks.append(t.text.strip())
                n += len(t.text)
                if n >= max_chars:
                    break
        if n >= max_chars:
            break
    s = " ".join(c for c in chunks if c)
    return s[:max_chars] if s else "(empty)"


def _block_has_text(block) -> bool:
    return bool(_block_pt_sizes(block))


def _scale_block_fonts(block, scale: float, min_pt: float = _FONT_MIN_PT) -> int:
    """Multiply every sz attr by ``scale``, floored at ``min_pt`` (centi-pt)."""
    floor = int(round(min_pt * 100))
    n = 0
    for el in block["xml_elements"]:
        for rpr in el.iter(A_RPR):
            sz = rpr.get("sz")
            if not sz:
                continue
            try:
                new = max(floor, int(round(int(sz) * scale)))
                if new != int(sz):
                    rpr.set("sz", str(new))
                    n += 1
            except ValueError:
                pass
    return n


def font_size_critic_pass(claude, prs, slides_blocks, original_pngs,
                          rendered_pngs) -> None:
    for i, (slide, blocks) in enumerate(zip(prs.slides, slides_blocks)):
        text_blocks = [b for b in blocks if _block_has_text(b)]
        if len(text_blocks) < 2:
            print(f"  [font-critic] slide {i+1}: {len(text_blocks)} text block(s), skip")
            continue
        if i >= len(original_pngs) or i >= len(rendered_pngs):
            print(f"  [font-critic] slide {i+1}: missing render; skip")
            continue
        orig_png = original_pngs[i]
        rend_png = rendered_pngs[i]
        if not orig_png.exists() or not rend_png.exists():
            print(f"  [font-critic] slide {i+1}: render path missing; skip")
            continue

        lines = []
        for j, b in enumerate(text_blocks):
            sizes = _block_pt_sizes(b)
            avg_pt = (sum(sizes) / len(sizes) / 100.0) if sizes else 12.0
            x0, y0, x1, y1 = b["bbox"]
            lines.append(
                f"  #{j+1}  '{_block_text_snippet(b)}'  "
                f"current_pt={avg_pt:.1f}  "
                f"bbox_px={int(x1-x0)}x{int(y1-y0)}"
            )
        blocks_str = "\n".join(lines)
        prompt = (FONT_SIZE_CRITIC_PROMPT
                  .replace("<<BLOCKS>>", blocks_str)
                  .replace("<<MIN_PT>>", f"{_FONT_MIN_PT:.0f}"))
        content = [
            {"type": "text",
             "text": "Image A -- ORIGINAL slide (pre-style-transfer):"},
            encode_image_b64(orig_png),
            {"type": "text",
             "text": "Image B -- CURRENT restyled slide:"},
            encode_image_b64(rend_png),
            {"type": "text", "text": prompt},
        ]
        try:
            raw = call_claude(claude, content, model=CLAUDE_OPUS, max_tokens=1024)
        except Exception as e:
            print(f"  [font-critic] slide {i+1}: claude error: {e}; keep")
            continue
        parsed = parse_json_loose(raw)
        if not isinstance(parsed, dict):
            print(f"  [font-critic] slide {i+1}: parse fail; keep")
            continue
        adjs = parsed.get("adjustments") or []
        if not adjs:
            print(f"  [font-critic] slide {i+1}: no adjustments")
            continue
        n_applied = 0
        for a in adjs:
            try:
                bid = int(a["block_id"])
                scale = float(a["scale"])
            except (KeyError, TypeError, ValueError):
                continue
            if not (1 <= bid <= len(text_blocks)):
                continue
            if not (_FONT_SCALE_MIN <= scale <= _FONT_SCALE_MAX):
                continue
            if abs(scale - 1.0) < 0.02:
                continue
            n_runs = _scale_block_fonts(text_blocks[bid - 1], scale)
            reason = (a.get("reason") or "").strip()[:80]
            print(f"  [font-critic] slide {i+1} block #{bid}: x{scale:.2f} "
                  f"-> {n_runs} run(s) ({reason})")
            n_applied += 1
        if n_applied == 0:
            print(f"  [font-critic] slide {i+1}: all adjustments rejected")


def critic_pass(claude, prs, slides_blocks, original_pngs,
                rendered_pngs) -> None:
    for i, (slide, blocks) in enumerate(zip(prs.slides, slides_blocks)):
        if len(blocks) < 2:
            print(f"  [critic] slide {i+1}: {len(blocks)} block(s), skip")
            continue
        if i >= len(original_pngs) or i >= len(rendered_pngs):
            print(f"  [critic] slide {i+1}: missing render; skip")
            continue
        orig_png = original_pngs[i]
        rend_png = rendered_pngs[i]
        if not orig_png.exists() or not rend_png.exists():
            print(f"  [critic] slide {i+1}: render path missing; skip")
            continue
        blocks_str = "\n".join(
            f"  #{j+1}  {b['label']}" for j, b in enumerate(blocks)
        )
        content = [
            {"type": "text",
             "text": "Image A -- ORIGINAL slide (pre-style-transfer):"},
            encode_image_b64(orig_png),
            {"type": "text",
             "text": "Image B -- CURRENT restyled slide:"},
            encode_image_b64(rend_png),
            {"type": "text",
             "text": CRITIC_PROMPT.replace("<<BLOCKS>>", blocks_str)},
        ]
        try:
            raw = call_claude(claude, content, model=CLAUDE_OPUS, max_tokens=512)
        except Exception as e:
            print(f"  [critic] slide {i+1}: claude error: {e}; keep")
            continue
        parsed = parse_json_loose(raw)
        if not isinstance(parsed, dict):
            print(f"  [critic] slide {i+1}: parse fail; keep. raw={raw[:160]!r}")
            continue
        if parsed.get("ok") is True:
            print(f"  [critic] slide {i+1}: ok ({len(blocks)} blocks)")
            continue
        new_order = parsed.get("new_order")
        if not isinstance(new_order, list):
            print(f"  [critic] slide {i+1}: no new_order; keep")
            continue
        try:
            new_order_int = [int(x) for x in new_order]
        except (TypeError, ValueError):
            print(f"  [critic] slide {i+1}: non-int order; keep")
            continue
        n = len(blocks)
        if sorted(new_order_int) != list(range(1, n + 1)):
            print(f"  [critic] slide {i+1}: not a permutation: {new_order_int}; keep")
            continue
        if new_order_int == list(range(1, n + 1)):
            print(f"  [critic] slide {i+1}: identity reorder; keep")
            continue
        apply_zorder(slide, blocks, new_order_int)
        reason = (parsed.get("reason") or "").strip()[:140]
        print(f"  [critic] slide {i+1}: reorder -> {new_order_int}  ({reason})")


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--case", required=True)
    ap.add_argument("--out-dir", default=None)
    ap.add_argument("--no-critic", action="store_true",
                    help="skip the post-compose Z-order self-critic pass")
    ap.add_argument("--no-font-critic", action="store_true",
                    help="skip the post-compose font-size self-critic pass")
    args = ap.parse_args()

    paths = case_paths(args.case)

    if args.out_dir is not None:
        out = Path(args.out_dir)
        paths = replace(
            paths,
            front_bg_dir=out / "front_bg",
            styled_bg_dir=out / "styled_bg",
            step3_class1_dir=out / "step3_class1",
            step3_class2_dir=out / "step3_class2",
            step3_class3_dir=out / "step3_class3",
            step3_class4_dir=out / "step3_class4",
            step3_class5_dir=out / "step3_class5",
            step4_dir=out / "step4_recolor",
            final_pptx=out / "final.pptx",
        )

    for r in ("opening", "content", "closing"):
        p = paths.styled_bg_dir / f"{r}.png"
        if not p.exists():
            print(f"styled background missing: {p}", file=sys.stderr)
            return 1

    pages = list_front_bg_pages(paths)
    if not pages:
        print(f"no front_bg pages under {paths.front_bg_dir}", file=sys.stderr)
        return 1

    page0_meta = json.loads((pages[0] / "metadata.json").read_text())
    slide_w = int(page0_meta["query_image_size"]["width"])
    slide_h = int(page0_meta["query_image_size"]["height"])
    print(f"Slide size: {slide_w} x {slide_h} px")

    prs = Presentation()
    prs.slide_width = px_to_emu(slide_w)
    prs.slide_height = px_to_emu(slide_h)

    n_pages = len(pages)
    slides_blocks = []
    composed_pages = []
    for i, page in enumerate(pages):
        cat_path = page / "categorization.json"
        meta_path = page / "metadata.json"
        if not cat_path.exists() or not meta_path.exists():
            print(f"[skip] {page.name}: missing categorization.json or metadata.json")
            continue
        blocks = build_slide(paths, prs, page, slide_w, slide_h,
                             page_role(i, n_pages))
        slides_blocks.append(blocks)
        composed_pages.append(page)

    paths.final_pptx.parent.mkdir(parents=True, exist_ok=True)

    if args.no_critic or not slides_blocks:
        prs.save(str(paths.final_pptx))
        print(f"\nSaved: {paths.final_pptx} (critic skipped)")
        return 0

    critic_dir = paths.final_pptx.parent / "_critic"
    critic_dir.mkdir(parents=True, exist_ok=True)
    pre_pptx = critic_dir / "pre_critic.pptx"
    prs.save(str(pre_pptx))
    print(f"\n=== Self-critic pass ===")
    print(f"  pre-critic pptx: {pre_pptx}")

    try:
        rendered_pngs = render_pptx_all_slides(
            pre_pptx, critic_dir / "renders", slide_w, slide_h
        )
    except Exception as e:
        print(f"  [critic] render failed: {e}; saving without critic")
        prs.save(str(paths.final_pptx))
        return 0

    original_pngs = [p / "original.png" for p in composed_pages]
    claude = anthropic_client()
    critic_pass(claude, prs, slides_blocks, original_pngs, rendered_pngs)

    if not args.no_font_critic:
        # Re-render after Z-order so the font critic sees the latest state.
        mid_pptx = critic_dir / "mid_critic.pptx"
        prs.save(str(mid_pptx))
        print(f"\n=== Self-critic pass (font sizes) ===")
        try:
            rendered_pngs_mid = render_pptx_all_slides(
                mid_pptx, critic_dir / "renders_mid", slide_w, slide_h
            )
        except Exception as e:
            print(f"  [font-critic] mid-render failed: {e}; skipping font pass")
            rendered_pngs_mid = rendered_pngs
        font_size_critic_pass(claude, prs, slides_blocks,
                              original_pngs, rendered_pngs_mid)

    prs.save(str(paths.final_pptx))
    print(f"\nSaved: {paths.final_pptx}")
    return 0


if __name__ == "__main__":
    sys.exit(main())

"""Step 3 (class 3): convert every class-3 element (typically a TABLE)
into a native, editable PPTX table.

Two-pass extraction so single-call output never truncates:
  A) STRUCTURE: ask Claude for n_rows/n_cols, col/row fractions, header
     flags, borders, and a global style anchor (body/header/first_col).
  B) CELLS in row batches: per-batch we pass the agreed structure +
     style anchor, then ask for cells for rows [row_start, row_end).

Then a single render -> verify -> adjust round with the table image.

Outputs: ``<page>/<stem>.pptx`` under the case's step3_class3 dir.
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from common import (
    CLAUDE_SONNET,
    DPI,
    EMU_PER_PX,
    MIN_READABLE_PT,
    anthropic_client,
    call_claude,
    case_paths,
    crop_bbox,
    encode_image_b64,
    hex_to_rgb,
    list_front_bg_pages,
    parse_json_loose,
    px_to_emu,
    render_pptx_to_png,
)


TARGET_CLASS = 3
MODEL = CLAUDE_SONNET
MAX_VERIFY_ROUNDS = 2
MAX_ROWS = 32
MAX_COLS = 16
STRUCTURE_MAX_TOKENS = 2048
CELLS_BATCH_MAX_TOKENS = 8192
VERIFY_MAX_TOKENS = 16384
BATCH_CELLS_TARGET = 48


STRUCTURE_PROMPT = """You are looking at one foreground element extracted from a PowerPoint slide.

The FIRST image is the ORIGINAL FULL SLIDE (for context).
The SECOND image is the FOREGROUND ELEMENT to encode as a NATIVE, EDITABLE PPTX TABLE.

The element's bbox on the slide is (x0, y0, x1, y1) = <<BBOX>> pixels.
Slide size is <<SLIDE_WIDTH>> x <<SLIDE_HEIGHT>> pixels.
The bbox is <<BBOX_W>> x <<BBOX_H>> pixels.

THIS PASS returns ONLY the table STRUCTURE plus a small global STYLE ANCHOR.

- Grid size: n_rows, n_cols (1 <= n_rows <= <<MAX_ROWS>>, 1 <= n_cols <= <<MAX_COLS>>).
- header_row: bool. first_col_header: bool.
- col_widths_frac (length n_cols, summing to ~1.0).
- row_heights_frac (length n_rows, summing to ~1.0).
- border_color_hex, border_width_pt.
- style_anchor: {body, header, first_col} each with font_size_pt,
  font_color_hex, bold, italic, h_align, v_align, fill_color_hex.

Hard constraint: the table is hard-clamped to the bbox.

STYLE / PALETTE CONTEXT: this deck has been restyled. Pick all colour
fields (border_color_hex, font_color_hex, fill_color_hex) from the
NEW DECK PALETTE below so the table harmonises with the restyled deck.
For text colours prefer entries from the FONT palette; for cell fills
prefer THEME palette entries.

<<PALETTE>>

Reply with STRICT JSON only:
{
  "n_rows": <int>, "n_cols": <int>,
  "header_row": <bool>, "first_col_header": <bool>,
  "col_widths_frac": [<floats>], "row_heights_frac": [<floats>],
  "border_color_hex": "#RRGGBB", "border_width_pt": <number>,
  "style_anchor": {
    "body":      {"font_size_pt": <n>, "font_color_hex": "#RRGGBB", "bold": <b>, "italic": <b>, "h_align": "...", "v_align": "...", "fill_color_hex": "#RRGGBB or none"},
    "header":    {...same shape...},
    "first_col": {...same shape...}
  }
}
"""


CELLS_BATCH_PROMPT = """You are looking at the SAME PPTX table foreground element as before.

The FIRST image is the ORIGINAL FULL SLIDE (for context).
The SECOND image is the FOREGROUND TABLE element.

The table STRUCTURE has ALREADY been agreed and is FIXED for this call:
  n_rows = <<N_ROWS>>
  n_cols = <<N_COLS>>
  header_row = <<HEADER_ROW>>
  first_col_header = <<FIRST_COL_HEADER>>
  col_widths_frac = <<COL_WIDTHS>>
  row_heights_frac = <<ROW_HEIGHTS>>

STYLE ANCHOR (carry-over context):
<<STYLE_ANCHOR>>

DECK PALETTE (use these for any fill_color_hex / font_color_hex you
emit so the table colours harmonise with the restyled deck):
<<PALETTE>>

THIS call: return CELL CONTENTS for rows [<<ROW_START>>, <<ROW_END>>) ONLY.

For each cell:
  text, font_size_pt, bold, italic, h_align, v_align, font_color_hex,
  fill_color_hex - defaulting to the style anchor unless the row really
  looks different.

Reply with STRICT JSON only:
{
  "row_start": <<ROW_START>>, "row_end": <<ROW_END>>,
  "rows": [
    [
      {"text": "...", "font_size_pt": <n>, "bold": <b>, "italic": <b>,
       "h_align": "...", "v_align": "...",
       "font_color_hex": "#RRGGBB", "fill_color_hex": "#RRGGBB or none"},
      ... (exactly n_cols entries)
    ],
    ... (exactly row_end - row_start entries)
  ]
}
"""


VERIFY_PROMPT = """You are verifying a PPTX table reconstruction.

Image A - the ORIGINAL foreground element.
Image B - the CURRENT table rendered from the PPTX, cropped to the bbox.

The rendered table must NOT extend beyond the bbox, and no cell text
may overflow its cell.

Current params: <<CURRENT_PARAMS>>

If Image B is a faithful, in-bounds reconstruction, reply EXACTLY: {"ok": true}

Otherwise reply with a COMPLETE replacement JSON object:
{
  "ok": false,
  "n_rows": <int>, "n_cols": <int>,
  "header_row": <b>, "first_col_header": <b>,
  "col_widths_frac": [<n_cols floats>],
  "row_heights_frac": [<n_rows floats>],
  "border_color_hex": "#RRGGBB", "border_width_pt": <n>,
  "cells": [
    [ {"text": "...", ...full cell schema...}, ... (n_cols) ],
    ... (n_rows)
  ]
}
Strict JSON, no markdown."""


_DEFAULT_STYLE = {
    "font_size_pt": 10, "font_color_hex": "#000000",
    "bold": False, "italic": False,
    "h_align": "left", "v_align": "middle",
    "fill_color_hex": "none",
}


def _normalize_style_group(d):
    out = dict(_DEFAULT_STYLE)
    if isinstance(d, dict):
        for k in _DEFAULT_STYLE:
            if k in d and d[k] is not None:
                out[k] = d[k]
    return out


def _normalize_anchor(sa):
    if not isinstance(sa, dict):
        sa = {}
    body = _normalize_style_group(sa.get("body"))
    header = _normalize_style_group(sa.get("header") or body)
    first_col = _normalize_style_group(sa.get("first_col") or body)
    return {"body": body, "header": header, "first_col": first_col}


def _cell_from_anchor(anchor, ri, ci, header_row, first_col_header):
    if header_row and ri == 0:
        base = anchor["header"]
    elif first_col_header and ci == 0:
        base = anchor["first_col"]
    else:
        base = anchor["body"]
    return {**base, "text": ""}


def _validate_structure(s):
    if not isinstance(s, dict):
        return False
    try:
        nr = int(s.get("n_rows"))
        nc = int(s.get("n_cols"))
    except (TypeError, ValueError):
        return False
    return 1 <= nr <= MAX_ROWS and 1 <= nc <= MAX_COLS


def extract_structure(client, original_path, component_path, bbox, slide_w, slide_h,
                      palette_summary: str = ""):
    bw = bbox[2] - bbox[0]
    bh = bbox[3] - bbox[1]
    prompt = (
        STRUCTURE_PROMPT
        .replace("<<BBOX>>", f"({bbox[0]:.1f}, {bbox[1]:.1f}, {bbox[2]:.1f}, {bbox[3]:.1f})")
        .replace("<<SLIDE_WIDTH>>", str(slide_w))
        .replace("<<SLIDE_HEIGHT>>", str(slide_h))
        .replace("<<BBOX_W>>", f"{bw:.1f}")
        .replace("<<BBOX_H>>", f"{bh:.1f}")
        .replace("<<MAX_ROWS>>", str(MAX_ROWS))
        .replace("<<MAX_COLS>>", str(MAX_COLS))
        .replace("<<PALETTE>>", palette_summary or "(no palette)")
    )
    content = [
        {"type": "text", "text": "Original slide:"},
        encode_image_b64(original_path),
        {"type": "text", "text": "Foreground table element (pass A: structure only):"},
        encode_image_b64(component_path),
        {"type": "text", "text": prompt},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=STRUCTURE_MAX_TOKENS)
    return raw, parse_json_loose(raw)


def extract_cells_batch(client, original_path, component_path,
                        structure, anchor, row_start, row_end,
                        palette_summary: str = ""):
    prompt = (
        CELLS_BATCH_PROMPT
        .replace("<<N_ROWS>>", str(structure["n_rows"]))
        .replace("<<N_COLS>>", str(structure["n_cols"]))
        .replace("<<HEADER_ROW>>", "true" if structure.get("header_row") else "false")
        .replace("<<FIRST_COL_HEADER>>", "true" if structure.get("first_col_header") else "false")
        .replace("<<COL_WIDTHS>>", json.dumps(structure.get("col_widths_frac")))
        .replace("<<ROW_HEIGHTS>>", json.dumps(structure.get("row_heights_frac")))
        .replace("<<STYLE_ANCHOR>>", json.dumps(anchor, ensure_ascii=False))
        .replace("<<ROW_START>>", str(row_start))
        .replace("<<ROW_END>>", str(row_end))
        .replace("<<PALETTE>>", palette_summary or "(no palette)")
    )
    content = [
        {"type": "text", "text": "Original slide:"},
        encode_image_b64(original_path),
        {"type": "text",
         "text": f"Foreground table element (pass B: cells [{row_start}, {row_end})):"},
        encode_image_b64(component_path),
        {"type": "text", "text": prompt},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=CELLS_BATCH_MAX_TOKENS)
    return raw, parse_json_loose(raw)


def extract_params(client, original_path, component_path, bbox, slide_w, slide_h,
                   palette_summary: str = ""):
    raw_struct, structure = extract_structure(
        client, original_path, component_path, bbox, slide_w, slide_h,
        palette_summary,
    )
    if not _validate_structure(structure):
        return {"raw_structure": raw_struct, "ok": False,
                "reason": "structure-parse-failed",
                "batches_total": 0, "batches_failed": 0}, None
    nr = int(structure["n_rows"])
    nc = int(structure["n_cols"])
    header_row = bool(structure.get("header_row", False))
    first_col_header = bool(structure.get("first_col_header", False))
    anchor = _normalize_anchor(structure.get("style_anchor"))
    cells_grid = [
        [_cell_from_anchor(anchor, ri, ci, header_row, first_col_header)
         for ci in range(nc)]
        for ri in range(nr)
    ]
    batch_rows = max(1, BATCH_CELLS_TARGET // max(1, nc))
    batches_total = batches_failed = 0
    for r0 in range(0, nr, batch_rows):
        r1 = min(nr, r0 + batch_rows)
        batches_total += 1
        _, batch = extract_cells_batch(
            client, original_path, component_path, structure, anchor, r0, r1,
            palette_summary,
        )
        rows = batch.get("rows") if isinstance(batch, dict) else None
        if not isinstance(rows, list) or not rows:
            batches_failed += 1
            print(f"    [extract] cells batch [{r0},{r1}) failed; using anchor defaults")
            continue
        for i, row in enumerate(rows):
            ri = r0 + i
            if ri >= nr or not isinstance(row, list):
                continue
            for ci in range(nc):
                seed = cells_grid[ri][ci]
                cell = row[ci] if ci < len(row) and isinstance(row[ci], dict) else {}
                cells_grid[ri][ci] = {**seed,
                                      **{k: v for k, v in cell.items() if v is not None}}
    params = {
        "n_rows": nr, "n_cols": nc,
        "header_row": header_row, "first_col_header": first_col_header,
        "col_widths_frac": structure.get("col_widths_frac"),
        "row_heights_frac": structure.get("row_heights_frac"),
        "border_color_hex": structure.get("border_color_hex", "#000000"),
        "border_width_pt": structure.get("border_width_pt", 0.5),
        "style_anchor": anchor,
        "cells": cells_grid,
    }
    info = {"raw_structure": raw_struct, "ok": True,
            "batches_total": batches_total, "batches_failed": batches_failed,
            "batch_rows": batch_rows}
    return info, params


def verify_and_adjust(client, component_path, rendered_crop, params):
    params_for_verify = {k: v for k, v in params.items() if k != "style_anchor"}
    prompt = VERIFY_PROMPT.replace(
        "<<CURRENT_PARAMS>>", json.dumps(params_for_verify, ensure_ascii=False)
    )
    content = [
        {"type": "text", "text": "Image A - original foreground:"},
        encode_image_b64(component_path),
        {"type": "text", "text": "Image B - current rendered table:"},
        encode_image_b64(rendered_crop),
        {"type": "text", "text": prompt},
    ]
    raw = call_claude(client, content, model=MODEL, max_tokens=VERIFY_MAX_TOKENS)
    return raw, parse_json_loose(raw)


def _normalize_fracs(fr, n):
    if not isinstance(fr, list) or len(fr) != n:
        return [1.0 / n] * n
    try:
        fr = [max(1e-6, float(x)) for x in fr]
    except (TypeError, ValueError):
        return [1.0 / n] * n
    s = sum(fr)
    return [x / s for x in fr] if s > 0 else [1.0 / n] * n


def _set_cell_border(cell, color_rgb, width_pt):
    if color_rgb is None:
        color_rgb = (0, 0, 0)
    width_emu = max(1, int(round(width_pt * 12700)))
    tc = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    for idx, tag in enumerate(("a:lnL", "a:lnR", "a:lnT", "a:lnB")):
        ln = etree.Element(qn(tag))
        ln.set("w", str(width_emu))
        ln.set("cap", "flat"); ln.set("cmpd", "sng"); ln.set("algn", "ctr")
        fill = etree.SubElement(ln, qn("a:solidFill"))
        srgb = etree.SubElement(fill, qn("a:srgbClr"))
        srgb.set("val", "%02X%02X%02X" % color_rgb)
        prst = etree.SubElement(ln, qn("a:prstDash"))
        prst.set("val", "solid")
        etree.SubElement(ln, qn("a:round"))
        tcPr.insert(idx, ln)


def _apply_cell(cell, cp, max_pt_for_row):
    fill_rgb = hex_to_rgb(cp.get("fill_color_hex"))
    if fill_rgb is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*fill_rgb)
    else:
        cell.fill.background()
    cell.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get((cp.get("v_align") or "middle").lower(), MSO_ANCHOR.MIDDLE)
    cell.margin_left = Emu(18000); cell.margin_right = Emu(18000)
    cell.margin_top = Emu(9000); cell.margin_bottom = Emu(9000)
    tf = cell.text_frame
    tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    p.alignment = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get((cp.get("h_align") or "left").lower(), PP_ALIGN.LEFT)
    text = "" if cp.get("text") is None else str(cp.get("text"))
    run = p.add_run()
    run.text = text
    try:
        pt = float(cp.get("font_size_pt") or 12)
    except (TypeError, ValueError):
        pt = 12.0
    pt = max(MIN_READABLE_PT, min(pt, max_pt_for_row))
    run.font.size = Pt(pt)
    run.font.bold = bool(cp.get("bold", False))
    run.font.italic = bool(cp.get("italic", False))
    color = hex_to_rgb(cp.get("font_color_hex")) or (0, 0, 0)
    run.font.color.rgb = RGBColor(*color)


def build_pptx(slide_w_px, slide_h_px, bbox, params, out_path):
    x0, y0, x1, y1 = bbox
    x0 = max(0.0, min(x0, slide_w_px))
    y0 = max(0.0, min(y0, slide_h_px))
    x1 = max(0.0, min(x1, slide_w_px))
    y1 = max(0.0, min(y1, slide_h_px))
    if x1 <= x0 or y1 <= y0:
        raise ValueError(f"degenerate bbox: {bbox}")
    bw = x1 - x0; bh = y1 - y0
    nr = max(1, min(int(params.get("n_rows") or 1), MAX_ROWS))
    nc = max(1, min(int(params.get("n_cols") or 1), MAX_COLS))
    col_fr = _normalize_fracs(params.get("col_widths_frac"), nc)
    row_fr = _normalize_fracs(params.get("row_heights_frac"), nr)
    prs = Presentation()
    prs.slide_width = px_to_emu(slide_w_px)
    prs.slide_height = px_to_emu(slide_h_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = px_to_emu(x0); top = px_to_emu(y0)
    width = px_to_emu(bw); height = px_to_emu(bh)
    tbl_shape = slide.shapes.add_table(nr, nc, left, top, width, height)
    tbl_shape.left, tbl_shape.top = left, top
    tbl_shape.width, tbl_shape.height = width, height
    table = tbl_shape.table
    table.first_row = bool(params.get("header_row", False))
    table.first_col = bool(params.get("first_col_header", False))
    table.horz_banding = False
    table.vert_banding = False
    cw = [int(round(f * bw * EMU_PER_PX)) for f in col_fr]
    cw[-1] += width - sum(cw)
    for ci in range(nc):
        table.columns[ci].width = max(1, cw[ci])
    rh = [int(round(f * bh * EMU_PER_PX)) for f in row_fr]
    rh[-1] += height - sum(rh)
    for ri in range(nr):
        table.rows[ri].height = max(1, rh[ri])
    border_rgb = hex_to_rgb(params.get("border_color_hex")) or (0, 0, 0)
    border_w = float(params.get("border_width_pt") or 0.5)
    cells = params.get("cells") or []
    capped = 0
    for ri in range(nr):
        max_pt = max(MIN_READABLE_PT, (row_fr[ri] * bh) / DPI * 72.0 - 1.0)
        row_cells = cells[ri] if ri < len(cells) else []
        for ci in range(nc):
            cp = row_cells[ci] if ci < len(row_cells) else {}
            if not isinstance(cp, dict):
                cp = {}
            try:
                requested = float(cp.get("font_size_pt") or 12)
                if requested > max_pt:
                    capped += 1
            except (TypeError, ValueError):
                pass
            _apply_cell(table.cell(ri, ci), cp, max_pt)
            _set_cell_border(table.cell(ri, ci), border_rgb, border_w)
    prs.save(str(out_path))
    return {"n_rows": nr, "n_cols": nc, "capped": capped}


def _validate_params(p):
    if not isinstance(p, dict):
        return False
    try:
        nr = int(p.get("n_rows"))
        nc = int(p.get("n_cols"))
    except (TypeError, ValueError):
        return False
    if not (1 <= nr <= MAX_ROWS and 1 <= nc <= MAX_COLS):
        return False
    cells = p.get("cells")
    return isinstance(cells, list) and len(cells) >= 1


def _default_params():
    return {
        "n_rows": 1, "n_cols": 1,
        "header_row": False, "first_col_header": False,
        "col_widths_frac": [1.0], "row_heights_frac": [1.0],
        "border_color_hex": "#000000", "border_width_pt": 0.5,
        "cells": [[{**_DEFAULT_STYLE, "text": ""}]],
    }


def process_component(client, page_dir, slide_w, slide_h, original_path,
                      comp_meta, comp_entry, out_root,
                      palette_summary: str = ""):
    page_name = page_dir.name
    cf = comp_entry["component_file"]
    bbox = comp_meta["bbox_xyxy"]
    stem = Path(cf).stem
    out_dir = out_root / page_name
    out_dir.mkdir(parents=True, exist_ok=True)
    component_path = page_dir / cf
    pptx_path = out_dir / f"{stem}.pptx"
    if pptx_path.exists():
        print(f"  [skip] {pptx_path.name} exists")
        return pptx_path
    print(f"--- [{page_name}] {stem} ---")
    info, params = extract_params(
        client, original_path, component_path, bbox, slide_w, slide_h,
        palette_summary,
    )
    if not _validate_params(params):
        print("  [warn] structure failed; using defaults")
        params = _default_params()
    else:
        print(f"  extract: batches={info['batches_total']} failed={info['batches_failed']}")
    style_anchor = params.get("style_anchor")
    summary = None
    for r in range(MAX_VERIFY_ROUNDS + 1):
        summary = build_pptx(slide_w, slide_h, bbox, params, pptx_path)
        if r == MAX_VERIFY_ROUNDS:
            print(f"  [round {r}] max rounds")
            break
        try:
            rendered = render_pptx_to_png(pptx_path, slide_w, slide_h)
        except (subprocess.TimeoutExpired, Exception) as e:
            print(f"  [warn] render failed: {e}")
            break
        crop = out_dir / f"{stem}_rendered_crop_r{r}.png"
        crop_bbox(rendered, bbox, crop)
        raw_v, adj = verify_and_adjust(client, component_path, crop, params)
        if not isinstance(adj, dict):
            print(f"  [round {r}] verify parse failed; stopping")
            break
        if adj.get("ok") is True:
            print(f"  [round {r}] OK")
            break
        if not _validate_params(adj):
            print(f"  [round {r}] invalid replacement; stopping")
            break
        params = adj
        if style_anchor is not None:
            params["style_anchor"] = style_anchor
    try:
        final = render_pptx_to_png(pptx_path, slide_w, slide_h)
        crop_bbox(final, bbox, out_dir / f"{stem}_rendered_crop_final.png")
    except Exception:
        pass
    print(f"  saved: {pptx_path}  (rows={summary['n_rows']} cols={summary['n_cols']})")
    return pptx_path


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--case", required=True)
    ap.add_argument("--out-dir", default=None,
                    help="override per-case output root (default: output/<case>/)")
    args = ap.parse_args()
    paths = case_paths(args.case, out_dir=args.out_dir)
    paths.step3_class3_dir.mkdir(parents=True, exist_ok=True)
    client = anthropic_client()
    # F-08: load the deck palette so every extract / cells call can pick
    # colours that already harmonise with the restyled deck.
    from step3_class5 import _load_palette, _palette_summary
    palette = _load_palette(paths)
    palette_summary = _palette_summary(palette)
    pages = list_front_bg_pages(paths)
    if not pages:
        print(f"no front_bg pages under {paths.front_bg_dir}", file=sys.stderr)
        return 1
    for page in pages:
        cat_path = page / "categorization.json"
        meta_path = page / "metadata.json"
        if not cat_path.exists() or not meta_path.exists():
            continue
        cat = json.loads(cat_path.read_text())
        meta = json.loads(meta_path.read_text())
        slide_w = int(meta["query_image_size"]["width"])
        slide_h = int(meta["query_image_size"]["height"])
        original_path = Path(meta["query_image_path"])
        bbox_lookup = {c["component_file"]: c for c in meta["components"]}
        targets = [r for r in cat["results"] if r.get("class") == TARGET_CLASS]
        if not targets:
            continue
        print(f"\n=== Page {page.name}: {len(targets)} class-3 ===")
        for r in targets:
            cm = bbox_lookup.get(r["component_file"])
            if cm is None:
                continue
            try:
                process_component(client, page, slide_w, slide_h, original_path,
                                  cm, r, paths.step3_class3_dir,
                                  palette_summary)
            except Exception as e:
                print(f"  [error] {r['component_file']}: {e}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
